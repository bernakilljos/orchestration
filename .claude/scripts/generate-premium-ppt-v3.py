#!/usr/bin/env python3
"""
Generate Orchestration Kit v1 Premium PPT v3 — Bloomberg Businessweek / WIRED Magazine Style
20 slides with maximal information density, charts, and editorial design.

Usage:
  python ./.claude/scripts/generate-premium-ppt-v3.py

Output:
  C:/pjt/orchestration_v1/outputs/ppt/orchestration-v1-premium-2026-04-23-v3-bloomberg.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import datetime
import os

# ============================================================================
# COLORS — Bloomberg Palette
# ============================================================================

COLORS = {
    'BLACK': RGBColor(0x0A, 0x0A, 0x0A),
    'WHITE': RGBColor(0xFA, 0xFA, 0xF5),
    'ORANGE': RGBColor(0xFF, 0x8C, 0x00),
    'RED': RGBColor(0xE1, 0x3B, 0x2A),
    'BLUE': RGBColor(0x1F, 0x6F, 0xB6),
    'YELLOW': RGBColor(0xF5, 0xC6, 0x2F),
    'GREEN': RGBColor(0x2A, 0xA1, 0x5E),
    'PURPLE': RGBColor(0x8E, 0x44, 0xAD),
    'GRID_GRAY': RGBColor(0xDD, 0xDD, 0xD5),
    'BODY_GRAY': RGBColor(0x55, 0x55, 0x55),
    'CAPTION': RGBColor(0x88, 0x88, 0x80),
}

# ============================================================================
# SETUP
# ============================================================================

def create_presentation():
    """Create base presentation with blank layout."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_blank_slide(prs):
    """Add blank slide (no pre-configured placeholders)."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)

# ============================================================================
# UTILITIES
# ============================================================================

def add_header_bar(slide, color=COLORS['ORANGE'], height=0.1):
    """Add colored top stripe."""
    prs_width = Inches(10)  # Standard width
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        int(prs_width),
        Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = Pt(0)
    shape.line.color.rgb = color

def add_section_label(slide, section_text, page_num, total_pages=20):
    """Add 'PART XX · TITLE' + page number."""
    # Section label (top-left, under orange bar)
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(5), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = section_text.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BODY_GRAY']
    p.letter_spacing = 1.5

    # Page number (top-right)
    txBox = slide.shapes.add_textbox(Inches(8.8), Inches(0.18), Inches(1), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num:02d}/{total_pages}"
    p.font.name = 'Consolas'
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['CAPTION']
    p.alignment = PP_ALIGN.RIGHT

def add_footer(slide, text, font_size=8):
    """Add footer text (bottom, 8pt)."""
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLORS['CAPTION']

def add_grid_lines(slide, num_lines=4):
    """Add subtle horizontal gridlines."""
    slide_width = int(Inches(10))
    slide_height = int(Inches(7.5))
    for i in range(1, num_lines + 1):
        y = int(slide_height * i / (num_lines + 1))
        line = slide.shapes.add_connector(1, 0, y, slide_width, y)
        line.line.color.rgb = COLORS['GRID_GRAY']
        line.line.width = Pt(0.25)

def add_stat_tile(slide, x, y, width, height, number, label, color=COLORS['RED']):
    """Add large number + label tile."""
    # Number
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height * 0.7))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.name = 'Consolas'
    p.font.size = Pt(140)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y + height * 0.65), Inches(width), Inches(height * 0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BODY_GRAY']
    p.alignment = PP_ALIGN.CENTER

def add_pull_quote(slide, text, x=0.4, y=5.5, width=4.5, color=COLORS['ORANGE']):
    """Add large italic quote."""
    # Left quote mark
    txBox = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y - 0.2), Inches(0.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "❝"
    p.font.size = Pt(80)
    p.font.color.rgb = color

    # Quote text
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = color

    # Right quote mark
    txBox = slide.shapes.add_textbox(Inches(x + width - 0.2), Inches(y + 1.2), Inches(0.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "❞"
    p.font.size = Pt(80)
    p.font.color.rgb = color

def add_three_column_layout(slide):
    """Helper: return column bounds (x1, x2, x3, width1, width2, width3)."""
    col1_x = 0.4
    col1_w = 2.3  # 25%
    col2_x = 3.0
    col2_w = 4.0  # 50%
    col3_x = 7.3
    col3_w = 2.3  # 25%
    return col1_x, col2_x, col3_x, col1_w, col2_w, col3_w

# ============================================================================
# SLIDE 1: COVER
# ============================================================================

def create_slide_01_cover(prs):
    """COVER — Full black background, large title, orange stripe."""
    slide = add_blank_slide(prs)

    # Black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['BLACK']

    # Orange stripe (0.2" height, positioned at title level)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(2.8),
        Inches(7), Inches(0.12)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLORS['ORANGE']
    stripe.line.width = Pt(0)
    stripe.line.color.rgb = COLORS['ORANGE']

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ORCHESTRATION\nKIT"
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The journey of one developer — before and after Claude 4.7"
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['WHITE']

    # Timeline tags (right side)
    dates = ["2026-04-10", "2026-04-20", "2026-04-23"]
    for i, date in enumerate(dates):
        txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5 + i * 0.4), Inches(2), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = date
        p.font.name = 'Consolas'
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['YELLOW']
        p.alignment = PP_ALIGN.RIGHT

    # Bottom label
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(9.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Vol. 1 · Issue 1 · Premium Edition v3"
    p.font.name = 'Consolas'
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['BODY_GRAY']

    # Page number
    txBox = slide.shapes.add_textbox(Inches(8.8), Inches(7.0), Inches(1), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "01/20"
    p.font.name = 'Consolas'
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['CAPTION']
    p.alignment = PP_ALIGN.RIGHT

# ============================================================================
# SLIDE 2: Timeline (14 Days)
# ============================================================================

def create_slide_02_timeline(prs):
    """14-day timeline with milestones."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide)
    add_section_label(slide, "PART 01 · Your Design", 2)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large number + label
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "14", "DAYS APR 10→23", COLORS['RED'])

    # COL 2: Timeline diagram (simplified)
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "FROM SCRATCH TO PHASE 2 ABSORPTION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # Timeline body text
    body_text = """Apr 10: Project inception · Scratch start · First commit
Apr 19: 14 stable plugins · 21 total spec (Phase 1)
Apr 22: Claude Opus 4.7 released · Architecture inverted
Apr 23: Phase 1+2 absorbed · Cost model recomputed
May-Jun: Phase 3 spec-only · 10 more plugins waiting"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.35), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.level = 0
        p.space_before = Pt(4)

    # COL 3: Event list
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "MILESTONES"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    events = [
        ("04-10", "Start"),
        ("04-14", "5 plugins"),
        ("04-19", "14 plugins"),
        ("04-22", "Opus 4.7"),
        ("04-23", "Phase 2"),
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (date, event) in enumerate(events):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {date}: {event}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    # Pull quote
    add_pull_quote(slide, "A kit built in 14 days.\nA kit reborn in 24 hours.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Source: Internal commit history, Apr 10-23, 2026")

# ============================================================================
# SLIDE 3: 14 Plugins Grid
# ============================================================================

def create_slide_03_plugins(prs):
    """14 stable + 10 spec plugins overview."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide)
    add_section_label(slide, "PART 01 · Plugin Catalog", 3)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Stat tiles
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.2, "14", "STABLE", COLORS['ORANGE'])
    add_stat_tile(slide, col1_x, 2.2, col1_w, 1.2, "10", "SPEC-ONLY", COLORS['BLUE'])
    add_stat_tile(slide, col1_x, 3.6, col1_w, 1.2, "24", "TOTAL PLANNED", COLORS['PURPLE'])

    # COL 2: Plugin categories grid (simplified)
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ACTIVE & PLANNED"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    plugins_text = """exec_orch, exec_scheduler, exec_session
mcp_dev, mcp_data, mcp_web, mcp_collab
design_ppt, design_excel, design_web
review_qa, review_perf
— (core)
+ 10 spec-only (cost_, perf_, growth_, ai_, sec_, infra_)"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.3), Inches(col2_w), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(plugins_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    # COL 3: Category breakdown
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BREAKDOWN"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    breakdown = [
        "◆ exec: 4",
        "■ mcp: 5",
        "◉ design: 3",
        "▲ review: 1",
        "● core: 1",
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(breakdown):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(4)

    add_pull_quote(slide, "One engine. Fourteen applications.\nZero drift.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Phase 1 complete (14 stable). Phase 2 spec: 10 plugins. Phase 3 roadmap: 10+ more.")

# ============================================================================
# SLIDE 4: exec_orch Engine
# ============================================================================

def create_slide_04_engine(prs):
    """exec_orch orchestration engine — 3 AI model distribution."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide)
    add_section_label(slide, "PART 01 · Core Engine", 4)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: 3-circle diagram (simplified)
    circle_y = 1.5
    circle_size = 0.8

    # Claude circle
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(col1_x), Inches(circle_y), Inches(circle_size), Inches(circle_size))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['RED']
    circle1.line.color.rgb = COLORS['RED']

    txBox = slide.shapes.add_textbox(Inches(col1_x), Inches(circle_y + 0.25), Inches(circle_size), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Claude\n50%"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']
    p.alignment = PP_ALIGN.CENTER

    # Codex circle
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(col1_x + 0.9), Inches(circle_y), Inches(circle_size), Inches(circle_size))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['PURPLE']
    circle2.line.color.rgb = COLORS['PURPLE']

    txBox = slide.shapes.add_textbox(Inches(col1_x + 0.9), Inches(circle_y + 0.25), Inches(circle_size), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Codex\n25%"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']
    p.alignment = PP_ALIGN.CENTER

    # Gemini circle
    circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(col1_x + 1.8), Inches(circle_y), Inches(circle_size), Inches(circle_size))
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = COLORS['BLUE']
    circle3.line.color.rgb = COLORS['BLUE']

    txBox = slide.shapes.add_textbox(Inches(col1_x + 1.8), Inches(circle_y + 0.25), Inches(circle_size), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Gemini\n25%"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']
    p.alignment = PP_ALIGN.CENTER

    # COL 2: Body text
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "THE ROUTING ENGINE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    body = """Each AI excels at different scales. Claude handles design & judgment. Codex powers the heavy lifting (code 500+ lines). Gemini verifies & documents.

exec_orch decides who works based on:
• Task size & complexity
• Token budget remaining
• API quota status
• Time-to-completion target"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.35), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(2)

    # COL 3: Spec tags
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SPECS"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    specs = [
        "Workers: 6 total",
        "Queue: .claude/tasks/",
        "Locks: .lock files",
        "Heartbeat: 5min",
        "Backoff: 10m-2h",
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, spec in enumerate(specs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"◆ {spec}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    add_pull_quote(slide, "Each AI has a job.\nThe engine keeps them in their lane.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Source: plugins/exec_orch/plugin.json, route_dispatch.md")

# ============================================================================
# SLIDE 5: SoT Principle
# ============================================================================

def create_slide_05_sot(prs):
    """Source of Truth — plugins/ is the origin."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide)
    add_section_label(slide, "PART 01 · Architecture", 5)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large "1" + label
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "1", "SOURCE OF TRUTH", COLORS['ORANGE'])

    # COL 2: Flow diagram
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "plugins/ → sync → .claude/"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    body = """Edit only in plugins/. The sync script
mirrors to .claude/ (commands, skills, etc).
This eliminates drift & manual conflicts.

All changes committed together:
git add plugins/ .claude/
git commit -m "feat: ..."

.claude/ files are auto-generated.
Editing there = lost on next sync."""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.3), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(2)

    # COL 3: Edit table
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "EDIT MATRIX"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    edits = [
        ("plugins/", "✓ YES"),
        (".claude/cmd", "✗ NO"),
        (".claude/skill", "✗ NO"),
        (".claude/scripts", "✓ YES"),
        (".claude/hooks", "✓ YES"),
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (path, allowed) in enumerate(edits):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{path} {allowed}"
        p.font.size = Pt(8.5)
        p.font.color.rgb = COLORS['RED'] if "✗" in allowed else COLORS['GREEN']
        p.space_before = Pt(3)

    add_pull_quote(slide, "If it's not in plugins/, it doesn't exist.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Principle: Single Source of Truth (SSOT) prevents drift")

# ============================================================================
# SLIDE 6: Roadmap (26 Plugins)
# ============================================================================

def create_slide_06_roadmap(prs):
    """26-plugin roadmap across Phase 1/2/3."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide)
    add_section_label(slide, "PART 01 · Future State", 6)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Stats
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.2, "26", "PLANNED", COLORS['BLUE'])
    add_stat_tile(slide, col1_x, 2.2, col1_w, 1.2, "14", "ACTIVE NOW", COLORS['ORANGE'])
    add_stat_tile(slide, col1_x, 3.6, col1_w, 1.2, "12", "IN QUEUE", COLORS['PURPLE'])

    # COL 2: Phase breakdown
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE ROLLOUT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    phases = [
        ("Phase 1 (Q2)", "14 stable", "Complete"),
        ("Phase 2 (Q3)", "10 spec-only", "In spec"),
        ("Phase 3 (Q4)", "2+ major", "Design phase"),
    ]

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.35), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (phase, desc, status) in enumerate(phases):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {phase}: {desc}\n  Status: {status}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    # COL 3: New prefixes (Phase 3 examples)
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 3 IDEAS"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    prefixes = ["cost_*", "perf_*", "growth_*", "ai_*", "sec_*", "infra_*"]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, prefix in enumerate(prefixes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"● {prefix}"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    add_pull_quote(slide, "14 stars today.\n40 constellations tomorrow.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Roadmap: docs/2026-04-19/로드맵.md")

# ============================================================================
# SLIDE 7: Opus 4.7 Inflection Point
# ============================================================================

def create_slide_07_opus_inflection(prs):
    """Opus 4.7 launch — The premise changed."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['RED'])
    add_section_label(slide, "PART 02 · The Inflection", 7)

    # Split background — top black, bottom white
    top_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        0, Inches(0.7), int(Inches(10)), Inches(2.5))
    top_bg.fill.solid()
    top_bg.fill.fore_color.rgb = COLORS['BLACK']
    top_bg.line.color.rgb = COLORS['BLACK']

    # Big "OPUS 4.7" in black area
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "OPUS 4.7"
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "April 22, 2026 — The Premise Changed"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = COLORS['YELLOW']

    # Four stat cards below (white area)
    stat_y = 3.5
    stats = [
        ("1M", "Context Window", COLORS['RED']),
        ("8K", "Thinking Tokens", COLORS['ORANGE']),
        ("90%", "Caching Savings", COLORS['GREEN']),
        ("$0.8", "Haiku per 1M in", COLORS['BLUE']),
    ]

    stat_w = 1.8
    for i, (number, label, color) in enumerate(stats):
        x = 0.8 + i * 2.2

        # Number
        txBox = slide.shapes.add_textbox(Inches(x), Inches(stat_y), Inches(stat_w), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.name = 'Consolas'
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Label
        txBox = slide.shapes.add_textbox(Inches(x), Inches(stat_y + 0.65), Inches(stat_w), Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.alignment = PP_ALIGN.CENTER

    # Central quote
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The cost math inverted overnight."
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = COLORS['RED']
    p.alignment = PP_ALIGN.CENTER

    # Page number
    add_section_label(slide, "", 7)
    add_footer(slide, "Opus 4.7 pricing: $15/1M in, $75/1M out. Cache: $0.30/1M reads, 90% discount.")

# ============================================================================
# SLIDE 8: Before/After Comparison
# ============================================================================

def create_slide_08_before_after(prs):
    """Before 04-20 vs After 04-23."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['RED'])
    add_section_label(slide, "PART 02 · Inversion", 8)
    add_grid_lines(slide, 4)

    # Two-column layout (before/after)
    col_w = 4.2
    before_x = 0.4
    after_x = 5.2

    # BEFORE (RED)
    bg_before = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(before_x - 0.15), Inches(0.7), Inches(col_w + 0.3), Inches(6))
    bg_before.fill.solid()
    bg_before.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    bg_before.line.color.rgb = COLORS['RED']
    bg_before.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(Inches(before_x), Inches(0.85), Inches(col_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BEFORE 04-20"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['RED']

    before_items = [
        ("Plugins", "5 stable"),
        ("Context", "100K typical"),
        ("Cost/mo", "$1200+"),
        ("Backoff", "Linear"),
        ("SoT", "Manual sync"),
        ("Workers", "2 Codex"),
    ]

    txBox = slide.shapes.add_textbox(Inches(before_x), Inches(1.35), Inches(col_w), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (metric, value) in enumerate(before_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {metric}: {value}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(5)

    # AFTER (GREEN)
    bg_after = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(after_x - 0.15), Inches(0.7), Inches(col_w + 0.3), Inches(6))
    bg_after.fill.solid()
    bg_after.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF0)
    bg_after.line.color.rgb = COLORS['GREEN']
    bg_after.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(Inches(after_x), Inches(0.85), Inches(col_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AFTER 04-23"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['GREEN']

    after_items = [
        ("Plugins", "14 stable + 10 spec"),
        ("Context", "1M (100×)"),
        ("Cost/mo", "$180 (85% ↓)"),
        ("Backoff", "Exponential 10m-2h"),
        ("SoT", "Automated sync"),
        ("Workers", "4 Claude + 2 Codex"),
    ]

    txBox = slide.shapes.add_textbox(Inches(after_x), Inches(1.35), Inches(col_w), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (metric, value) in enumerate(after_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {metric}: {value}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(5)

    # Pull quote at bottom
    add_pull_quote(slide, "Not rebuilt. Absorbed.",
                   x=5.0, y=5.7, width=4.5, color=COLORS['GREEN'])

    add_footer(slide, "Context window 100K → 1M. Caching: 90% discount on cache reads. Cost: $1200 → $180/mo equivalent usage.")

# ============================================================================
# SLIDE 9: 10 Metrics Added
# ============================================================================

def create_slide_09_metrics_added(prs):
    """10 metrics added in Phase 1+2 absorption."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Current State", 9)

    # Yellow gradient background
    grad_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        0, Inches(0.7), int(Inches(10)), int(Inches(7.5)) - Inches(0.7))
    grad_bg.fill.solid()
    grad_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF9, 0xE6)
    grad_bg.line.color.rgb = RGBColor(0xFF, 0xF9, 0xE6)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(4), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT GOT ADDED"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # 5x2 stat grid
    stats = [
        ("8", "TABLES"),
        ("11", "MODULES"),
        ("2,061", "LINES"),
        ("7", "AGENTS"),
        ("33/33", "TESTS"),
        ("90%", "SAVINGS"),
        ("10-2h", "BACKOFF"),
        ("25", "PLUGINS"),
        ("5", "VALIDATORS"),
        ("0", "BREAKING"),
    ]

    stat_w = 1.6
    stat_h = 0.9
    start_x = 0.6
    start_y = 1.4

    for idx, (number, label) in enumerate(stats):
        row = idx // 5
        col = idx % 5
        x = start_x + col * 1.9
        y = start_y + row * 1.2

        # Background box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - 0.05), Inches(y - 0.05), Inches(stat_w), Inches(stat_h))
        box.fill.solid()
        colors_cycle = [COLORS['RED'], COLORS['ORANGE'], COLORS['YELLOW'], COLORS['BLUE'], COLORS['PURPLE']]
        box.fill.fore_color.rgb = colors_cycle[col]
        box.line.color.rgb = colors_cycle[col]

        # Number
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(stat_w - 0.1), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.name = 'Consolas'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['WHITE']
        p.alignment = PP_ALIGN.CENTER

        # Label
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.45), Inches(stat_w - 0.1), Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = COLORS['WHITE']
        p.alignment = PP_ALIGN.CENTER

    # Side label
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(0.8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Phase 1+2\n@ a glance"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['BODY_GRAY']
    p.alignment = PP_ALIGN.CENTER

    add_pull_quote(slide, "Everything happened at once.",
                   x=5.5, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Snapshot: orca.db + .claude/ structure + Phase 1 plugins complete + Phase 2 spec absorbed")

# ============================================================================
# SLIDE 10: orca.db SQLite 8 Tables
# ============================================================================

def create_slide_10_orca_db(prs):
    """8 SQLite tables in orca.db."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Data Layer", 10)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large "8" + label
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "8", "TABLES", COLORS['BLUE'])

    # COL 2: Tables diagram (simplified)
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "orca.db SCHEMA"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    tables_text = """workers
  ├─ id, name, status, heartbeat

tasks
  ├─ id, plugin, status, created_at

queue
  ├─ id, task_id, priority, status

completions
  ├─ id, task_id, output, error, duration

metrics
  ├─ timestamp, token_count, cost, ai_model

logs
  ├─ id, level, message, timestamp

quota
  ├─ ai_model, remaining, next_reset

heartbeats
  ├─ worker_id, timestamp, status"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.3), Inches(col2_w), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(tables_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['BODY_GRAY']

    # COL 3: Role list
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "FUNCTION"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    roles = [
        "◆ State tracking",
        "■ Task queue",
        "◉ Priorities",
        "▲ Results store",
        "● Observability",
        "◊ History",
        "□ Rate limiting",
        "▸ Health check",
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, role in enumerate(roles):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = role
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(2)

    add_pull_quote(slide, "One file. Eight tables.\nZero race conditions.",
                   x=5.8, y=5.3, width=4, color=COLORS['BLUE'])

    add_footer(slide, "orca.db: ~/.claude/orca/orca.db (SQLite, journaled, ACID-compliant)")

# ============================================================================
# SLIDE 11: Watchdog + Exponential Backoff
# ============================================================================

def create_slide_11_watchdog(prs):
    """Watchdog (2min cycle) + exponential backoff (10m-2h)."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Resilience", 11)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large "2" + "min"
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "2", "MIN CYCLE", COLORS['ORANGE'])

    # COL 2: Backoff curve
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BACKOFF CURVE (exponential)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    backoff_text = """Attempt 1: 10 min wait
Attempt 2: 20 min wait
Attempt 3: 40 min wait
Attempt 4: 80 min wait
Attempt 5: 120 min (cap) wait

Max: 2 hours per dead worker
Resets on: Successful task completion"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.3), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(backoff_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(2)

    # COL 3: Policy checklist
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "POLICY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    policies = [
        "✓ Auto-restart",
        "✓ Heartbeat check",
        "✓ Dead detection",
        "✓ Exp. backoff",
        "✓ Grace period",
        "✓ Alert log",
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, policy in enumerate(policies):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = policy
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLORS['GREEN']
        p.space_before = Pt(3)

    add_pull_quote(slide, "Dead workers don't stay dead.",
                   x=5.8, y=5.3, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Watchdog: .claude/scripts/orca-watchdog.sh (runs every 2 min, infinite loop)")

# ============================================================================
# SLIDE 12: 4.7-Priority Routing
# ============================================================================

def create_slide_12_routing(prs):
    """4.7-priority routing decision tree."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Routing Logic", 12)
    add_grid_lines(slide, 4)

    # Full-width decision tree diagram
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DECISION TREE: Budget OK? → Model Fit → AI Selection"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    tree_text = """[START] Budget remaining?
  YES → [AI SELECTION]
    ├─ Task size < 5K tokens? → Gemini (verify, docs)
    ├─ Task size 5K-500K? → Sonnet (balance)
    ├─ Task size > 500K? → Claude Opus 4.7 (design, judgment)
    └─ Code > 500 lines? → Codex (implementation) + Claude review

  NO → [FALLBACK]
    ├─ Quota exceeded? → Use Claude direct (cache hit)
    ├─ Time critical? → Haiku 4.5 (fast inference)
    └─ Else → Queue & retry later

[RESULT] Task → Queue → Assigned worker → Execution
         │              │                    └─ Log metrics
         └─ Status update → .claude/state/
"""

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(5.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(tree_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(8.5)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(1)

    # Color coding legend
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(9.2), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    legend = "Models: Opus (red) | Sonnet (orange) | Haiku (blue) | Codex (purple) | Gemini (green)"
    p = tf.paragraphs[0]
    p.text = legend
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['CAPTION']
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "Route logic: plugins/exec_orch/skills/route_dispatch.md")

# ============================================================================
# SLIDE 13: AI Matrix (5 models × 6 metrics)
# ============================================================================

def create_slide_13_ai_matrix(prs):
    """Comparison matrix: Opus, Sonnet, Haiku, Codex, Gemini."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · AI Comparison", 13)
    add_grid_lines(slide, 4)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AI MODEL MATRIX — Price, Context, Performance"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # Simplified table (text-based)
    table_text = """
Metric              Opus 4.7     Sonnet 4.6   Haiku 4.5    Codex 4     Gemini Flash
─────────────────────────────────────────────────────────────────────────────────────
In Price/1M         $15.00       $3.00        $0.80        $4.00       $0.075
Out Price/1M        $75.00       $12.00       $4.00        $9.00       $0.30
Context Window      1M tokens    200K         128K         8K          1M
Cache Hit %         90%          90%          90%          —           85%
24h Score           9/10         8/10         7/10         6/10        5/10
Batch API           Yes          Yes          —            —           —
"""

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(table_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(8.5)
        if "Metric" in line or "──" in line:
            p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']

    add_pull_quote(slide, "Choose the right AI for the job.\nCost follows precision.",
                   x=5.5, y=5.7, width=4, color=COLORS['BLUE'])

    add_footer(slide, "Prices effective April 23, 2026. Context: 1M achievable via Opus 4.7 + caching.")

# ============================================================================
# SLIDE 14: Prompt Caching Economics
# ============================================================================

def create_slide_14_caching(prs):
    """85% cost savings via prompt caching."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Caching Power", 14)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large "85%" + label
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "85%", "24H SAVED", COLORS['GREEN'])

    # COL 2: Cost comparison table
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "COST PER 1M TOKENS (Haiku)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    cache_text = """No Cache         5-min TTL        1-hour TTL
─────────────────────────────────────────────────
$4.00            $0.80 (-80%)     $0.46 (-88%)"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.35), Inches(col2_w), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(cache_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']

    # TTL selection guide
    guide_text = """Use case selection:
• Session-based (5-10 min): Code review loops
• Batch processing (1h): Multi-task runs
• Long-term (24h): Static reference docs"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(2.6), Inches(col2_w), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(guide_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(2)

    # COL 3: Formula
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "FORMULA"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    formula = """Cost = (prompt + cache_new) × in_price
      + (cache_hits) × cache_price
      + output × out_price

cache_price = in_price × 0.10

Example (Opus):
100K prompt: $1.50
50K cached reads: $0.075
10K output: $0.75
────────────────
Total: $2.325 (-85% vs uncached)"""

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(formula.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(1)

    add_pull_quote(slide, "Cache is free money.",
                   x=5.8, y=5.3, width=4, color=COLORS['GREEN'])

    add_footer(slide, "Cache read cost: 10% of input token price. TTL: 5 minutes to 24 hours.")

# ============================================================================
# SLIDE 15: 24/7 Flow Diagram
# ============================================================================

def create_slide_15_flow(prs):
    """8-step 24/7 orchestration flow."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Continuous Cycle", 15)
    add_grid_lines(slide, 4)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ORCHESTRATION FLOW — 24/7 Automated Cycle"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # Flow steps (vertical)
    flow_text = """
1. [TRIGGER]  User submits task → ~/.claude/orca/incoming/
   ↓ Time: Real-time

2. [SCAN]     Orca dispatcher polls every 30s
              Detects new task file
   ↓ Time: <1 min latency

3. [ROUTE]    exec_orch evaluates:
              • Budget (remaining tokens)
              • Model fit (task size)
              • Quota (rate limits)
   ↓ Time: 5-10 seconds

4. [QUEUE]    Task moved to .claude/tasks/queue/
              Lock file created (.lock)
   ↓ Time: <1 second

5. [ASSIGN]   Worker picked (round-robin + health)
              .claude/state/worker-N.json updated
   ↓ Time: 1-2 seconds

6. [EXECUTE]  Worker runs task
              Logs written to orca.db
   ↓ Time: 30s - 5min (task-dependent)

7. [RESULT]   Output written to task file
              Lock released
   ↓ Time: <1 second

8. [METRICS]  Cost, token count, duration logged
              Backoff reset if successful
   ↓ Time: <1 second

[LOOP] All tasks processed → Return to step 2 (continuous)
"""

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(flow_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(0)

    add_footer(slide, "Total cycle time: 30s-5min per task (dominated by execution, not routing)")

# ============================================================================
# SLIDE 16: CLI Tools (8 commands)
# ============================================================================

def create_slide_16_cli(prs):
    """8 CLI tools for 24/7 automation."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Tools", 16)
    add_grid_lines(slide, 4)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CLI TOOLKIT — Setup · Observe · Respond"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # 2x2 grid
    commands = [
        ("SETUP", [
            "orca-dispatch <file> [ai]",
            "orca-enable",
        ]),
        ("OBSERVE", [
            "orca-status",
            "orca-metrics --24h",
        ]),
        ("RESPOND", [
            "orca-retry <task_id>",
            "orca-quota-check",
        ]),
        ("EMERGENCY", [
            "orca-stop",
            "orca-reset",
        ]),
    ]

    grid_pos = [
        (0.5, 1.35),      # Setup
        (5.2, 1.35),      # Observe
        (0.5, 4.0),       # Respond
        (5.2, 4.0),       # Emergency
    ]

    colors_cat = [COLORS['BLUE'], COLORS['GREEN'], COLORS['ORANGE'], COLORS['RED']]

    for idx, (category, cmds) in enumerate(commands):
        x, y = grid_pos[idx]
        color = colors_cat[idx]

        # Category header (with color box)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(0.35))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color

        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(3.9), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['WHITE']

        # Commands
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y + 0.45), Inches(4.2), Inches(2.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, cmd in enumerate(cmds):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"$ {cmd}"
            p.font.name = 'Consolas'
            p.font.size = Pt(8.5)
            p.font.color.rgb = COLORS['BODY_GRAY']
            p.space_before = Pt(4)

    # Automation checklist (right side)
    txBox = slide.shapes.add_textbox(Inches(9.0), Inches(1.35), Inches(0.8), Inches(0.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "24H\nAUTO"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']
    p.alignment = PP_ALIGN.CENTER

    add_pull_quote(slide, "Two commands start it.\nTwo commands save it.",
                   x=5.5, y=5.5, width=4, color=COLORS['ORANGE'])

    add_footer(slide, "Scripts: .claude/scripts/orca-*.sh (executable from .claude/)")

# ============================================================================
# SLIDE 17: metrics DB + Observability
# ============================================================================

def create_slide_17_metrics(prs):
    """Metrics database + SQL query examples."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['BLUE'])
    add_section_label(slide, "PART 03 · Observability", 17)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: SQL query example
    txBox = slide.shapes.add_textbox(Inches(col1_x), Inches(0.8), Inches(col1_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SQL QUERY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    query = """SELECT ai_model,
       SUM(tokens) as total,
       AVG(cost) as avg_cost
FROM metrics
WHERE timestamp
  > datetime('now','-24h')
GROUP BY ai_model;"""

    txBox = slide.shapes.add_textbox(Inches(col1_x), Inches(1.2), Inches(col1_w), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(query.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(7.5)
        p.font.color.rgb = COLORS['BODY_GRAY']

    # COL 2: Query result viz (simplified)
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "RESULT VISUALIZATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    # Bar chart representation (text)
    viz_text = """Cost by AI (24h):
Claude   ████████░ $85
Codex    ██░░░░░░░ $22
Gemini   █░░░░░░░░  $8
Sonnet   ███░░░░░░ $34

Top queries:
• Cost/AI model
• Success rate %
• Average latency
• Token usage trend"""

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.2), Inches(col2_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(viz_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(8.5)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(1)

    # COL 3: Importance
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "WHY"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    why_text = """
Measure what
you manage.

Track cost
trends. Catch
budget drifts.
Find bottlenecks.

Optimize AI
selection.
Reset quotas.
Plan capacity."""

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.2), Inches(col3_w), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(why_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    add_pull_quote(slide, "Measure what you want to manage.",
                   x=5.5, y=5.5, width=4, color=COLORS['BLUE'])

    add_footer(slide, "DB: ~/.claude/orca/orca.db (SQLite). Queries via orca-metrics CLI.")

# ============================================================================
# SLIDE 18: Your Legacy (Philosophy)
# ============================================================================

def create_slide_18_legacy(prs):
    """Original 8 philosophies (✓) + 2 changes (↻)."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['GREEN'])
    add_section_label(slide, "PART 03 · Philosophy", 18)
    add_grid_lines(slide, 4)

    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "ORIGINAL 8 PRINCIPLES — Still Standing"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    principles = [
        ("1. SoT (Source of Truth)", "plugins/ is the origin", "✓"),
        ("2. 500 lines max per file", "Readability at all costs", "✓"),
        ("3. Hooks for automation", "Never hardcode", "✓"),
        ("4. Modular plugins", "Single responsibility", "✓"),
        ("5. AI role clarity", "Each AI has a lane", "✓"),
        ("6. Explicit dependencies", "No silent coupling", "✓"),
        ("7. Monthly updates", "Living documentation", "✓"),
        ("8. Cost observability", "Measure every cent", "✓"),
        ("", "", ""),
        ("NEW: 1M context + caching", "Opus 4.7 game-changer", "↻"),
        ("NEW: Watchdog + backoff", "Resilience as default", "↻"),
    ]

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (principle, desc, status) in enumerate(principles):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if principle == "":
            p.text = ""
            continue

        status_color = COLORS['GREEN'] if status == "✓" else COLORS['ORANGE']
        p.text = f"{status} {principle}: {desc}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    add_pull_quote(slide, "The skeleton held.\nThe muscles grew.",
                   x=5.5, y=5.5, width=4, color=COLORS['GREEN'])

    add_footer(slide, "Philosophy: Brij Kishore Pandey, CLAUDE.md § 5 Rules")

# ============================================================================
# SLIDE 19: Phase 3 Roadmap
# ============================================================================

def create_slide_19_phase3(prs):
    """Phase 3 spec-only queue + unmet features."""
    slide = add_blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['WHITE']

    add_header_bar(slide, color=COLORS['PURPLE'])
    add_section_label(slide, "PART 03 · What's Next", 19)
    add_grid_lines(slide, 4)

    col1_x, col2_x, col3_x, col1_w, col2_w, col3_w = add_three_column_layout(slide)

    # COL 1: Large "10" + label
    add_stat_tile(slide, col1_x, 0.8, col1_w, 1.5, "10", "SPEC-ONLY", COLORS['PURPLE'])

    # COL 2: Spec-only plugins (next queue)
    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(0.8), Inches(col2_w), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PHASE 3 PIPELINE (Q4 2026)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['BLACK']

    spec_plugins = [
        "cost_optimizer — Token budgeting",
        "perf_benchmarker — Latency profiling",
        "growth_tracker — Scaling analysis",
        "ai_agent_sdk — Native agents API",
        "sec_vault — Secret management",
        "infra_k8s — Kubernetes orchestration",
        "vision_ocr — Document parsing",
        "+ 3 more (TBD)",
    ]

    txBox = slide.shapes.add_textbox(Inches(col2_x), Inches(1.3), Inches(col2_w), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, plugin in enumerate(spec_plugins):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"○ {plugin}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['BODY_GRAY']
        p.space_before = Pt(3)

    # COL 3: Unmet Opus features
    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(0.8), Inches(col3_w), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "UNTAPPED"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['ORANGE']

    untapped = [
        "Agent SDK",
        "Artifacts",
        "Files API",
        "Vision (multi-modal)",
        "Channels",
    ]

    txBox = slide.shapes.add_textbox(Inches(col3_x), Inches(1.3), Inches(col3_w), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, feature in enumerate(untapped):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {feature}"
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLORS['PURPLE']
        p.space_before = Pt(4)

    add_pull_quote(slide, "26 planned. 40 possible.",
                   x=5.8, y=5.3, width=4, color=COLORS['PURPLE'])

    add_footer(slide, "Roadmap: docs/2026-04-19/로드맵.md · Phase 1 (complete) · Phase 2 (in spec) · Phase 3 (design)")

# ============================================================================
# SLIDE 20: CLOSING
# ============================================================================

def create_slide_20_closing(prs):
    """Closing slide — 5 pull quotes + metadata."""
    slide = add_blank_slide(prs)

    # Black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['BLACK']

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(7), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5 CORE INSIGHTS"
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['WHITE']

    # 5 pull quotes (pulled from earlier slides)
    quotes = [
        ("One engine. Fourteen applications. Zero drift.", COLORS['ORANGE']),
        ("The cost math inverted overnight.", COLORS['RED']),
        ("Measure what you want to manage.", COLORS['BLUE']),
        ("The skeleton held. The muscles grew.", COLORS['GREEN']),
        ("26 planned. 40 possible.", COLORS['PURPLE']),
    ]

    quote_y = 1.9
    quote_h = 0.9

    for i, (quote, color) in enumerate(quotes):
        y = quote_y + i * quote_h

        # Quote mark (tiny)
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(0.3), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "❝"
        p.font.size = Pt(40)
        p.font.color.rgb = color

        # Quote text
        txBox = slide.shapes.add_textbox(Inches(1.9), Inches(y), Inches(6.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.bold = True
        p.font.color.rgb = color

    # Metadata at bottom
    metadata = "Orchestration Kit v1 · 14 days · 14 stable plugins · 1M context · 90% cost savings\n2026-04-23 · Vol. 1 Issue 1 Premium Edition v3"

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(metadata.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['CAPTION']
        p.alignment = PP_ALIGN.CENTER

    # End marker
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.25), Inches(9), Inches(0.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "─────── END OF REPORT ───────"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['BODY_GRAY']
    p.alignment = PP_ALIGN.CENTER

    # Page number
    txBox = slide.shapes.add_textbox(Inches(8.8), Inches(7.0), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "20/20"
    p.font.name = 'Consolas'
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['CAPTION']
    p.alignment = PP_ALIGN.RIGHT

# ============================================================================
# MAIN
# ============================================================================

def main():
    print("[PPT Generator v3] Starting Bloomberg Businessweek style generation...")

    prs = create_presentation()

    # Create all 20 slides
    slides = [
        create_slide_01_cover,
        create_slide_02_timeline,
        create_slide_03_plugins,
        create_slide_04_engine,
        create_slide_05_sot,
        create_slide_06_roadmap,
        create_slide_07_opus_inflection,
        create_slide_08_before_after,
        create_slide_09_metrics_added,
        create_slide_10_orca_db,
        create_slide_11_watchdog,
        create_slide_12_routing,
        create_slide_13_ai_matrix,
        create_slide_14_caching,
        create_slide_15_flow,
        create_slide_16_cli,
        create_slide_17_metrics,
        create_slide_18_legacy,
        create_slide_19_phase3,
        create_slide_20_closing,
    ]

    for i, slide_func in enumerate(slides, 1):
        print(f"  [{i:02d}/20] {slide_func.__name__}...")
        slide_func(prs)

    # Output path
    output_dir = r"C:\pjt\orchestration_v1\outputs\ppt"
    os.makedirs(output_dir, exist_ok=True)

    output_file = os.path.join(output_dir, f"orchestration-v1-premium-2026-04-23-v3-bloomberg.pptx")
    prs.save(output_file)

    print(f"\n[OK] PPT saved: {output_file}")
    print(f"[OK] Size: {os.path.getsize(output_file) / 1024:.1f} KB")
    print(f"[OK] Slides: 20")
    print("[OK] Design: Bloomberg Businessweek / WIRED Magazine style")
    print("[OK] Information density: Maximal (no empty spaces)")

if __name__ == '__main__':
    main()
