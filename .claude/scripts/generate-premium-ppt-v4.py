#!/usr/bin/env python3
"""
Orchestration Kit v1 - Premium PPT v4 Generator
Style: Apple Keynote + Stripe Docs + Luxury Minimal
Design: Hero element + micro details, 20 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os

# ============================================================================
# COLOR PALETTE - Luxury Minimal
# ============================================================================
LUX_CREAM = RGBColor(0xF8, 0xF4, 0xEC)         # Main bg
LUX_IVORY = RGBColor(0xFC, 0xF9, 0xF2)         # Secondary bg
LUX_INK = RGBColor(0x0C, 0x0F, 0x15)           # Text

LUX_GOLD = RGBColor(0xB8, 0x8B, 0x4A)          # Primary accent
LUX_DEEP_GOLD = RGBColor(0x8B, 0x6A, 0x32)     # Deep gold
LUX_CHAMPAGNE = RGBColor(0xE8, 0xD8, 0xB0)     # Soft highlight

LUX_STONE = RGBColor(0x6E, 0x68, 0x5C)         # Body text
LUX_FOG = RGBColor(0xC7, 0xC0, 0xB0)           # Divider
LUX_SAND = RGBColor(0xE5, 0xDD, 0xCA)          # Light box

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_gradient_fill(shape, color1, color2):
    """Add linear gradient fill to shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_shadow(shape, blur=3, distance=0.5, opacity=0.15):
    """Add outer shadow to shape"""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.blur_radius = Pt(blur)
        shadow.distance = Pt(distance)
        shadow.alpha = opacity
    except:
        pass  # Shadow not supported on this shape type

def add_eyebrow_header(slide, text, page_num, total_pages=20):
    """Add luxury header with page number"""
    # Eyebrow text
    tf = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(4), Inches(0.25))
    p = tf.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = LUX_GOLD
    p.font.name = "Malgun Gothic"

    # Page number (right aligned)
    tf_page = slide.shapes.add_textbox(Inches(8.5), Inches(0.35), Inches(1.5), Inches(0.25))
    p_page = tf_page.text_frame.paragraphs[0]
    p_page.text = "  " + str(page_num) + " / " + str(total_pages)
    p_page.font.size = Pt(10)
    p_page.font.color.rgb = LUX_STONE
    p_page.font.name = "Consolas"
    p_page.alignment = PP_ALIGN.RIGHT

    # Thin gold divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.65), Inches(8.4), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LUX_GOLD
    line.line.color.rgb = LUX_GOLD

def add_hero_number(slide, number, label, x, y, color=LUX_GOLD, size=180):
    """Add large hero number"""
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(1.2))
    p = tf.text_frame.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER

    # Label below
    tf_label = slide.shapes.add_textbox(Inches(x - 0.5), Inches(y + 1.0), Inches(2.5), Inches(0.3))
    p_label = tf_label.text_frame.paragraphs[0]
    p_label.text = label.upper()
    p_label.font.size = Pt(10)
    p_label.font.bold = True
    p_label.font.color.rgb = LUX_GOLD
    p_label.font.name = "Malgun Gothic"
    p_label.alignment = PP_ALIGN.CENTER

def add_luxury_box(slide, x, y, w, h, text="", color_bg=LUX_SAND, add_shadow_flag=True):
    """Add luxury box with shadow"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color_bg
    box.line.color.rgb = LUX_FOG
    box.line.width = Pt(0.5)

    if add_shadow_flag:
        add_shadow(box, blur=3, distance=0.5, opacity=0.12)

    if text:
        tf = box.text_frame
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = LUX_INK
        p.font.name = "Malgun Gothic"

    return box

def add_title_section(slide, title, subtitle=""):
    """Add section title (H1 style)"""
    tf_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.8))
    p_title = tf_title.text_frame.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(52)
    p_title.font.bold = True
    p_title.font.color.rgb = LUX_INK
    p_title.font.name = "Malgun Gothic"
    p_title.line_spacing = 1.2

    if subtitle:
        tf_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(0.35))
        p_sub = tf_sub.text_frame.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = LUX_STONE
        p_sub.font.name = "Malgun Gothic"

def add_body_text(slide, text, x, y, w, h, size=13, lead=False):
    """Add body text (lead or regular)"""
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tf.text_frame
    p.text = text
    p.word_wrap = True
    p.line_spacing = 1.55 if not lead else 1.6

    para = p.paragraphs[0]
    para.font.size = Pt(size if not lead else 16)
    para.font.color.rgb = LUX_INK
    para.font.name = "Malgun Gothic"

    return tf

# ============================================================================
# PPT GENERATION
# ============================================================================

def create_premium_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    add_gradient_fill(background, LUX_IVORY, LUX_CHAMPAGNE)
    background.line.fill.background()

    # Hero title
    tf_hero = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p_hero = tf_hero.text_frame.paragraphs[0]
    p_hero.text = "Orchestration Kit"
    p_hero.font.size = Pt(160)
    p_hero.font.bold = True
    p_hero.font.color.rgb = LUX_INK
    p_hero.font.name = "Malgun Gothic"
    p_hero.alignment = PP_ALIGN.CENTER

    # Subtitle
    tf_sub = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.5))
    p_sub = tf_sub.text_frame.paragraphs[0]
    p_sub.text = "Multi-AI Orchestration Framework v1"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = LUX_STONE
    p_sub.font.name = "Malgun Gothic"
    p_sub.alignment = PP_ALIGN.CENTER

    # Footer meta
    tf_footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    p_footer = tf_footer.text_frame.paragraphs[0]
    p_footer.text = "Claude 4.7 * Codex * Gemini  -  2026-04-23  -  v1.0"
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = LUX_GOLD
    p_footer.font.name = "Consolas"
    p_footer.alignment = PP_ALIGN.CENTER

    # Slide 2: Timeline
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 01 - GENESIS", 2)
    add_title_section(slide, "The Journey Begins")

    # Timeline narrative
    narrative = "Started April 2026 as a response to single-AI limits.\nClaude designs, Codex implements, Gemini validates.\nInspired by orchestration complexity in real teams.\n\n14 plugins in v1.0 - 25+ planned for v2/v3."

    add_body_text(slide, narrative, 0.8, 2.2, 8.4, 4.5, size=14, lead=True)

    # Slide 3: 14 Stable Plugins
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 01 - PLUGINS", 3)

    # Hero number
    add_hero_number(slide, 14, "Stable Plugins", 4.5, 1.2, color=LUX_GOLD)

    # Category boxes (left)
    categories = [
        ("exec_", "4", "Orchestration"),
        ("mcp_", "6", "Connectors"),
        ("design_", "3", "Output"),
        ("review_", "1", "Validation")
    ]

    y_start = 2.5
    for i, (prefix, count, name) in enumerate(categories):
        y = y_start + i * 0.9
        add_luxury_box(slide, 0.8, y, 2.0, 0.7, count + " " + name, LUX_SAND)

    # Quote (right)
    quote = "14 stable plugins\ncreate the substrate\nfor 40-plugin vision."
    tf_quote = slide.shapes.add_textbox(Inches(3.2), Inches(2.5), Inches(5.5), Inches(3.5))
    p_quote = tf_quote.text_frame.paragraphs[0]
    p_quote.text = quote
    p_quote.font.size = Pt(16)
    p_quote.font.color.rgb = LUX_STONE
    p_quote.font.name = "Malgun Gothic"
    p_quote.font.italic = True
    p_quote.line_spacing = 1.8

    # Slide 4: exec_orch Engine
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 02 - ENGINE", 4)
    add_title_section(slide, "exec_orch: Multi-AI Routing")

    # AI trio boxes
    ais = [
        ("Claude", "Design\nComplex", 1.0),
        ("Codex", "Code\nImpl", 4.2),
        ("Gemini", "Verify\nScale", 7.4)
    ]

    for name, role, x in ais:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.2), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = LUX_SAND
        box.line.color.rgb = LUX_FOG
        add_shadow(box, blur=3)

        tf = box.text_frame
        tf.margin_top = Inches(0.3)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name + "\n\n" + role
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = LUX_INK
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Malgun Gothic"

    # Bottom caption
    add_body_text(slide, "Router decides AI per task: design - Opus 4.7 | code <200L - Sonnet 4.6 | code 500L+ - Codex | verify - Haiku + Gemini",
                  0.8, 4.5, 8.4, 2, size=12)

    # Slide 5: SoT Principle
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 02 - ARCHITECTURE", 5)

    # Hero: "1 of Truth"
    tf_hero = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(1.5))
    p_hero = tf_hero.text_frame.paragraphs[0]
    p_hero.text = "1"
    p_hero.font.size = Pt(180)
    p_hero.font.bold = True
    p_hero.font.color.rgb = LUX_GOLD
    p_hero.font.name = "Consolas"
    p_hero.alignment = PP_ALIGN.CENTER

    tf_of = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(7), Inches(0.8))
    p_of = tf_of.text_frame.paragraphs[0]
    p_of.text = "Source of Truth"
    p_of.font.size = Pt(48)
    p_of.font.bold = True
    p_of.font.color.rgb = LUX_INK
    p_of.alignment = PP_ALIGN.CENTER
    p_of.font.name = "Malgun Gothic"

    # Explanation
    explanation = "plugins/ is canonical. .claude/ syncs automatically.\nNo manual edits to generated files. Prevents drift."
    add_body_text(slide, explanation, 0.8, 4.0, 8.4, 3, size=13)

    # Slide 6: Roadmap Matrix
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 03 - EVOLUTION", 6)
    add_title_section(slide, "Roadmap: 3 Phases")

    phases = [
        ("Phase 1", "4 plugins\n1-2 months", LUX_GOLD),
        ("Phase 2", "13 plugins\n3-6 months", LUX_STONE),
        ("Phase 3", "9 plugins\n6+ months", LUX_FOG)
    ]

    for i, (phase, info, color) in enumerate(phases):
        x = 0.8 + i * 3.0
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        add_shadow(box, blur=2, opacity=0.1)

        tf = box.text_frame
        tf.margin_top = Inches(0.3)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = LUX_INK if color != LUX_FOG else LUX_STONE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "\n" + info
        p2.font.size = Pt(11)
        p2.font.color.rgb = LUX_INK if color != LUX_FOG else LUX_STONE
        p2.alignment = PP_ALIGN.CENTER

    # Slide 7: Premise Changed - 4.7 Era
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    add_gradient_fill(background, LUX_CREAM, LUX_SAND)
    background.line.fill.background()

    # Hero "4.7"
    tf_hero = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    p_hero = tf_hero.text_frame.paragraphs[0]
    p_hero.text = "4.7"
    p_hero.font.size = Pt(180)
    p_hero.font.bold = True
    p_hero.font.color.rgb = LUX_GOLD
    p_hero.font.name = "Consolas"
    p_hero.alignment = PP_ALIGN.CENTER

    # Subtitle
    tf_sub = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    p_sub = tf_sub.text_frame.paragraphs[0]
    p_sub.text = "The Premise Changed"
    p_sub.font.size = Pt(32)
    p_sub.font.color.rgb = LUX_INK
    p_sub.font.name = "Malgun Gothic"
    p_sub.alignment = PP_ALIGN.CENTER

    # Metrics (small)
    metrics_text = "1M context  *  Extended thinking  *  Prompt caching  *  90% cost savings"
    add_body_text(slide, metrics_text, 0.8, 4.5, 8.4, 2.5, size=13)

    # Slide 8: Before/After
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 03 - PARADIGM", 8)
    add_title_section(slide, "Single Model - Multi-Model")

    # Before box
    before_text = "Limited context\nNo caching\nTask isolation\nManual routing\nHigh latency"
    add_luxury_box(slide, 0.8, 2.0, 3.5, 3.8, before_text, LUX_SAND)

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(3.5), Inches(1), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LUX_GOLD
    arrow.line.fill.background()

    # After box
    after_text = "1M context window\n90% cost via caching\nIntelligent routing\nParallel workers\nSub-second response"
    after_box = add_luxury_box(slide, 5.7, 2.0, 3.5, 3.8, after_text, LUX_CHAMPAGNE)

    # Slide 9: 10 New Metrics
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 03 - INSTRUMENTATION", 9)
    add_title_section(slide, "10 New Observability Metrics")

    metrics = [
        ("Cache Hit %", "85-90%"),
        ("Quota Status", "Real-time"),
        ("Worker Health", "8 gauges"),
        ("Task Latency", "p50/p95/p99"),
        ("Cost/Session", "USD tracked"),
        ("Token Spend", "Per AI model"),
        ("Error Rate", "By component"),
        ("Route Accuracy", "Decision audit"),
        ("Budget Burndown", "Daily limit"),
        ("Orca Uptime", "24/7 monitor")
    ]

    for i, (metric, value) in enumerate(metrics):
        row = i // 5
        col = i % 5
        x = 0.8 + col * 1.8
        y = 2.0 + row * 2.0

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = LUX_SAND
        box.line.color.rgb = LUX_FOG
        add_shadow(box, blur=2, opacity=0.1)

        tf = box.text_frame
        tf.margin_top = Inches(0.2)
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = metric
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = LUX_INK
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "\n" + value
        p2.font.size = Pt(9)
        p2.font.color.rgb = LUX_GOLD
        p2.alignment = PP_ALIGN.CENTER

    # Slide 10: orca.db Persistence
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 04 - STATE", 10)
    add_title_section(slide, "orca.db: Unified State Store")

    tables = "8 tables: workers | tasks | metrics | quota | budget | session | routes | heartbeat"
    add_body_text(slide, tables, 0.8, 2.0, 8.4, 0.5, size=12)

    # Table schema visualization
    table_info = "workers: id, status, pid, task_id, started_at, last_heartbeat\ntasks: id, instruction, status, priority, ai_model, created_at\nmetrics: timestamp, metric_name, value, worker_id\nquota: ai_model, quota_limit, used_today, reset_at\nbudget: daily_limit_usd, spent_today, project_id\nsession: session_id, created_at, snapshot_path\nroutes: task_id, ai_selected, latency_ms, token_count\nheartbeat: worker_id, timestamp, cpu%, memory%, task_count"

    add_body_text(slide, table_info, 1.0, 2.8, 8.0, 4.0, size=11)

    # Slide 11: Watchdog 2-minute Interval
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 04 - RELIABILITY", 11)

    # Hero "2 min"
    add_hero_number(slide, "2", "min interval", 4.2, 1.2, color=LUX_GOLD)

    # Explanation
    watchdog_text = "Watchdog checks every 2 minutes:\n  * Worker process alive?\n  * Heartbeat fresh? (within 3 min)\n  * No zombie tasks?\n\nBackoff: 10m - 20m - 40m - 2h\nFallback to Claude if Codex/Gemini quota exceeded."

    add_body_text(slide, watchdog_text, 0.8, 2.2, 8.4, 4.5, size=13)

    # Slide 12: Routing Tree
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 04 - ROUTING", 12)
    add_title_section(slide, "Intelligent Task Router")

    routing_tree = "Task arrives\n  Design/architecture? - Claude Opus 4.7\n  Code <200L? - Claude Sonnet 4.6\n  Code 500L+? - Codex (4 parallel) + Polish\n  Verification? - Haiku (90% cache) + Gemini (>500k)\n  All: Check quota. If exceeded - backoff/fallback\n\nBudget tier: Daily limit (USD) - pause if exceeded"

    add_body_text(slide, routing_tree, 0.8, 2.0, 8.4, 4.8, size=11)

    # Slide 13: AI Cost Matrix
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 04 - ECONOMICS", 13)
    add_title_section(slide, "Claude 4.7 + Cohort Pricing")

    # AI models grid (simplified)
    models = [
        ("Opus 4.7", "$15/$75", "1M", "yes"),
        ("Sonnet 4.6", "$3/$15", "200k", "yes"),
        ("Haiku 4.5", "$0.8/$4", "200k", "yes"),
        ("Codex", "$2.5/$10", "128k", "no"),
        ("Gemini", "$0.08/$0.32", "1M+", "no")
    ]

    y_start = 2.0
    for i, (model, pricing, context, cache) in enumerate(models):
        y = y_start + i * 0.9

        # Model name
        tf = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(1.8), Inches(0.7))
        p = tf.text_frame.paragraphs[0]
        p.text = model
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = LUX_INK

        # Pricing
        tf = slide.shapes.add_textbox(Inches(2.7), Inches(y), Inches(1.8), Inches(0.7))
        p = tf.text_frame.paragraphs[0]
        p.text = pricing
        p.font.size = Pt(10)
        p.font.color.rgb = LUX_GOLD

        # Context
        tf = slide.shapes.add_textbox(Inches(4.7), Inches(y), Inches(1.5), Inches(0.7))
        p = tf.text_frame.paragraphs[0]
        p.text = context
        p.font.size = Pt(10)
        p.font.color.rgb = LUX_STONE

        # Cache
        tf = slide.shapes.add_textbox(Inches(6.4), Inches(y), Inches(1.0), Inches(0.7))
        p = tf.text_frame.paragraphs[0]
        p.text = "yes" if cache == "yes" else "no"
        p.font.size = Pt(10)
        p.font.color.rgb = LUX_GOLD

    # Slide 14: Caching = 85% Cost Savings
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 05 - EFFICIENCY", 14)

    # Hero "85%"
    add_hero_number(slide, "85%", "Cache Savings", 4.0, 1.2, color=LUX_GOLD)

    caching_explanation = "Prompt caching TTL: 5 minutes\n* System prompt cached: shared across workers\n* Request history cached: 3-message lookback\n* Plugin specs cached: 24-hour TTL\n\nResult: 85-90% cost reduction on cached tasks\nSession strategy: cluster similar tasks"

    add_body_text(slide, caching_explanation, 0.8, 2.3, 8.4, 4.5, size=13)

    # Slide 15: 24/7 Automation Flow
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 05 - OPERATIONS", 15)
    add_title_section(slide, "24/7 Automation Loop")

    flow_text = "1. Task arrives - queued to .claude/tasks/\n2. Router decides AI (Opus/Sonnet/Codex/Gemini)\n3. AI executes - result - .claude/tasks/done/\n4. Metrics recorded - orca.db\n5. Watchdog checks every 2 min - restart if needed\n6. Nightly: snapshot - .claude/context-cache/\n7. Budget check: daily limit enforced\n8. Loop: next task - repeat\n\nAll 100% stateless * No manual intervention * Full audit trail"

    add_body_text(slide, flow_text, 0.8, 2.0, 8.4, 4.8, size=12)

    # Slide 16: CLI Tools Ecosystem
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 06 - DEVELOPER EXPERIENCE", 16)
    add_title_section(slide, "CLI Tools & Validation")

    cli_categories = [
        ("sync", "sync-plugins.sh\nvalidate-plugin-schema.py", 1.0, 2.0),
        ("observe", "check-agents\nmetrics-query.py", 4.2, 2.0),
        ("run", "codex-auto [N]\nhaiku-auto [N]", 7.4, 2.0),
        ("debug", "orca-status\nworker-health.sh", 1.0, 4.5)
    ]

    for cat, tools, x, y in cli_categories:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = LUX_SAND
        box.line.color.rgb = LUX_FOG
        add_shadow(box)

        tf = box.text_frame
        tf.margin_top = Inches(0.2)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = cat.upper()
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = LUX_GOLD
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "\n" + tools
        p2.font.size = Pt(9)
        p2.font.color.rgb = LUX_INK
        p2.alignment = PP_ALIGN.CENTER

    # Slide 17: Metrics Query Database
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 06 - INSIGHTS", 17)
    add_title_section(slide, "Metrics: Every Question Answerable")

    queries_text = "SELECT * FROM routes WHERE ai_selected='Codex' AND token_count > 50000;\n\nSELECT ai_model, SUM(token_count) FROM routes WHERE DATE(created_at) = TODAY GROUP BY ai_model;\n\nSELECT COUNT(*) FROM tasks WHERE status='retry_backoff' AND DATE(created_at) = TODAY;"

    add_body_text(slide, queries_text, 0.8, 2.0, 8.4, 4.8, size=10)

    # Slide 18: Legacy Narrative
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    add_gradient_fill(background, LUX_CREAM, LUX_SAND)
    background.line.fill.background()

    # Hero quotation
    quote = "Not rebuilt.\nAbsorbed."

    tf_quote = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(2.5))
    p_quote = tf_quote.text_frame.paragraphs[0]
    p_quote.text = quote
    p_quote.font.size = Pt(80)
    p_quote.font.bold = True
    p_quote.font.color.rgb = LUX_INK
    p_quote.font.name = "Malgun Gothic"
    p_quote.alignment = PP_ALIGN.CENTER
    p_quote.line_spacing = 1.3

    # Attribution
    tf_attr = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    p_attr = tf_attr.text_frame.paragraphs[0]
    p_attr.text = "- v1.0 Philosophy"
    p_attr.font.size = Pt(18)
    p_attr.font.color.rgb = LUX_STONE
    p_attr.font.name = "Malgun Gothic"
    p_attr.alignment = PP_ALIGN.CENTER

    # Slide 19: Phase 3 Roadmap
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LUX_CREAM
    add_eyebrow_header(slide, "PART 07 - FUTURE", 19)
    add_title_section(slide, "Phase 3: 9 Additional Plugins")

    phase3_text = "cost_instagram - multi-channel monetization\ncost_blog - AdSense automation\nperf_monitor - Prometheus + Grafana\ndesign_pdf - generation & signing\ndesign_video - Shorts & thumbnails\nmcp_payment - Toss, Stripe, KakaoPay\nmcp_analytics - GA4, Mixpanel\nexec_scheduler - cron & workflows\ngrowth_seo - long-form rank tracking\n\nTarget: 26 total plugins by end of 2026"

    add_body_text(slide, phase3_text, 0.8, 2.0, 8.4, 4.8, size=12)

    # Slide 20: Closing
    # ===================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    add_gradient_fill(background, LUX_IVORY, LUX_CHAMPAGNE)
    background.line.fill.background()

    # Main message
    tf_msg = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
    p_msg = tf_msg.text_frame.paragraphs[0]
    p_msg.text = "One Prompt,\nMultiple Minds,\nInfinite Possibilities"
    p_msg.font.size = Pt(44)
    p_msg.font.bold = True
    p_msg.font.color.rgb = LUX_INK
    p_msg.font.name = "Malgun Gothic"
    p_msg.alignment = PP_ALIGN.CENTER
    p_msg.line_spacing = 1.4

    # Footer
    tf_footer = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
    p_footer = tf_footer.text_frame.paragraphs[0]
    p_footer.text = "orchestration-v1-premium-2026-04-23  **  Claude Code  **  github.com/bernakilljos"
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = LUX_GOLD
    p_footer.font.name = "Consolas"
    p_footer.alignment = PP_ALIGN.CENTER

    # Save
    output_dir = r"C:\pjt\orchestration_v1\outputs\ppt"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "orchestration-v1-premium-2026-04-23-v4-luxury.pptx")
    prs.save(output_path)

    file_size = os.path.getsize(output_path) / 1024
    print(f"SUCCESS: Premium PPT v4 created")
    print(f"  File: {output_path}")
    print(f"  Size: {file_size:.1f} KB")
    print(f"  Slides: 20 (Luxury minimal design)")
    print(f"  Color palette: Cream/Ivory/Gold/Stone (LUX_*)")
    print(f"  Typography: 9pt-180pt (Malgun Gothic + Consolas)")
    print(f"  Patterns: Hero, Side Narrative, Full Bleed, Grid Matrix")
    print(f"  Luxury details: Gradients, shadows, thin gold dividers")

    return output_path

if __name__ == "__main__":
    create_premium_ppt()
