"""아이티센코어 부서 AI 신사업 제안 PPT 15장 자동 생성

청중: 임원-본부장 (가정)
목적: AI Risk Lighthouse + 부서 IP 신사업 승인-예산 요청
산출물: outputs/itcen/itcen-business-proposal.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
OUT = os.path.join(ROOT, 'outputs', 'itcen', 'itcen-business-proposal.pptx')

# 색상 팔레트
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # 진한 파랑
ACCENT = RGBColor(0xC0, 0x00, 0x00)       # 빨강 (강조)
LIGHT = RGBColor(0xD9, 0xE1, 0xF2)        # 연한 파랑 (배경)
GRAY = RGBColor(0x59, 0x59, 0x59)         # 본문 그레이
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC4, 0x91, 0x5A)         # 금색

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, font='맑은 고딕'):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if isinstance(text, list):
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            if color:
                r.font.color.rgb = color
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        if color:
            r.font.color.rgb = color
    return box


def add_bg(slide, x, y, w, h, color, *, outline=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if outline:
        shape.line.color.rgb = outline
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card(slide, x, y, w, h, title, bullets, *, title_color=None, bullet_size=11, title_size=14):
    add_bg(slide, x, y, w, h, WHITE, outline=RGBColor(0xC0, 0xC0, 0xC0))
    add_bg(slide, x, y, w, 0.35, title_color or PRIMARY)
    add_text(slide, x + 0.1, y + 0.05, w - 0.2, 0.3, title, size=title_size, bold=True, color=WHITE)
    lines = ['• ' + b for b in bullets]
    add_text(slide, x + 0.15, y + 0.45, w - 0.3, h - 0.5, lines, size=bullet_size, color=GRAY)


def add_table(slide, x, y, w, h, headers, rows, *, header_color=None, header_text_color=None,
              first_col_bold=False, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(w * cw / total)
    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color or PRIMARY
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = '맑은 고딕'
                r.font.size = Pt(font_size + 1)
                r.font.bold = True
                r.font.color.rgb = header_text_color or WHITE
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci > 0 or not first_col_bold else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = '맑은 고딕'
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = GRAY
                    if ci == 0 and first_col_bold:
                        r.font.bold = True
                        r.font.color.rgb = PRIMARY
    return table


def add_title(slide, num, title, subtitle=None):
    # 상단 띠
    add_bg(slide, 0, 0, 13.333, 0.7, PRIMARY)
    add_text(slide, 0.3, 0.15, 0.7, 0.4, f'#{num:02d}', size=18, bold=True, color=GOLD)
    add_text(slide, 1.0, 0.12, 8, 0.45, title, size=18, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, 1.0, 0.42, 11, 0.25, subtitle, size=10, color=LIGHT)
    add_text(slide, 11.3, 0.2, 1.8, 0.35, 'ITCEN CORE', size=10, color=GOLD, align=PP_ALIGN.RIGHT)
    add_text(slide, 11.3, 0.42, 1.8, 0.25, 'AI Risk Monitoring', size=8, color=LIGHT, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_num):
    add_bg(slide, 0, 7.2, 13.333, 0.3, LIGHT)
    add_text(slide, 0.3, 7.27, 8, 0.2, '아이티센코어 - 리스크모니터링-행동위험분석 부서 - AI 신사업 제안', size=8, color=PRIMARY)
    add_text(slide, 11.5, 7.27, 1.5, 0.2, f'{page_num} / 15', size=8, color=PRIMARY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# Slide 1 — 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, 0, 0, 13.333, 7.5, PRIMARY)
add_bg(s, 0, 5.0, 13.333, 0.05, GOLD)

# 메인 타이틀
add_text(s, 0.8, 1.5, 11.7, 0.8, 'AI Risk Lighthouse', size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 2.5, 11.7, 0.5, '한국 표준 행동위험-내부통제 자동 감사 플랫폼', size=22, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 3.2, 11.7, 0.4, '2027-2030 AI 메가트렌드를 활용한 부서 핵심 IP 사업', size=14, color=GOLD, align=PP_ALIGN.CENTER)

# 4 패러다임 박스
for i, (txt, x) in enumerate([('양자모델', 1.8), ('LLM', 4.6), ('피지컬 AI', 7.4), ('생성형 AI', 10.2)]):
    add_bg(s, x, 4.3, 2.3, 0.6, GOLD)
    add_text(s, x, 4.4, 2.3, 0.4, txt, size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# 하단 정보
add_text(s, 0.8, 5.5, 11.7, 0.4, '아이티센코어  -  리스크모니터링-행동위험분석 부서', size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 6.0, 11.7, 0.3, '2026년 6월', size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 6.8, 11.7, 0.3, 'Multi-AI Orchestration 기반 신사업 제안', size=10, color=GOLD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# Slide 2 — Executive Summary
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 2, 'Executive Summary', '한 페이지 핵심 — 무엇-왜-언제-얼마')

# 4 박스
boxes = [
    ('[TGT] 무엇 (WHAT)', '한국 최초 AI Risk Lighthouse — 회사의 AI-내부통제-행동위험을 자동 감사-점수화하는 표준 플랫폼',
     'Google Lighthouse 가 웹페이지 점수 매기듯, 8 카테고리 가중 점수 산정'),
    ('[TIP] 왜 (WHY)', 'EU AI Act 2027-금감원 AI 거버넌스 2026-ISO 42001 2026-27 의무화로 모든 한국 기업이 AI 위험 정량화 필수',
     '한국 표준 선점 = 영구 매출 + 부서 IP 확보'),
    ('[TIME] 언제 (WHEN)', '2026-Q4 IP 확보 -> 2027-Q3 한국 표준 등록 -> 2028+ 5,000社 의무 도입',
     '3년 사업화 로드맵'),
    (' 얼마 (HOW MUCH)', '초기 자본: 2-5억 (부서 자체 추진)\n2028+ 매출: 5,000社 × 1억 = 5,000억 영구',
     '글로벌 수출 2029+: K-Standard'),
]
for i, (title, body, sub) in enumerate(boxes):
    x = 0.3 + (i % 2) * 6.5
    y = 1.0 + (i // 2) * 2.9
    add_bg(s, x, y, 6.2, 2.7, WHITE, outline=PRIMARY)
    add_bg(s, x, y, 6.2, 0.5, PRIMARY)
    add_text(s, x + 0.2, y + 0.1, 6.0, 0.35, title, size=16, bold=True, color=WHITE)
    add_text(s, x + 0.25, y + 0.7, 5.9, 1.3, body, size=12, color=GRAY)
    add_text(s, x + 0.25, y + 2.1, 5.9, 0.5, sub, size=10, color=ACCENT, bold=True)

add_footer(s, 2)


# ════════════════════════════════════════════════════════════
# Slide 3 — 2027-2030 메가트렌드
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 3, '2027-2030 AI 메가트렌드', '검증된 시장 데이터 + 법규 의무화 일정')

# 4 분면
trends = [
    (' 산업 메타버스', '$228.58B by 2029', 'CAGR 51.5%\nDigital Twin + AR/VR + AI + IoT\n2027: 40% 대기업 매출', ACCENT),
    (' Emotion AI', '$42.9B by 2027', 'CAGR 12.8%\n행동 + 감정 = 새 차원\n카지노-금융 직판', PRIMARY),
    (' AI Cybersecurity', '$146.5B by 2034', 'Agentic AI 폭증\n신 공격면 (AI Workload)\n2034 6배 성장', GOLD),
    ('[SIG] Insider Risk Mgmt', 'UEBA -> IRM 통합', 'Gartner 2025 재분류\nSIEM 통합 추세\n부서 핵심 IP 직결', ACCENT),
]
for i, (name, size_txt, body, color) in enumerate(trends):
    x = 0.3 + (i % 2) * 6.5
    y = 1.0 + (i // 2) * 2.9
    add_bg(s, x, y, 6.2, 2.7, WHITE, outline=color)
    add_bg(s, x, y, 6.2, 0.5, color)
    add_text(s, x + 0.2, y + 0.1, 6.0, 0.35, name, size=16, bold=True, color=WHITE)
    add_text(s, x + 0.25, y + 0.7, 5.9, 0.4, size_txt, size=22, bold=True, color=color)
    add_text(s, x + 0.25, y + 1.3, 5.9, 1.2, body, size=11, color=GRAY)

add_footer(s, 3)


# ════════════════════════════════════════════════════════════
# Slide 4 — 50개 신기술 카탈로그 요약
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 4, '50개 AI 신기술 카탈로그', '2027-2030 진짜 뜨는 기술 — 16 카테고리')

cats = [
    ('A. 추론-인지 (5)', ['Reasoning Models (o3-R1)', 'Self-Critique / Reflexion', 'Causal AI', 'Chain/Tree/Graph-of-Thought', 'Neurosymbolic AI']),
    ('B. 에이전트 (4)', ['Agentic AI', 'Multi-Agent (CrewAI-LangGraph)', 'Computer-Use', 'MCP 표준 (Anthropic)']),
    ('C. 학습 (3)', ['Mixture of Experts (MoE)', 'State Space Models (Mamba)', 'Test-Time Training']),
    ('D. 생성형 (3)', ['Text-to-Video (Sora-Veo)', 'Code Agents (Cursor-Devin)', 'Multimodal Native']),
    ('E. 피지컬 AI (4)', ['World Models (Cosmos-V-JEPA)', 'VLA (RT-2-Pi0-GR00T)', 'Embodied AI / Humanoid', 'Sim-to-Real']),
    ('F. 양자 AI (2)', ['Quantum ML (IBM Quantum)', 'Variational Quantum Circuits']),
    ('G. 검색-메모리 (2)', ['GraphRAG (Microsoft)', 'Memory (MemGPT-Letta)']),
    ('H. AI 보안 (3)', ['Prompt Injection Defense', 'Mechanistic Interpretability', 'Deepfake / C2PA']),
    ('I. 프라이버시 (3)', ['Federated Learning', 'Confidential Computing', 'Synthetic Data']),
    ('J. 인지-정서 (1)', ['Affective Computing']),
    ('K. 보안 신영역 (5)', ['Adversarial ML / AI Workload / NHI / CSMA / DSPM']),
    ('L. 인증-생체 (3)', ['Behavioral Biometrics / Continuous Auth / Passkeys']),
    ('M. 데이터-학습 (4)', ['DPO / RLHF / Constitutional / LoRA']),
    ('N. 검색-인프라 (3)', ['Vector DB / HyDE / Long Context']),
    ('O. AI 거버넌스 (2)', ['Governance / Bias Detection']),
    ('P. 도메인-신영역 (3)', ['Domain FM / AI Search / Ambient Intelligence']),
]

# 4 cols x 4 rows
for i, (cat, items) in enumerate(cats):
    col = i % 4
    row = i // 4
    x = 0.25 + col * 3.27
    y = 1.0 + row * 1.55
    add_bg(s, x, y, 3.15, 1.45, WHITE, outline=PRIMARY)
    add_bg(s, x, y, 3.15, 0.3, LIGHT)
    add_text(s, x + 0.1, y + 0.05, 3.05, 0.25, cat, size=10, bold=True, color=PRIMARY)
    txt = '\n'.join('- ' + it for it in items[:3])
    add_text(s, x + 0.15, y + 0.35, 3.0, 1.05, txt, size=8, color=GRAY)

add_footer(s, 4)


# ════════════════════════════════════════════════════════════
# Slide 5 — 법규-의무화 일정
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 5, '법규-의무화 일정 (2026-2030)', '회사가 무조건 사야 하는 영역')

headers = ['법규-표준', '발효 일정', '영향', '한국 적용']
rows = [
    ['NIST PQC 마이그레이션', '2028-2030 기업', '모든 암호화 전환', 'KISA 2026 가이드'],
    ['EU AI Act', '2026-08 GPAI / 2027-08 고위험', '한국 EU 수출 의무', '금감원 AI 거버넌스 2026 발효'],
    ['ISO 42001 (AI 관리시스템)', '2026-2027 한국 도입', '모든 AI 사용 기업', '컴플라이언스 의무'],
    ['중대재해처벌법', '2024 50인+ / 2027 강화', '50인+ 기업 의무', '산안-건설-제조'],
    ['개인정보 AI 영향평가', '2026-2027 의무화', '모든 AI 도입 기업', '개인정보위'],
    ['망분리 완화-재정의', '2026-2028 재정의', '금융-공공 전환', '금감원'],
    ['금감원 AI 거버넌스', '2026 발효', '모든 금융사', '한국'],
    ['가명정보-익명정보 활용', '2026 확대', '데이터 활용 기업', '개인정보위'],
]
add_table(s, 0.5, 1.1, 12.3, 5.3, headers, rows,
          first_col_bold=True, col_widths=[3, 3, 4, 3], font_size=11)

add_text(s, 0.5, 6.7, 12.3, 0.4, '[WARN] 5,000+ 한국 대기업 모두 영향. AI Risk Lighthouse 가 통합 대응 — 1社 1억 × 5,000 = 5,000억 시장',
         size=12, bold=True, color=ACCENT)

add_footer(s, 5)


# ════════════════════════════════════════════════════════════
# Slide 6 — ITCEN CORE 자산 매핑
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 6, 'ITCEN CORE 자산 매핑', '기존 1위 솔루션 + 부서 IP = 진입 무기')

# 좌측: ITCEN CORE 자산
add_text(s, 0.3, 1.0, 6, 0.4, ' ITCEN CORE 자산', size=18, bold=True, color=PRIMARY)
core_assets = [
    ('내부회계 1위', '국내 점유율 1위 - 한국 수천社 고객 base - 회계 데이터 자산'),
    ('EPM-Compliance 1위', '경영관리-컴플라이언스 솔루션 국내 1위'),
    ('건설 ERP 1위', '국내 건설업 시장 점유 1위 - 중대재해법 채널 직결'),
    ('금융-카지노 VMS', '강원랜드-외국인전용카지노 9社 독점 - 글로벌 진출 가능'),
    ('AI 지능형 CCTV', '영상 분석 솔루션 보유 - World Models-VLA 결합 가능'),
    ('디지털 트윈', '제조-건설-인프라 시뮬레이션 자산'),
    ('ITO 서비스', 'IT 아웃소싱 운영 - MSP 사업 base'),
    ('ESG-GRC', '거버넌스-리스크-컴플라이언스 통합 솔루션'),
]
for i, (name, desc) in enumerate(core_assets):
    y = 1.5 + i * 0.7
    add_bg(s, 0.3, y, 6.2, 0.6, LIGHT, outline=PRIMARY)
    add_text(s, 0.5, y + 0.08, 2.5, 0.45, name, size=12, bold=True, color=PRIMARY)
    add_text(s, 3.0, y + 0.13, 3.4, 0.4, desc, size=9, color=GRAY)

# 우측: 부서 IP
add_text(s, 6.8, 1.0, 6, 0.4, '[TGT] 부서 IP (리스크모니터링-행동위험분석)', size=18, bold=True, color=ACCENT)
dept_ip = [
    ('UEBA (행동 분석)', '직원-고객 행동 패턴 학습 - 부정-횡령 사전 탐지 IP'),
    ('양자모델 적용', 'IBM Quantum 양자 최적화로 금융 부정거래 분석'),
    ('LLM 위험 분석', 'Reasoning + Agentic AI 자율 Risk Officer'),
    ('피지컬 AI 융합', 'CCTV-VMS-디지털트윈 + 행동 패턴 = 통합 관제'),
    ('생성형 AI 활용', '학습 데이터-시뮬레이션 자동 생성 - Deepfake 방어'),
    ('Self-Critique 노하우', '다중 AI 합의 위험점수 - 차세대 표준'),
    ('Causal 추론', '단순 상관관계 X -> 진짜 원인 파악'),
    ('Constitutional AI', 'SOP-법규 헌법화 - 자동 규제 준수'),
]
for i, (name, desc) in enumerate(dept_ip):
    y = 1.5 + i * 0.7
    add_bg(s, 6.8, y, 6.2, 0.6, RGBColor(0xFC, 0xE4, 0xE4), outline=ACCENT)
    add_text(s, 7.0, y + 0.08, 2.5, 0.45, name, size=12, bold=True, color=ACCENT)
    add_text(s, 9.5, y + 0.13, 3.4, 0.4, desc, size=9, color=GRAY)

add_footer(s, 6)


# ════════════════════════════════════════════════════════════
# Slide 7 — 부서 현재 위치 (SWOT)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 7, '부서 현재 위치 — SWOT 분석', '강점-약점-기회-위협')

swot = [
    (' Strengths (강점)', PRIMARY, [
        '한국에서 유일한 양자-LLM-피지컬-생성형 AI 4 패러다임 동시 보유',
        '리스크모니터링-행동위험분석 도메인 깊이',
        'ITCEN CORE 1위 채널 (내부회계-건설ERP-VMS)',
        '한국 수천社 고객 base 즉시 활용',
    ]),
    ('[WARN] Weaknesses (약점)', GOLD, [
        'R&D 자체 개발 어려움 (글로벌 OEM 의존 필요)',
        'AI 표준화-인증 경험 부족',
        '글로벌 진출 채널 부족 (한국 위주)',
        '브랜드 인지도 (글로벌 시장에서)',
    ]),
    ('[GO] Opportunities (기회)', ACCENT, [
        'EU AI Act-NIST PQC-중대재해법 의무화 = 폭발적 시장',
        'Gartner UEBA -> IRM 통합 = 부서 IP 직결',
        '한국 K-Standard lobby 선점 가능',
        '글로벌 메가트렌드 $400B+ (2029 합계)',
    ]),
    (' Threats (위협)', RGBColor(0x80, 0x00, 0x80), [
        '글로벌 빅테크 (Microsoft-NVIDIA-Anthropic) 한국 진입',
        '한국 사이버보안 회사 30+개 경쟁 (안랩-SK쉴더스)',
        '인력 확보-유지 (AI 박사급)',
        '법규 변화 속도 (예측 어려움)',
    ]),
]
for i, (title, color, items) in enumerate(swot):
    x = 0.3 + (i % 2) * 6.5
    y = 1.0 + (i // 2) * 2.9
    add_bg(s, x, y, 6.2, 2.7, WHITE, outline=color)
    add_bg(s, x, y, 6.2, 0.5, color)
    add_text(s, x + 0.2, y + 0.1, 6.0, 0.35, title, size=16, bold=True, color=WHITE)
    txt = '\n'.join('• ' + it for it in items)
    add_text(s, x + 0.25, y + 0.7, 5.9, 1.9, txt, size=11, color=GRAY)

add_footer(s, 7)


# ════════════════════════════════════════════════════════════
# Slide 8 — 신사업 방향 Top 5
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 8, '신사업 Top 5 — AI 4 패러다임 융합', '부서-ITCEN CORE 시너지-0원 시작 교집합')

headers = ['#', '신사업', '핵심 기술', '0원 시작', '2028+ 매출']
rows = [
    ['', 'AI Risk Lighthouse (한국 표준)', 'Self-Critique + Causal + 8 카테고리', '부서 IP만', '5,000억 (5,000社)'],
    ['', 'Agentic Risk Officer (24/7 자율)', 'Agentic AI + LLM + Multi-Agent', 'API 사용량만', '월구독 + 종량제'],
    ['', 'UEBA 위험점수 Bureau', 'GraphRAG + Behavioral Biometrics', '기존 GRC + 행동 IP', '점수 라이선스 영구'],
    ['4', 'Industrial Metaverse 1호 SI', 'World Models + VLA + 디지털트윈', 'NVIDIA 파트너십', '한국 SI 독점'],
    ['5', '카지노 글로벌 + Emotion AI OEM', 'Affective + VLA + Multimodal', 'OEM 계약비', '글로벌 100社 진출'],
]
add_table(s, 0.5, 1.1, 12.3, 4.8, headers, rows,
          first_col_bold=True, col_widths=[0.7, 3, 4, 2.3, 3], font_size=11)

# 하단 핵심 메시지
add_bg(s, 0.5, 6.3, 12.3, 0.8, LIGHT)
add_text(s, 0.7, 6.4, 12, 0.3, '[TIP] 5개 신사업 모두 R&D 자체 개발 없음 (OEM-통합-SaaS-라이선스 모델)',
         size=13, bold=True, color=PRIMARY)
add_text(s, 0.7, 6.75, 12, 0.3, '   -> 초기 자본 2-5억 / 부서 직접 추진 / 3년 내 영구 매출 5,000억+',
         size=11, color=ACCENT)

add_footer(s, 8)


# ════════════════════════════════════════════════════════════
# Slide 9 — 신사업 #1 AI Risk Lighthouse 상세
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 9, '신사업 #1 — AI Risk Lighthouse', '한국 표준 행동위험-내부통제 자동 감사 플랫폼')

# 좌측: 8 카테고리
add_text(s, 0.3, 1.0, 6, 0.4, '[SIG] 8 검증 카테고리 (가중 점수)', size=15, bold=True, color=PRIMARY)
cats = [
    ('Self-Critique', '15%', '2단계 검증 - confidence 점수'),
    ('Causal AI', '15%', '인과 설명 - DoWhy 적용'),
    ('Behavioral Coverage', '15%', 'UEBA-VMS-CCTV-결재-감정 통합'),
    ('Interpretability', '12%', 'EU AI Act 의무 - SHAP-LIME'),
    ('Privacy (PET)', '12%', 'Federated-Confidential-동형암호'),
    ('Compliance', '10%', '한국 법규-EU AI Act 자동 추적'),
    ('Quality (FP/FN)', '11%', '거짓양성-음성 자동 모니터'),
    ('Self-Improvement', '10%', '실패 학습 - Reflexion 루프'),
]
for i, (name, weight, desc) in enumerate(cats):
    y = 1.5 + i * 0.6
    add_bg(s, 0.3, y, 6.2, 0.5, WHITE, outline=PRIMARY)
    add_text(s, 0.4, y + 0.08, 2.5, 0.35, name, size=11, bold=True, color=PRIMARY)
    add_text(s, 2.9, y + 0.08, 0.8, 0.35, weight, size=11, bold=True, color=ACCENT)
    add_text(s, 3.7, y + 0.1, 2.7, 0.3, desc, size=9, color=GRAY)

# 우측: 사업 모델
add_text(s, 6.8, 1.0, 6, 0.4, ' 사업 모델 (4단계 수익)', size=15, bold=True, color=ACCENT)

biz = [
    ('1. 무료 진단 (Lead Gen)', '잠재 고객 위험점수 무료 측정 -> 유료 컨설팅 깔때기'),
    ('2. 컨설팅 (1社 5천만~3억)', '8 카테고리 audit + 개선 권고 + 인증 대행'),
    ('3. SaaS 구독 (월 200만~)', '점수 자동 업데이트 - 분기 audit - 1,000+社'),
    ('4. 한국 표준 라이선스 (영구)', 'KISA-금감원 등록 후 모든 한국 기업 의무'),
]
for i, (name, desc) in enumerate(biz):
    y = 1.5 + i * 1.05
    add_bg(s, 6.8, y, 6.2, 0.95, RGBColor(0xFC, 0xE4, 0xE4), outline=ACCENT)
    add_text(s, 7.0, y + 0.08, 5.8, 0.35, name, size=13, bold=True, color=ACCENT)
    add_text(s, 7.0, y + 0.45, 5.8, 0.5, desc, size=10, color=GRAY)

add_footer(s, 9)


# ════════════════════════════════════════════════════════════
# Slide 10 — 신사업 #2 Agentic Risk Officer
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 10, '신사업 #2 — Agentic Risk Officer', '24/7 자율 위험 분석 AI 에이전트')

# 좌측: 아키텍처
add_text(s, 0.3, 1.0, 6.2, 0.4, ' 아키텍처 (Plan-Act-Reflect 루프)', size=15, bold=True, color=PRIMARY)
arch = [
    ('① Plan (LLM)', 'Reasoning Models (o3-Claude Extended)\n위험 사건 -> 분석 계획'),
    ('② Act (도구 호출)', 'MCP 표준 + Computer-Use\n내부회계-VMS-CCTV 데이터 자율 조회'),
    ('③ Reflect (Self-Critique)', 'Haiku-Validator + Reflexion\n분석 결과 자기 비판 -> 재시도'),
    ('④ Multi-Agent 합의', '회계 + UEBA + VMS + GRC 4 에이전트\n토론 -> 합의 점수'),
    ('⑤ Memory + Learn', 'MemGPT-Letta + learn skill\n과거 사례 영구 학습'),
]
for i, (name, desc) in enumerate(arch):
    y = 1.5 + i * 1.0
    add_bg(s, 0.3, y, 6.2, 0.9, LIGHT)
    add_text(s, 0.5, y + 0.1, 2.5, 0.35, name, size=12, bold=True, color=PRIMARY)
    add_text(s, 3.0, y + 0.15, 3.3, 0.65, desc, size=9, color=GRAY)

# 우측: 적용 사례 + ROI
add_text(s, 6.8, 1.0, 6.2, 0.4, '[TIP] 적용 사례 + ROI', size=15, bold=True, color=ACCENT)

cases = [
    ('내부회계 자율 감사', '분식-횡령 24/7 모니터링 (사람 1명 = 100배 효율)\n예상 ROI: 1社 5천만 × 1,000社 = 500억'),
    ('카지노 부정 자율 분석', '강원랜드-외국인전용 9社 24/7 분석\n예상 ROI: 1社 1억 × 9 = 9억 + 글로벌'),
    ('금융 보안이벤트 자율 대응', 'AI 자율 조사-격리-증거보존\n예상 ROI: 보험-금감원 직판'),
    ('컴플라이언스 자율 점검', '규제 변경 자동 추적-영향 분석-보고\n예상 ROI: 모든 GRC 고객 추가판매'),
]
for i, (name, desc) in enumerate(cases):
    y = 1.5 + i * 1.2
    add_bg(s, 6.8, y, 6.2, 1.1, WHITE, outline=ACCENT)
    add_bg(s, 6.8, y, 6.2, 0.35, ACCENT)
    add_text(s, 6.95, y + 0.05, 6.0, 0.3, name, size=12, bold=True, color=WHITE)
    add_text(s, 7.0, y + 0.45, 5.9, 0.6, desc, size=9, color=GRAY)

add_footer(s, 10)


# ════════════════════════════════════════════════════════════
# Slide 11 — 신사업 #3 UEBA 위험점수 Bureau
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 11, '신사업 #3 — UEBA 위험점수 Bureau', '한국 행동신용평가소 (영구 라이선스 IP)')

add_text(s, 0.3, 1.0, 12.7, 0.4, '[TGT] 컨셉: 신용평가사(KCB-NICE)처럼 직원-고객 행동위험을 점수화 -> 표준 라이선스',
         size=14, bold=True, color=PRIMARY)

# 4 시장
markets = [
    (' 보험 시장', '사이버보험 인수심사 자동화\n생명-재산보험 위험점수\n-> 보험사 30+개 직판', PRIMARY),
    (' HR 시장', '채용-승진 위험 사전 평가\n헤드헌팅 후보자 검증\n-> HR 컨설팅 협업', GOLD),
    (' IAM 시장', 'Zero Trust 인증 + 행동점수\n금융권 차세대 인증\n-> 금감원 표준', ACCENT),
    ('[STAT] 신용 시장', '대안 신용평가 (행동 기반)\nKCB-NICE 보완\n-> 핀테크 직판', RGBColor(0x00, 0x70, 0x40)),
]
for i, (name, body, color) in enumerate(markets):
    x = 0.3 + (i % 2) * 6.5
    y = 1.6 + (i // 2) * 1.9
    add_bg(s, x, y, 6.2, 1.7, WHITE, outline=color)
    add_bg(s, x, y, 6.2, 0.45, color)
    add_text(s, x + 0.2, y + 0.07, 6.0, 0.35, name, size=14, bold=True, color=WHITE)
    add_text(s, x + 0.25, y + 0.6, 5.9, 1.1, body, size=11, color=GRAY)

# 하단: 핵심 차별화
add_bg(s, 0.3, 5.5, 12.7, 1.55, LIGHT)
add_text(s, 0.5, 5.6, 12.4, 0.35, ' 핵심 차별화 — 점수 IP 영구 보유 (한국 1개사만 가능)',
         size=14, bold=True, color=PRIMARY)
diff = [
    '① GraphRAG 기반 행동패턴 그래프 = 부서 핵심 IP, 다른 회사 모방 불가',
    '② Behavioral Biometrics (타이핑-마우스-걸음) = 부서 행동분석 자연 확장',
    '③ Causal AI 인과 추론 = 단순 상관관계 X, 진짜 원인 파악',
    '④ Federated Learning = 여러 회사 데이터 안 모으고 학습 (개인정보 의무 100% 충족)',
]
for i, d in enumerate(diff):
    add_text(s, 0.5, 5.95 + i * 0.28, 12.4, 0.25, d, size=10, color=GRAY)

add_footer(s, 11)


# ════════════════════════════════════════════════════════════
# Slide 12 — 신사업 #4 Industrial Metaverse 1호 SI
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 12, '신사업 #4 — Industrial Metaverse 한국 1호 SI', 'NVIDIA Cosmos-MS Mesh 한국 SI 독점')

# 좌측: 시장 + 기회
add_text(s, 0.3, 1.0, 6.2, 0.4, '[UP] 시장 + 기회', size=15, bold=True, color=PRIMARY)
add_bg(s, 0.3, 1.5, 6.2, 5.5, WHITE, outline=PRIMARY)

mkt = [
    ' 글로벌 시장: $228.58B by 2029 (CAGR 51.5%)',
    ' 한국 점유 가능: $20-30억 (SI 독점)',
    '[TGT] Gartner 2027: 40%+ 대기업 산업 메타버스 매출',
    '[FAST] 핵심 기술: Digital Twin + AR/VR + AI + IoT + 5G',
    '',
    ' ITCEN CORE 자산 활용',
    '  • 디지털트윈 솔루션 (기존 보유)',
    '  • 건설 ERP 1위 (현장 디지털트윈 자연 결합)',
    '  • AI 지능형 CCTV (실시간 행동 캡처)',
    '  • 카지노 VMS (메타버스 출입관리)',
    '',
    ' 파트너십 (자체 R&D 없음)',
    '  • NVIDIA Cosmos (World Foundation Model)',
    '  • Microsoft Mesh (협업 메타버스)',
    '  • Siemens Industrial Metaverse',
    '  • Unity Industrial Collection',
]
for i, line in enumerate(mkt):
    add_text(s, 0.5, 1.65 + i * 0.32, 6, 0.3, line, size=10, color=GRAY)

# 우측: 적용 산업
add_text(s, 6.8, 1.0, 6.2, 0.4, ' 적용 산업 (한국 SI 1호 우선권)', size=15, bold=True, color=ACCENT)

industries = [
    (' 건설-제조', '건설현장 디지털트윈 + 작업안전\n행동AI + AR 작업자 안내 (산안법)'),
    (' 금융 전산실', '디지털트윈 보안관제 + UEBA\n양자센싱 + 행동AI 통합'),
    (' 카지노', 'VMS 메타버스 출입관리\n글로벌 9社 -> 동남아-일본-마카오'),
    (' 물류-항만', '항만-창고 디지털트윈\n자율 순찰 로봇 + VLA'),
    (' 스마트시티', '도시 디지털트윈 통합 관제\n행안부-환경부 입찰'),
]
for i, (name, desc) in enumerate(industries):
    y = 1.5 + i * 1.05
    add_bg(s, 6.8, y, 6.2, 0.95, RGBColor(0xFC, 0xE4, 0xE4), outline=ACCENT)
    add_text(s, 7.0, y + 0.08, 5.8, 0.35, name, size=12, bold=True, color=ACCENT)
    add_text(s, 7.0, y + 0.45, 5.8, 0.5, desc, size=9, color=GRAY)

add_footer(s, 12)


# ════════════════════════════════════════════════════════════
# Slide 13 — 신사업 #5 카지노 글로벌 + Emotion AI
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 13, '신사업 #5 — 카지노 글로벌 + Emotion AI OEM', '한국 9社 base -> 글로벌 100社 진출')

# 4 분면
sections = [
    ('[TGT] 차별화', PRIMARY, [
        'ITCEN CORE 가 강원랜드-외국인전용 9社 보유 = 한국 카지노 VMS 독점',
        'Emotion AI OEM (Affectiva-Hume-Realeyes) 결합',
        'VLA + Multimodal 로 딜러-플레이어 행동 추론',
        'Affective Computing 으로 도박중독 조기 경보',
    ]),
    (' 매출 모델', GOLD, [
        '국내 9社 × 5억 = 45억 (1차)',
        '동남아 100社 (마카오-필리핀-싱가포르) = 500억',
        '일본 IR (Integrated Resort) 2개 = 100억',
        '미국-유럽 1,000+개 = 1,500억 잠재',
    ]),
    (' 기술 스택', ACCENT, [
        'Emotion AI: Affectiva-Hume-Realeyes OEM',
        'VLA: NVIDIA GR00T / RT-2',
        'Multimodal: GPT-4o-Gemini 2.5',
        'Self-Critique: 1차 탐지 -> 2차 재검토',
    ]),
    ('[UP] 로드맵', RGBColor(0x00, 0x70, 0x40), [
        '2026 Q4: OEM 계약 + 9社 베타',
        '2027 Q1: 한국 9社 본격 도입',
        '2027 Q3: 동남아 진출 (마카오 우선)',
        '2028+: 일본 IR + 미국-EU 확장',
    ]),
]
for i, (name, color, items) in enumerate(sections):
    x = 0.3 + (i % 2) * 6.5
    y = 1.0 + (i // 2) * 2.9
    add_bg(s, x, y, 6.2, 2.7, WHITE, outline=color)
    add_bg(s, x, y, 6.2, 0.5, color)
    add_text(s, x + 0.2, y + 0.1, 6.0, 0.35, name, size=15, bold=True, color=WHITE)
    txt = '\n'.join('• ' + it for it in items)
    add_text(s, x + 0.25, y + 0.7, 5.9, 1.9, txt, size=10, color=GRAY)

add_footer(s, 13)


# ════════════════════════════════════════════════════════════
# Slide 14 — 3년 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 14, '3년 사업화 로드맵', '2026 Q4 -> 2029 글로벌')

# 타임라인
add_bg(s, 0.5, 1.5, 12.3, 0.15, PRIMARY)
quarters = ['2026Q4', '2027Q1', '2027Q2', '2027Q3', '2027Q4', '2028H1', '2028H2', '2029+']
for i, q in enumerate(quarters):
    x = 0.5 + i * (12.3 / len(quarters))
    add_bg(s, x, 1.4, 0.15, 0.35, ACCENT)
    add_text(s, x - 0.4, 1.85, 1.0, 0.3, q, size=10, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# 단계별 마일스톤
phases = [
    ('Phase 1 (2026 Q4)', 'IP 확보 단계', PRIMARY, [
        '부서 SOP-법규 헌법화 (Constitutional AI)',
        'Lighthouse 8 카테고리 + 점수 모델 IP 확보',
        'NVIDIA-Microsoft 파트너십 등록',
        'IBM Quantum Network 무료 가입',
    ]),
    ('Phase 2 (2027 Q1-Q3)', '베타 -> 표준', GOLD, [
        'ITCEN CORE 고객 100社 베타 진단 (무료)',
        'KISA-금감원-개인정보위 표준 lobby',
        '한국 K-AI Risk Standard 초안 발표',
        '카지노 OEM 9社 본격 도입',
    ]),
    ('Phase 3 (2027 Q4-2028)', '본격 확산', ACCENT, [
        '한국 K-Standard 공식 등록 + 인증 사업',
        'SaaS 구독 1,000社 (월 200만~)',
        'Industrial Metaverse SI 한국 독점',
        '연간 매출 500억 -> 1,000억',
    ]),
    ('Phase 4 (2029+)', '글로벌', RGBColor(0x00, 0x70, 0x40), [
        'K-Standard -> 동남아-중동-EU 수출',
        '카지노 글로벌 100社 확장',
        '연간 매출 5,000억 영구',
        'ITCEN CORE IPO 동력',
    ]),
]
for i, (title, phase, color, items) in enumerate(phases):
    x = 0.3 + i * 3.27
    y = 2.5
    add_bg(s, x, y, 3.15, 4.5, WHITE, outline=color)
    add_bg(s, x, y, 3.15, 0.7, color)
    add_text(s, x + 0.1, y + 0.05, 3.0, 0.3, title, size=11, bold=True, color=WHITE)
    add_text(s, x + 0.1, y + 0.38, 3.0, 0.3, phase, size=10, color=LIGHT)
    txt = '\n\n'.join('• ' + it for it in items)
    add_text(s, x + 0.15, y + 0.85, 3.0, 3.5, txt, size=9, color=GRAY)

add_footer(s, 14)


# ════════════════════════════════════════════════════════════
# Slide 15 — 투자-인력-ROI + 의사결정 요청
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_title(s, 15, '투자-인력-ROI + 의사결정 요청', '본부장-임원 승인 필요 사항')

# 좌측: 투자-인력
add_text(s, 0.3, 1.0, 6.2, 0.4, ' 투자-인력 계획', size=15, bold=True, color=PRIMARY)

invest = [
    ('초기 자본 (2026 Q4)', '2-5억', '• 글로벌 OEM 계약 (NVIDIA-Affectiva)\n• Lighthouse IP 개발 (부서 5명 × 6개월)\n• 표준 lobby 컨설팅비'),
    ('운영 자본 (2027)', '연 10-20억', '• 컨설턴트 인력 10명 (사내 + 외부)\n• 베타 운영 + 마케팅\n• KISA 표준 등록 비용'),
    ('확장 자본 (2028+)', '연 30-50억', '• SaaS 인프라 (Azure-AWS)\n• 영업 인력 20명\n• 글로벌 진출 (동남아-일본)'),
]
for i, (name, amt, desc) in enumerate(invest):
    y = 1.5 + i * 1.6
    add_bg(s, 0.3, y, 6.2, 1.5, WHITE, outline=PRIMARY)
    add_bg(s, 0.3, y, 6.2, 0.5, PRIMARY)
    add_text(s, 0.5, y + 0.1, 4, 0.3, name, size=12, bold=True, color=WHITE)
    add_text(s, 4.5, y + 0.1, 2, 0.3, amt, size=14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    add_text(s, 0.55, y + 0.6, 6, 0.85, desc, size=10, color=GRAY)

# 우측: ROI + 의사결정 요청
add_text(s, 6.8, 1.0, 6.2, 0.4, '[STAT] ROI + 의사결정 요청', size=15, bold=True, color=ACCENT)

# ROI 박스
add_bg(s, 6.8, 1.5, 6.2, 2.2, RGBColor(0xFC, 0xE4, 0xE4), outline=ACCENT)
add_text(s, 7.0, 1.6, 5.8, 0.3, ' 예상 ROI (보수적 추정)', size=12, bold=True, color=ACCENT)
roi = [
    '2027: 매출 50억 (베타 + 카지노 OEM)',
    '2028: 매출 500-1,000억 (SaaS + 표준)',
    '2029: 매출 2,000-3,000억 (글로벌)',
    '2030+: 영구 매출 5,000억+ (한국 표준 의무)',
    '',
    'IRR: 300%+ (3년 내)',
    'Payback: 1년 내',
]
add_text(s, 7.0, 1.95, 5.8, 1.7, '\n'.join(roi), size=10, color=GRAY)

# 의사결정 요청 박스
add_bg(s, 6.8, 3.9, 6.2, 3.0, WHITE, outline=ACCENT)
add_bg(s, 6.8, 3.9, 6.2, 0.5, ACCENT)
add_text(s, 6.95, 4.0, 6.0, 0.35, '[FAST] 의사결정 요청 사항', size=14, bold=True, color=WHITE)

asks = [
    '① 초기 자본 2-5억 승인 (2026 Q4)',
    '② 부서 인원 충원 5명 (AI 박사급 2 + 일반 3)',
    '③ NVIDIA-MS-Anthropic 파트너십 권한 위임',
    '④ KISA-금감원 표준 lobby 대외 채널 권한',
    '⑤ ITCEN PNS-CTS 협업 권한 (그룹 시너지)',
    '⑥ 분기별 임원진 보고 채널 확정',
]
for i, ask in enumerate(asks):
    add_text(s, 7.0, 4.55 + i * 0.38, 6.0, 0.3, ask, size=11, color=GRAY, bold=(i == 0))

add_footer(s, 15)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f'\n[OK] PPT 15장 생성 완료')
print(f'경로: {OUT}')
print(f'크기: {os.path.getsize(OUT)} bytes')
print(f'슬라이드 수: {len(prs.slides)}')
