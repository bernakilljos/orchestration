#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Orchestration Kit v1 — Premium PPT Generator v5 (Cyberpunk Theme)
Author: Claude Code
Date: 2026-04-23

Style: Cyberpunk / Futuristic Tech Dark Theme
- 20 slides with neon glow effects
- Terminal-style logs, tech dashboard aesthetic
- Dark backgrounds (navy/black) + neon accents (cyan/pink/gold/green)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import random

# ============================================================================
# COLOR PALETTE — CYBERPUNK
# ============================================================================

CYBER_SPACE = RGBColor(0x0A, 0x0D, 0x14)           # 심해 블랙
CYBER_NIGHT = RGBColor(0x12, 0x16, 0x22)           # 네이비 블랙
CYBER_MIST = RGBColor(0x1C, 0x23, 0x32)            # 연한 네이비 (박스 배경)

NEON_CYAN = RGBColor(0x00, 0xE5, 0xFF)             # 네온 시안
NEON_PINK = RGBColor(0xFF, 0x2E, 0x7A)             # 네온 핑크
NEON_GREEN = RGBColor(0x39, 0xFF, 0x6C)            # 네온 그린
NEON_GOLD = RGBColor(0xFF, 0xD1, 0x4F)             # 앰버 골드
NEON_PURPLE = RGBColor(0xB8, 0x6B, 0xFF)           # 네온 퍼플

TEXT_BRIGHT = RGBColor(0xF0, 0xF2, 0xF8)           # 본문 (밝은 회색)
TEXT_MID = RGBColor(0xA8, 0xB0, 0xC0)              # 보조 텍스트
TEXT_DIM = RGBColor(0x5A, 0x68, 0x7A)              # 라벨·캡션

GRID_LINE = RGBColor(0x2A, 0x33, 0x42)             # 배경 그리드
DIVIDER = RGBColor(0x3A, 0x46, 0x5A)               # 섹션 구분

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_dark_background(slide, color=CYBER_SPACE):
    """Add dark background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_scan_line_decoration(slide, y_inches=0.5, color=NEON_CYAN):
    """Add thin glowing scan line at top"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(y_inches),
        Inches(10), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    line.line.width = Pt(0.5)

def add_breadcrumb_nav(slide, text, x=0.3, y=0.15):
    """Add breadcrumb navigation text (top-left)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6), Inches(0.3))
    tf = tb.text_frame
    tf.text = text
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.font.name = "Consolas"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_DIM
    p.font.bold = False

def add_slide_id(slide, slide_num, total=20):
    """Add slide ID tag (top-right)"""
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(0.15), Inches(1.2), Inches(0.3))
    tf = tb.text_frame
    tf.text = f"[ID:{slide_num:03d}/{total}]"
    p = tf.paragraphs[0]
    p.font.name = "Consolas"
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_DIM
    p.alignment = PP_ALIGN.RIGHT

def add_glow_rectangle(slide, left, top, width, height, color=NEON_CYAN, fill_color=CYBER_MIST):
    """Add rectangle with neon glow effect"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)

    # Add outer glow shadow via XML
    try:
        sp = shape._element
        spPr = sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if spPr is not None:
            from lxml import etree
            odt_ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
            effectLst = spPr.find(f'.//{odt_ns}effectLst')
            if effectLst is None:
                effectLst = etree.SubElement(spPr, f'{odt_ns}effectLst')
            outerShdw = etree.SubElement(effectLst, f'{odt_ns}outerShdw')
            outerShdw.set('blurRad', '200000')
            outerShdw.set('dist', '0')
            outerShdw.set('dir', '2700000')
            outerShdw.set('algn', 'ctr')
            outerShdw.set('rotWithShape', '0')
            srgbClr = etree.SubElement(outerShdw, f'{odt_ns}srgbClr')
            srgbClr.set('val', color.rgb[1:].upper() if hasattr(color, 'rgb') else 'FFFFFF')
            alpha = etree.SubElement(srgbClr, f'{odt_ns}alpha')
            alpha.set('val', '50000')
    except:
        pass

    return shape

def add_large_text(slide, text, x, y, width, height, size=120, color=NEON_CYAN, bold=True, font="Consolas"):
    """Add large prominent text"""
    tb = slide.shapes.add_textbox(x, y, width, height)
    tf = tb.text_frame
    tf.text = text
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return tb

def add_terminal_log(slide, lines, x, y, width, height):
    """Add terminal-style log area"""
    # Background box
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x, y, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1A)
    bg.line.color.rgb = NEON_GREEN
    bg.line.width = Pt(1)

    # Log text
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), width - Inches(0.2), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = NEON_GREEN
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.1

def add_metric_box(slide, label, value, x, y, size_w=1.5, size_h=1.2, accent_color=NEON_CYAN):
    """Add metric display box"""
    # Box (size_w and size_h are already floats in inches)
    box = add_glow_rectangle(slide, x, y, Inches(float(size_w)), Inches(float(size_h)), accent_color, CYBER_MIST)

    # Label
    tb_label = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), Inches(size_w - 0.2), Inches(0.25))
    tf_label = tb_label.text_frame
    tf_label.text = label.upper()
    p_label = tf_label.paragraphs[0]
    p_label.font.name = "Consolas"
    p_label.font.size = Pt(7)
    p_label.font.color.rgb = TEXT_DIM
    p_label.font.bold = True

    # Value
    tb_value = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.5), Inches(size_w - 0.2), Inches(0.6))
    tf_value = tb_value.text_frame
    tf_value.text = value
    p_value = tf_value.paragraphs[0]
    p_value.font.name = "Consolas"
    p_value.font.size = Pt(22)
    p_value.font.bold = True
    p_value.font.color.rgb = accent_color
    p_value.alignment = PP_ALIGN.CENTER

def add_body_text(slide, text, x, y, width, height):
    """Add body text block"""
    tb = slide.shapes.add_textbox(x, y, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_BRIGHT
    p.line_spacing = 1.4

# ============================================================================
# SLIDE 1: COVER
# ============================================================================

def slide_01_cover(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_SPACE)
    add_scan_line_decoration(slide, 0.5, NEON_CYAN)

    # Title
    add_large_text(slide, "ORCHESTRATION.KIT", Inches(0.5), Inches(2.5), Inches(9), Inches(1.5), size=120, color=NEON_CYAN, font="Consolas")

    # Subtitle line
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.3))
    tf = tb.text_frame
    tf.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
    p = tf.paragraphs[0]
    p.font.name = "Consolas"
    p.font.size = Pt(10)
    p.font.color.rgb = NEON_CYAN
    p.alignment = PP_ALIGN.CENTER

    # Version info
    add_large_text(slide, "v1.0 · 2026-04-23", Inches(0.5), Inches(4.6), Inches(9), Inches(0.4), size=16, color=TEXT_MID, bold=False, font="Consolas")

    # Bottom meta
    meta_text = "CLAUDE.OPUS.4_7 × CODEX × HAIKU.4_5 × GEMINI.FLASH"
    tb_meta = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf_meta = tb_meta.text_frame
    tf_meta.text = meta_text
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.font.name = "Consolas"
    p_meta.font.size = Pt(10)
    p_meta.font.color.rgb = TEXT_DIM
    p_meta.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: TIMELINE
# ============================================================================

def slide_02_timeline(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.01 > TIMELINE")
    add_slide_id(slide, 2)

    # Metrics row
    add_metric_box(slide, "Days Elapsed", "14", Inches(0.5), Inches(0.7), accent_color=NEON_CYAN)
    add_metric_box(slide, "Commits", "3,500+", Inches(2.3), Inches(0.7), accent_color=NEON_CYAN)
    add_metric_box(slide, "Upgrade Events", "1", Inches(4.1), Inches(0.7), accent_color=NEON_PINK)

    # Timeline bar
    timeline_y = Inches(2.2)
    timeline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), timeline_y,
        Inches(9), Inches(0.08)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = DIVIDER
    timeline.line.width = Pt(0)

    # Timeline points
    events = [
        (Inches(0.8), "04-10", "KIT.INITIATED", NEON_CYAN),
        (Inches(2.3), "04-19", "14.PLUGINS", NEON_CYAN),
        (Inches(4.0), "04-20", "SNAPSHOT", TEXT_MID),
        (Inches(5.7), "04-22", "OPUS.4.7", NEON_PINK),
        (Inches(7.4), "04-23", "COMPLETE", NEON_GREEN),
    ]

    for x, date, event, color in events:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, timeline_y - Inches(0.15), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = color
        dot.line.width = Pt(1)

        tb = slide.shapes.add_textbox(x - Inches(0.3), timeline_y + Inches(0.25), Inches(0.6), Inches(0.2))
        tf = tb.text_frame
        tf.text = date
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DIM
        p.alignment = PP_ALIGN.CENTER

        tb_event = slide.shapes.add_textbox(x - Inches(0.4), timeline_y + Inches(0.5), Inches(0.8), Inches(0.3))
        tf_event = tb_event.text_frame
        tf_event.text = event
        tf_event.word_wrap = True
        p_event = tf_event.paragraphs[0]
        p_event.font.name = "Consolas"
        p_event.font.size = Pt(7)
        p_event.font.color.rgb = color
        p_event.alignment = PP_ALIGN.CENTER

    # Log feed
    add_terminal_log(slide, [
        "[2026-04-10 08:30] system.initialize completed",
        "[2026-04-19 14:32] plugins.stable = 14 ✓",
        "[2026-04-20 12:00] snapshot.created",
        "[2026-04-22 19:15] opus.4.7.released",
        "[2026-04-23 09:00] phase1+2.absorption.complete",
    ], Inches(0.5), Inches(4.2), Inches(9), Inches(1.8))

# ============================================================================
# SLIDE 3: 14 PLUGINS MATRIX
# ============================================================================

def slide_03_plugins_matrix(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.01 > PLUGINS")
    add_slide_id(slide, 3)

    # Metrics
    add_metric_box(slide, "Stable", "14", Inches(0.5), Inches(0.7), accent_color=NEON_CYAN)
    add_metric_box(slide, "Spec-Only", "10", Inches(2.3), Inches(0.7), accent_color=NEON_PINK)
    add_metric_box(slide, "Total", "25", Inches(4.1), Inches(0.7), accent_color=NEON_GOLD)

    # 4x4 grid (14 visible)
    plugins = [
        "exec_orch", "ai_rag", "bundles_", "mcp_dev",
        "ai_coding", "bundles_", "mcp_github", "ai_pptx",
        "bundles_", "ai_excel", "mcp_figma", "ai_video",
        "ai_audio", "exec_watch", "", ""
    ]

    start_y = Inches(2.0)
    box_size = 0.8
    spacing = 0.15

    for idx, plugin in enumerate(plugins):
        row = idx // 4
        col = idx % 4
        x = Inches(0.6 + col * (box_size + spacing))
        y = start_y + Inches(row * (box_size + spacing))

        if plugin:
            box = add_glow_rectangle(slide, x, y, Inches(box_size), Inches(box_size), NEON_CYAN, CYBER_MIST)
            tb = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(0.15), Inches(box_size - 0.1), Inches(box_size - 0.3))
            tf = tb.text_frame
            tf.text = plugin
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.name = "Consolas"
            p.font.size = Pt(7)
            p.font.color.rgb = NEON_CYAN
            p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 4: EXEC_ORCH ENGINE
# ============================================================================

def slide_04_exec_orch(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.01 > ENGINE")
    add_slide_id(slide, 4)

    # 3 hexagon nodes (Claude, Codex, Gemini)
    center_y = Inches(2.0)
    nodes = [
        (Inches(1.5), "CLAUDE\nOPUS", NEON_CYAN),
        (Inches(4.5), "CODEX\n×4", NEON_PINK),
        (Inches(7.5), "GEMINI\n×2", NEON_GREEN),
    ]

    for x, label, color in nodes:
        # Circle node
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.4), center_y - Inches(0.4), Inches(0.8), Inches(0.8))
        node.fill.solid()
        node.fill.fore_color.rgb = CYBER_MIST
        node.line.color.rgb = color
        node.line.width = Pt(2)

        # Label
        tb = slide.shapes.add_textbox(x - Inches(0.5), center_y - Inches(0.25), Inches(1), Inches(0.5))
        tf = tb.text_frame
        tf.text = label
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Connection lines
    line1 = slide.shapes.add_connector(1, Inches(1.9), center_y, Inches(4.1), center_y)
    line1.line.color.rgb = NEON_CYAN
    line1.line.width = Pt(2)

    line2 = slide.shapes.add_connector(1, Inches(4.9), center_y, Inches(7.1), center_y)
    line2.line.color.rgb = NEON_PINK
    line2.line.width = Pt(2)

    # Terminal log
    add_terminal_log(slide, [
        "[2026-04-23 10:15] orca-dispatch started",
        "[2026-04-23 10:15] codex-auto spawned 4 workers",
        "[2026-04-23 10:15] gemini-auto spawned 2 workers",
        "[2026-04-23 10:17] task-001.md picked by codex-1",
        "[2026-04-23 10:18] task-001.md completed ✓ 2.3s",
    ], Inches(0.5), Inches(3.8), Inches(9), Inches(2))

# ============================================================================
# SLIDE 5: SOURCE OF TRUTH
# ============================================================================

def slide_05_source_of_truth(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.01 > ARCHITECTURE")
    add_slide_id(slide, 5)

    # Left box: plugins/
    left_box = add_glow_rectangle(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.5), NEON_CYAN)
    tb_left = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(3.1), Inches(0.3))
    tf_left = tb_left.text_frame
    tf_left.text = "SOURCE OF TRUTH"
    p_left = tf_left.paragraphs[0]
    p_left.font.name = "Consolas"
    p_left.font.size = Pt(10)
    p_left.font.bold = True
    p_left.font.color.rgb = NEON_CYAN
    p_left.alignment = PP_ALIGN.CENTER

    tb_plugins = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(3.1), Inches(1.5))
    tf_plugins = tb_plugins.text_frame
    tf_plugins.word_wrap = True
    tf_plugins.text = "plugins/\n\n14 stable\n10 spec-only\n\nAll plugins\nencapsulated here"
    for p in tf_plugins.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_BRIGHT
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    # Arrow
    arrow = slide.shapes.add_connector(1, Inches(4.5), Inches(3.25), Inches(5.5), Inches(3.25))
    arrow.line.color.rgb = NEON_CYAN
    arrow.line.width = Pt(2)

    # Add arrow head manually (add shape at end)
    arrow_head = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.3), Inches(3.1), Inches(0.4), Inches(0.3))
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = NEON_CYAN
    arrow_head.line.color.rgb = NEON_CYAN
    arrow_head.line.width = Pt(0.5)

    # Label on arrow
    tb_sync = slide.shapes.add_textbox(Inches(4.6), Inches(2.7), Inches(1.3), Inches(0.3))
    tf_sync = tb_sync.text_frame
    tf_sync.text = "sync →"
    p_sync = tf_sync.paragraphs[0]
    p_sync.font.name = "Consolas"
    p_sync.font.size = Pt(8)
    p_sync.font.color.rgb = TEXT_MID
    p_sync.alignment = PP_ALIGN.CENTER

    # Right box: .claude/
    right_box = add_glow_rectangle(slide, Inches(6.2), Inches(2.0), Inches(3.5), Inches(2.5), DIVIDER)
    tb_right = slide.shapes.add_textbox(Inches(6.4), Inches(2.2), Inches(3.1), Inches(0.3))
    tf_right = tb_right.text_frame
    tf_right.text = "DERIVED"
    p_right = tf_right.paragraphs[0]
    p_right.font.name = "Consolas"
    p_right.font.size = Pt(10)
    p_right.font.bold = True
    p_right.font.color.rgb = TEXT_MID
    p_right.alignment = PP_ALIGN.CENTER

    tb_derived = slide.shapes.add_textbox(Inches(6.4), Inches(2.8), Inches(3.1), Inches(1.5))
    tf_derived = tb_derived.text_frame
    tf_derived.word_wrap = True
    tf_derived.text = ".claude/\n\ncommands/\nskills/\nscripts/\n\nAuto-generated"
    for p in tf_derived.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MID
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    # Rules
    rules = [
        "✓ Edit in plugins/ only",
        "✓ .claude/ auto-regenerated",
        "✓ No drift (validate-plugin-schema)",
        "✓ Monthly audit",
    ]

    for i, rule in enumerate(rules):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.0 + i * 0.3), Inches(8.4), Inches(0.25))
        tf = tb.text_frame
        tf.text = rule
        p = tf.paragraphs[0]
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(10)
        p.font.color.rgb = NEON_CYAN if i % 2 == 0 else TEXT_BRIGHT

# ============================================================================
# SLIDE 6: ROADMAP STAR MAP
# ============================================================================

def slide_06_roadmap(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_SPACE)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.01 > ROADMAP")
    add_slide_id(slide, 6)

    # Title
    add_large_text(slide, "PHASE CONSTELLATION", Inches(0.5), Inches(0.8), Inches(9), Inches(0.4), size=28, color=NEON_CYAN, font="Malgun Gothic")

    # Phase boxes
    phases = [
        (Inches(0.8), "PHASE 1\n14 stable\n2,061 lines", NEON_CYAN),
        (Inches(3.5), "PHASE 2\n7 agents\n+24h pipeline", NEON_PINK),
        (Inches(6.2), "PHASE 3\n10 spec-only\n roadmap", NEON_GOLD),
    ]

    for x, text, color in phases:
        box = add_glow_rectangle(slide, x, Inches(1.8), Inches(2.3), Inches(4), color)
        tb = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(1.9), Inches(3))
        tf = tb.text_frame
        tf.text = text
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.name = "Consolas"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(2)
            p.space_after = Pt(2)

# ============================================================================
# SLIDE 7: "PREMISE.CHANGED"
# ============================================================================

def slide_07_premise_changed(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide, 0.3, NEON_PINK)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.02 > INFLECTION")
    add_slide_id(slide, 7)

    # Title
    add_large_text(slide, "OPUS 4.7", Inches(0.5), Inches(1.5), Inches(9), Inches(1), size=150, color=NEON_CYAN, font="Consolas")

    # Subtitle
    subtitle_text = "Released 2026-04-22"
    tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.3))
    tf_sub = tb_sub.text_frame
    tf_sub.text = subtitle_text
    p_sub = tf_sub.paragraphs[0]
    p_sub.font.name = "Consolas"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MID
    p_sub.alignment = PP_ALIGN.CENTER

    # 4 specs
    specs = [
        (Inches(0.8), "1M CTX", NEON_CYAN),
        (Inches(2.8), "8K THINK", NEON_PINK),
        (Inches(4.8), "90% CACHE", NEON_GOLD),
        (Inches(6.8), "HAIKU.4.5", NEON_GREEN),
    ]

    for x, label, color in specs:
        add_metric_box(slide, label, "✓", x, Inches(3.5), size_w=1.4, size_h=1.0, accent_color=color)

    # Quote
    quote = "`> The cost math inverted overnight.`"
    tb_quote = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.8))
    tf_quote = tb_quote.text_frame
    tf_quote.text = quote
    tf_quote.word_wrap = True
    p_quote = tf_quote.paragraphs[0]
    p_quote.font.name = "Consolas"
    p_quote.font.size = Pt(14)
    p_quote.font.color.rgb = NEON_GREEN
    p_quote.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 8: BEFORE / AFTER DIFF
# ============================================================================

def slide_08_before_after(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.02 > DIFF")
    add_slide_id(slide, 8)

    # Left: BEFORE
    left_box = add_glow_rectangle(slide, Inches(0.5), Inches(1.0), Inches(4), Inches(5.5), NEON_PINK)
    tb_before = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(3.6), Inches(0.3))
    tf_before = tb_before.text_frame
    tf_before.text = "BEFORE"
    p_before = tf_before.paragraphs[0]
    p_before.font.name = "Consolas"
    p_before.font.size = Pt(11)
    p_before.font.bold = True
    p_before.font.color.rgb = NEON_PINK
    p_before.alignment = PP_ALIGN.CENTER

    # Right: AFTER
    right_box = add_glow_rectangle(slide, Inches(5.5), Inches(1.0), Inches(4), Inches(5.5), NEON_CYAN)
    tb_after = slide.shapes.add_textbox(Inches(5.7), Inches(1.2), Inches(3.6), Inches(0.3))
    tf_after = tb_after.text_frame
    tf_after.text = "AFTER"
    p_after = tf_after.paragraphs[0]
    p_after.font.name = "Consolas"
    p_after.font.size = Pt(11)
    p_after.font.bold = True
    p_after.font.color.rgb = NEON_CYAN
    p_after.alignment = PP_ALIGN.CENTER

    # Diff items
    diffs = [
        ("- 7 agents", "+ 7 agents spec"),
        ("- manual CLI", "+ exec_orca auto"),
        ("- no orch", "+ multi-AI router"),
        ("- codex only", "+ codex×4, gemini×2"),
        ("- no cache", "+ 85% savings"),
    ]

    before_y = Inches(1.8)
    after_y = Inches(1.8)

    for before_text, after_text in diffs:
        # Before
        tb = slide.shapes.add_textbox(Inches(0.7), before_y, Inches(3.6), Inches(0.35))
        tf = tb.text_frame
        tf.text = before_text
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xFF, 0x80, 0x80)  # Soft red
        before_y += Inches(0.5)

        # After
        tb2 = slide.shapes.add_textbox(Inches(5.7), after_y, Inches(3.6), Inches(0.35))
        tf2 = tb2.text_frame
        tf2.text = after_text
        p2 = tf2.paragraphs[0]
        p2.font.name = "Consolas"
        p2.font.size = Pt(9)
        p2.font.color.rgb = NEON_GREEN
        after_y += Inches(0.5)

# ============================================================================
# SLIDE 9: UPGRADE METRICS
# ============================================================================

def slide_09_upgrade_metrics(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > METRICS")
    add_slide_id(slide, 9)

    # Title
    add_large_text(slide, "PHASE 1+2 DELTA", Inches(0.5), Inches(0.7), Inches(9), Inches(0.4), size=24, color=NEON_CYAN, font="Malgun Gothic")

    # 10 metric boxes (2x5 grid)
    metrics = [
        ("8 Tables", "DB", NEON_CYAN),
        ("11 Modules", "Python", NEON_CYAN),
        ("2,061 Lines", "Insert", NEON_GOLD),
        ("7 Agents", "Spec", NEON_PINK),
        ("33/33 Tests", "PASS", NEON_GREEN),
        ("90% Cache", "Savings", NEON_GOLD),
        ("10m-2h", "Backoff", NEON_CYAN),
        ("25 Plugins", "Total", NEON_CYAN),
        ("5 Validators", "Active", NEON_CYAN),
        ("0 Breaking", "Changes", NEON_GREEN),
    ]

    for idx, (value, label, color) in enumerate(metrics):
        col = idx % 5
        row = idx // 5
        x = Inches(0.6 + col * 1.7)
        y = Inches(1.5 + row * 2.2)

        add_metric_box(slide, label, value, x, y, size_w=1.5, size_h=1.8, accent_color=color)

# ============================================================================
# SLIDE 10: ORCA.DB SCHEMA
# ============================================================================

def slide_10_schema(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > DATABASE")
    add_slide_id(slide, 10)

    # Title
    add_large_text(slide, "ORCA.DB", Inches(0.5), Inches(0.7), Inches(5), Inches(0.4), size=22, color=NEON_CYAN, font="Consolas")

    # 8 table boxes (rough ER layout)
    tables = [
        ("workers", Inches(0.8), Inches(1.5), ["id", "name", "status", "cpu"]),
        ("tasks", Inches(3.0), Inches(1.5), ["id", "plugin", "status", "result"]),
        ("logs", Inches(5.2), Inches(1.5), ["id", "task_id", "level", "msg"]),
        ("state", Inches(7.4), Inches(1.5), ["key", "value", "ttl"]),
        ("plugins", Inches(0.8), Inches(3.8), ["prefix", "name", "status"]),
        ("heartbeats", Inches(3.0), Inches(3.8), ["worker_id", "ts", "alive"]),
        ("quota", Inches(5.2), Inches(3.8), ["ai_model", "used", "limit"]),
        ("cache", Inches(7.4), Inches(3.8), ["hash", "tokens", "hit"]),
    ]

    for name, x, y, cols in tables:
        # Table box
        box = add_glow_rectangle(slide, x, y, Inches(1.8), Inches(1.8), NEON_CYAN, CYBER_MIST)

        # Table name
        tb_name = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(1.6), Inches(0.25))
        tf_name = tb_name.text_frame
        tf_name.text = name.upper()
        p_name = tf_name.paragraphs[0]
        p_name.font.name = "Consolas"
        p_name.font.size = Pt(8)
        p_name.font.bold = True
        p_name.font.color.rgb = NEON_CYAN
        p_name.alignment = PP_ALIGN.CENTER

        # Columns
        tb_cols = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.4), Inches(1.6), Inches(1.2))
        tf_cols = tb_cols.text_frame
        tf_cols.word_wrap = True
        for col in cols:
            if tf_cols.paragraphs:
                tf_cols.add_paragraph()
            idx = len(tf_cols.paragraphs) - 1
            p = tf_cols.paragraphs[idx]
            p.text = f"• {col}"
            p.font.name = "Consolas"
            p.font.size = Pt(6)
            p.font.color.rgb = TEXT_DIM
            p.space_before = Pt(0)
            p.space_after = Pt(0)

    # Query example
    add_terminal_log(slide, [
        "-- worker heartbeat check",
        "SELECT worker_id, last_heartbeat FROM workers",
        "WHERE status='running' ORDER BY last_heartbeat DESC;",
        "",
        "worker_id  | last_heartbeat | status",
        "──────────────────────────────────────",
        "codex-1    | 2026-04-23 10:18:42 | active",
        "gemini-1   | 2026-04-23 10:17:30 | active",
    ], Inches(6.2), Inches(3.5), Inches(3.3), Inches(2.3))

# ============================================================================
# SLIDE 11: WATCHDOG TELEMETRY
# ============================================================================

def slide_11_watchdog(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > WATCHDOG")
    add_slide_id(slide, 11)

    # Status indicator
    status_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.7), Inches(2.5), Inches(0.35))
    status_box.fill.solid()
    status_box.fill.fore_color.rgb = CYBER_MIST
    status_box.line.color.rgb = NEON_GREEN
    status_box.line.width = Pt(1.5)

    tb_status = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(2.3), Inches(0.25))
    tf_status = tb_status.text_frame
    tf_status.text = "✓ WATCHDOG.STATUS: ACTIVE"
    p_status = tf_status.paragraphs[0]
    p_status.font.name = "Consolas"
    p_status.font.size = Pt(10)
    p_status.font.bold = True
    p_status.font.color.rgb = NEON_GREEN

    # Backoff timeline
    add_large_text(slide, "Exponential Backoff", Inches(0.5), Inches(1.4), Inches(9), Inches(0.3), size=14, color=NEON_CYAN, font="Consolas")

    backoff_levels = [
        ("10m", Inches(0.8), NEON_CYAN),
        ("20m", Inches(2.8), NEON_CYAN),
        ("40m", Inches(4.8), NEON_PINK),
        ("2h", Inches(6.8), NEON_GOLD),
    ]

    for level, x, color in backoff_levels:
        # Bar
        bar_height = Inches(1.5)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(1.5), bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = CYBER_MIST
        bar.line.color.rgb = color
        bar.line.width = Pt(1.5)

        # Label
        tb = slide.shapes.add_textbox(x, Inches(3.9), Inches(1.5), Inches(0.25))
        tf = tb.text_frame
        tf.text = level
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Worker status panel
    add_terminal_log(slide, [
        "[WORKER STATUS] 2026-04-23 10:18:45",
        "codex-1      | RUNNING | 1m ago",
        "codex-2      | RUNNING | 2m ago",
        "gemini-1     | QUOTA_WAIT | 30m remaining",
        "haiku-2      | DEAD | reviving in 10m",
    ], Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))

# ============================================================================
# SLIDE 12: ROUTING DECISION TREE
# ============================================================================

def slide_12_routing(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > ROUTING")
    add_slide_id(slide, 12)

    # Title
    add_large_text(slide, "DISPATCH FLOWCHART", Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), size=18, color=NEON_CYAN, font="Consolas")

    # Decision tree
    # Top: START
    start = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.3), Inches(1.6), Inches(0.35))
    start.fill.solid()
    start.fill.fore_color.rgb = CYBER_MIST
    start.line.color.rgb = NEON_CYAN
    start.line.width = Pt(1)
    tb_start = slide.shapes.add_textbox(Inches(4.3), Inches(1.35), Inches(1.4), Inches(0.25))
    tf_start = tb_start.text_frame
    tf_start.text = "task.received"
    p_start = tf_start.paragraphs[0]
    p_start.font.name = "Consolas"
    p_start.font.size = Pt(8)
    p_start.font.color.rgb = NEON_CYAN
    p_start.alignment = PP_ALIGN.CENTER

    # Decision 1
    decision1_y = Inches(2.2)
    decision1 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(3.8), decision1_y, Inches(2.4), Inches(0.7))
    decision1.fill.solid()
    decision1.fill.fore_color.rgb = CYBER_MIST
    decision1.line.color.rgb = NEON_CYAN
    decision1.line.width = Pt(1)
    tb_d1 = slide.shapes.add_textbox(Inches(3.9), decision1_y + Inches(0.1), Inches(2.2), Inches(0.5))
    tf_d1 = tb_d1.text_frame
    tf_d1.text = "lines > 500\nor complex?"
    tf_d1.word_wrap = True
    for p in tf_d1.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(7)
        p.font.color.rgb = TEXT_DIM
        p.alignment = PP_ALIGN.CENTER

    # Branch YES: Codex
    arrow_yes = slide.shapes.add_connector(1, Inches(3.8), decision1_y + Inches(0.5), Inches(2.5), decision1_y + Inches(1.5))
    arrow_yes.line.color.rgb = NEON_PINK
    arrow_yes.line.width = Pt(1.5)
    tb_yes = slide.shapes.add_textbox(Inches(2.8), decision1_y + Inches(0.7), Inches(0.5), Inches(0.2))
    tf_yes = tb_yes.text_frame
    tf_yes.text = "YES"
    p_yes = tf_yes.paragraphs[0]
    p_yes.font.name = "Consolas"
    p_yes.font.size = Pt(7)
    p_yes.font.color.rgb = NEON_PINK

    codex_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), decision1_y + Inches(1.5), Inches(2.6), Inches(0.5))
    codex_box.fill.solid()
    codex_box.fill.fore_color.rgb = CYBER_MIST
    codex_box.line.color.rgb = NEON_PINK
    codex_box.line.width = Pt(1.5)
    tb_codex = slide.shapes.add_textbox(Inches(1.4), decision1_y + Inches(1.6), Inches(2.2), Inches(0.3))
    tf_codex = tb_codex.text_frame
    tf_codex.text = "CODEX (×4 parallel)"
    p_codex = tf_codex.paragraphs[0]
    p_codex.font.name = "Consolas"
    p_codex.font.size = Pt(8)
    p_codex.font.bold = True
    p_codex.font.color.rgb = NEON_PINK
    p_codex.alignment = PP_ALIGN.CENTER

    # Branch NO: Claude
    arrow_no = slide.shapes.add_connector(1, Inches(6.2), decision1_y + Inches(0.5), Inches(7.5), decision1_y + Inches(1.5))
    arrow_no.line.color.rgb = NEON_CYAN
    arrow_no.line.width = Pt(1.5)
    tb_no = slide.shapes.add_textbox(Inches(6.8), decision1_y + Inches(0.7), Inches(0.5), Inches(0.2))
    tf_no = tb_no.text_frame
    tf_no.text = "NO"
    p_no = tf_no.paragraphs[0]
    p_no.font.name = "Consolas"
    p_no.font.size = Pt(7)
    p_no.font.color.rgb = NEON_CYAN

    claude_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), decision1_y + Inches(1.5), Inches(2.6), Inches(0.5))
    claude_box.fill.solid()
    claude_box.fill.fore_color.rgb = CYBER_MIST
    claude_box.line.color.rgb = NEON_CYAN
    claude_box.line.width = Pt(1.5)
    tb_claude = slide.shapes.add_textbox(Inches(6.4), decision1_y + Inches(1.6), Inches(2.2), Inches(0.3))
    tf_claude = tb_claude.text_frame
    tf_claude.text = "CLAUDE (direct)"
    p_claude = tf_claude.paragraphs[0]
    p_claude.font.name = "Consolas"
    p_claude.font.size = Pt(8)
    p_claude.font.bold = True
    p_claude.font.color.rgb = NEON_CYAN
    p_claude.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 13: AI MATRIX
# ============================================================================

def slide_13_ai_matrix(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > AI.MATRIX")
    add_slide_id(slide, 13)

    # Title
    add_large_text(slide, "MODEL COMPARISON MATRIX", Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), size=18, color=NEON_CYAN, font="Consolas")

    # 5 model cards
    models = [
        ("OPUS 4.7", "$15/$75", "1M", "4h", NEON_CYAN),
        ("SONNET 4.6", "$3/$15", "200K", "none", NEON_PINK),
        ("HAIKU 4.5", "$0.8/$4", "200K", "5m", NEON_GREEN),
        ("CODEX", "var", "8K", "none", NEON_GOLD),
        ("GEMINI", "var", "1M", "10m", NEON_PURPLE),
    ]

    for idx, (name, cost, context, cache, color) in enumerate(models):
        x = Inches(0.6 + idx * 1.8)
        y = Inches(1.4)

        card = add_glow_rectangle(slide, x, y, Inches(1.6), Inches(3), color)

        # Model name
        tb_name = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.2), Inches(1.4), Inches(0.3))
        tf_name = tb_name.text_frame
        tf_name.text = name
        tf_name.word_wrap = True
        for p in tf_name.paragraphs:
            p.font.name = "Consolas"
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

        # Specs
        specs_text = f"Cost: {cost}\nContext: {context}\nCache: {cache}"
        tb_specs = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.65), Inches(1.4), Inches(2.2))
        tf_specs = tb_specs.text_frame
        tf_specs.text = specs_text
        tf_specs.word_wrap = True
        for p in tf_specs.paragraphs:
            p.font.name = "Consolas"
            p.font.size = Pt(7)
            p.font.color.rgb = TEXT_DIM
            p.space_before = Pt(2)
            p.space_after = Pt(2)

# ============================================================================
# SLIDE 14: PROMPT CACHE SAVINGS
# ============================================================================

def slide_14_cache_savings(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > CACHE")
    add_slide_id(slide, 14)

    # Large percentage
    add_large_text(slide, "85%", Inches(0.5), Inches(1.2), Inches(9), Inches(1.2), size=180, color=NEON_GOLD, font="Consolas")

    # Subtitle
    subtitle = "24H × 100 TASKS SIMULATION"
    tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.3))
    tf_sub = tb_sub.text_frame
    tf_sub.text = subtitle
    p_sub = tf_sub.paragraphs[0]
    p_sub.font.name = "Consolas"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MID
    p_sub.alignment = PP_ALIGN.CENTER

    # Cost comparison bars
    scenarios = [
        ("NO CACHE", "$1,250", Inches(1.0)),
        ("5MIN TTL", "$325", Inches(2.5)),
        ("1H TTL", "$185", Inches(4.0)),
    ]

    max_cost = 1250
    for scenario, cost_str, x in scenarios:
        cost_val = int(cost_str.replace("$", "").replace(",", ""))
        bar_width = Inches(2.5 * cost_val / max_cost)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.5), bar_width, Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NEON_CYAN
        bar.line.color.rgb = NEON_CYAN
        bar.line.width = Pt(1)

        tb_scenario = slide.shapes.add_textbox(x, Inches(3.2), Inches(1.2), Inches(0.25))
        tf_scenario = tb_scenario.text_frame
        tf_scenario.text = scenario
        p_scenario = tf_scenario.paragraphs[0]
        p_scenario.font.name = "Consolas"
        p_scenario.font.size = Pt(9)
        p_scenario.font.bold = True
        p_scenario.font.color.rgb = TEXT_DIM

        tb_cost = slide.shapes.add_textbox(x + bar_width + Inches(0.1), Inches(3.5), Inches(0.8), Inches(0.6))
        tf_cost = tb_cost.text_frame
        tf_cost.text = cost_str
        p_cost = tf_cost.paragraphs[0]
        p_cost.font.name = "Consolas"
        p_cost.font.size = Pt(10)
        p_cost.font.bold = True
        p_cost.font.color.rgb = NEON_CYAN

    # Formula
    formula = "cost = base × (1 + 0.25 × write) × (1 - 0.9 × hit_rate)"
    tb_formula = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1.5))
    tf_formula = tb_formula.text_frame
    tf_formula.text = formula
    tf_formula.word_wrap = True
    p_formula = tf_formula.paragraphs[0]
    p_formula.font.name = "Consolas"
    p_formula.font.size = Pt(11)
    p_formula.font.color.rgb = NEON_GREEN
    p_formula.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 15: 24/7 PIPELINE
# ============================================================================

def slide_15_pipeline(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > PIPELINE")
    add_slide_id(slide, 15)

    # Title
    add_large_text(slide, "24/7 EXECUTION", Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), size=18, color=NEON_CYAN, font="Consolas")

    # 8-step vertical pipeline
    steps = [
        ("1. ENQUEUE", "task.md arrives", NEON_CYAN),
        ("2. PARSE", "parse spec", NEON_CYAN),
        ("3. ROUTE", "choose AI", NEON_PINK),
        ("4. EXECUTE", "run code", NEON_GOLD),
        ("5. VALIDATE", "test result", NEON_CYAN),
        ("6. LOG", "write orca.db", NEON_CYAN),
        ("7. CACHE", "store prompt", NEON_CYAN),
        ("8. NOTIFY", "webhook", NEON_GREEN),
    ]

    y = Inches(1.4)
    for step, desc, color in steps:
        # Step box
        box = add_glow_rectangle(slide, Inches(1.5), y, Inches(7), Inches(0.35), color)

        # Step text
        tb = slide.shapes.add_textbox(Inches(1.7), y + Inches(0.05), Inches(2), Inches(0.25))
        tf = tb.text_frame
        tf.text = step
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color

        # Description
        tb_desc = slide.shapes.add_textbox(Inches(3.8), y + Inches(0.05), Inches(4), Inches(0.25))
        tf_desc = tb_desc.text_frame
        tf_desc.text = f"→ {desc}"
        p_desc = tf_desc.paragraphs[0]
        p_desc.font.name = "Consolas"
        p_desc.font.size = Pt(8)
        p_desc.font.color.rgb = TEXT_DIM

        y += Inches(0.5)

    # Metrics sidebar
    add_metric_box(slide, "Requests/Min", "14.2", Inches(8.5), Inches(1.5), size_w=1.2, size_h=1.0, accent_color=NEON_CYAN)
    add_metric_box(slide, "Avg Latency", "2.3s", Inches(8.5), Inches(2.8), size_w=1.2, size_h=1.0, accent_color=NEON_CYAN)
    add_metric_box(slide, "Success Rate", "99.7%", Inches(8.5), Inches(4.1), size_w=1.2, size_h=1.0, accent_color=NEON_GREEN)

# ============================================================================
# SLIDE 16: CLI COMMAND MAP
# ============================================================================

def slide_16_cli(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > CLI")
    add_slide_id(slide, 16)

    # Title
    add_large_text(slide, "COMMAND REFERENCE", Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), size=18, color=NEON_CYAN, font="Consolas")

    # 4 terminal sections
    sections = [
        ("INIT", [
            "bash .claude/scripts/install.sh",
            "python validate-plugin-schema.py",
        ], Inches(0.5), Inches(1.5), NEON_CYAN),
        ("RUN", [
            "orca-dispatch task.md",
            "codex-auto --parallel 4",
        ], Inches(4.7), Inches(1.5), NEON_PINK),
        ("MONITOR", [
            "orca-status --watch",
            "worker-health-check.sh",
        ], Inches(0.5), Inches(3.8), NEON_GOLD),
        ("EMERGENCY", [
            "touch ~/.claude/orca/stop",
            "kill-workers --force",
        ], Inches(4.7), Inches(3.8), NEON_PINK),
    ]

    for section, commands, x, y, color in sections:
        # Box
        box = add_glow_rectangle(slide, x, y, Inches(3.8), Inches(1.8), color)

        # Title
        tb_section = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(3.6), Inches(0.25))
        tf_section = tb_section.text_frame
        tf_section.text = f">> {section}"
        p_section = tf_section.paragraphs[0]
        p_section.font.name = "Consolas"
        p_section.font.size = Pt(9)
        p_section.font.bold = True
        p_section.font.color.rgb = color

        # Commands
        tb_cmds = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.45), Inches(3.5), Inches(1.2))
        tf_cmds = tb_cmds.text_frame
        tf_cmds.word_wrap = True
        for i, cmd in enumerate(commands):
            if i > 0:
                tf_cmds.add_paragraph()
            p = tf_cmds.paragraphs[i]
            p.text = f"$ {cmd}"
            p.font.name = "Consolas"
            p.font.size = Pt(8)
            p.font.color.rgb = NEON_GREEN
            p.space_before = Pt(1)
            p.space_after = Pt(1)

# ============================================================================
# SLIDE 17: METRICS DB QUERY
# ============================================================================

def slide_17_query(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > QUERY")
    add_slide_id(slide, 17)

    # SQL Query on left
    query = """-- cost analysis by model
SELECT model, COUNT(*) as calls,
       SUM(tokens) as total_tokens,
       SUM(cost) as total_cost
FROM task_log
WHERE date >= date('now', '-1 day')
GROUP BY model
ORDER BY total_cost DESC;"""

    add_terminal_log(slide, query.split('\n'), Inches(0.5), Inches(1.0), Inches(4.8), Inches(4.5))

    # Results on right
    results = [
        "model   | calls | tokens  | cost",
        "────────────────────────────────────",
        "opus    | 142   | 856K    | $18.50",
        "codex   | 234   | 2.3M    | $24.10",
        "gemini  | 89    | 450K    | $2.25",
        "haiku   | 412   | 1.8M    | $0.85",
    ]

    add_terminal_log(slide, results, Inches(5.5), Inches(1.0), Inches(4.0), Inches(4.5))

    # Mini charts bottom
    add_metric_box(slide, "Cost Distribution", "Opus 38%", Inches(0.5), Inches(5.8), size_w=1.5, size_h=0.8, accent_color=NEON_CYAN)
    add_metric_box(slide, "Total Tokens", "5.4M", Inches(2.3), Inches(5.8), size_w=1.5, size_h=0.8, accent_color=NEON_GOLD)
    add_metric_box(slide, "Daily Cost", "$45.70", Inches(4.1), Inches(5.8), size_w=1.5, size_h=0.8, accent_color=NEON_PINK)

# ============================================================================
# SLIDE 18: LEGACY PRESERVED
# ============================================================================

def slide_18_legacy(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > COMPATIBILITY")
    add_slide_id(slide, 18)

    # Title
    add_large_text(slide, "BACKWARD COMPATIBILITY", Inches(0.5), Inches(0.7), Inches(9), Inches(0.35), size=18, color=NEON_CYAN, font="Consolas")

    # Checklist
    items = [
        ("api.yaml", "✓ preserved", NEON_GREEN),
        ("task schema", "✓ extended", NEON_CYAN),
        ("db.sql", "✓ migrated", NEON_CYAN),
        ("cli flags", "✓ backward compat", NEON_GREEN),
        ("worker protocol", "✓ v2→v3 wrapper", NEON_CYAN),
        ("cache format", "◆ mutated", NEON_PINK),
        ("auth system", "✓ upgraded", NEON_GOLD),
        ("log schema", "◆ expanded", NEON_PINK),
    ]

    y = Inches(1.5)
    for item, status, color in items:
        tb = slide.shapes.add_textbox(Inches(1.5), y, Inches(7), Inches(0.3))
        tf = tb.text_frame
        tf.text = f"{item.ljust(20)} {status}"
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(10)
        p.font.color.rgb = color
        y += Inches(0.45)

    # Quote
    quote = "`> Skeleton preserved. Muscles regenerated.`"
    tb_quote = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    tf_quote = tb_quote.text_frame
    tf_quote.text = quote
    tf_quote.word_wrap = True
    p_quote = tf_quote.paragraphs[0]
    p_quote.font.name = "Consolas"
    p_quote.font.size = Pt(13)
    p_quote.font.color.rgb = NEON_GREEN
    p_quote.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 19: PHASE 3 ROADMAP
# ============================================================================

def slide_19_phase3(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_NIGHT)
    add_scan_line_decoration(slide)
    add_breadcrumb_nav(slide, "ORCHESTRATION.KIT > PART.03 > FUTURE")
    add_slide_id(slide, 19)

    # Large "10"
    add_large_text(slide, "10", Inches(0.5), Inches(0.8), Inches(2), Inches(0.8), size=120, color=NEON_PINK, font="Consolas")

    # Subtitle
    subtitle = "SPEC-ONLY PLUGINS (PHASE 3)"
    tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(2.5), Inches(0.3))
    tf_sub = tb_sub.text_frame
    tf_sub.text = subtitle
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.font.name = "Consolas"
    p_sub.font.size = Pt(8)
    p_sub.font.color.rgb = TEXT_DIM

    # Plugin list
    plugins_phase3 = [
        "ai_music", "ai_voice", "ai_social", "ai_crm",
        "ai_mlops", "ai_web3", "bundles_team", "bundles_enterprise",
        "mcp_slack", "mcp_linear"
    ]

    y = Inches(1.2)
    for i, plugin in enumerate(plugins_phase3):
        if i == 5:
            y = Inches(1.2)
            x = Inches(3.5)
        else:
            x = Inches(0.5)

        tb = slide.shapes.add_textbox(x, y, Inches(2.8), Inches(0.25))
        tf = tb.text_frame
        tf.text = f"▪ {plugin}"
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = NEON_PINK

        if i < 5:
            y += Inches(0.35)
        else:
            y += Inches(0.35)

    # Unused Opus features
    add_large_text(slide, "Untapped Capabilities", Inches(6.2), Inches(1.2), Inches(3.3), Inches(0.3), size=11, color=NEON_GOLD, font="Consolas")

    features = [
        "Agent SDK", "Artifacts", "Files API", "Vision", "Channel"
    ]

    y = Inches(1.7)
    for feat in features:
        tb = slide.shapes.add_textbox(Inches(6.2), y, Inches(3.3), Inches(0.25))
        tf = tb.text_frame
        tf.text = f"◇ {feat}"
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DIM
        y += Inches(0.35)

    # Progress bars
    add_large_text(slide, "PHASE PROGRESS", Inches(0.5), Inches(4.2), Inches(9), Inches(0.25), size=11, color=NEON_CYAN, font="Consolas")

    phases_progress = [
        ("Phase 1", "60%", NEON_CYAN),
        ("Phase 2", "20%", NEON_PINK),
        ("Phase 3", "0%", NEON_GOLD),
    ]

    y = Inches(4.6)
    for phase, pct, color in phases_progress:
        # Label
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.0), Inches(0.2))
        tf = tb.text_frame
        tf.text = phase
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_DIM

        # Progress bar
        pct_val = int(pct.replace("%", ""))
        bar_width = Inches(3.5 * pct_val / 100)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), y, bar_width, Inches(0.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.width = Pt(0)

        # Percent
        tb_pct = slide.shapes.add_textbox(Inches(5.5), y, Inches(0.5), Inches(0.2))
        tf_pct = tb_pct.text_frame
        tf_pct.text = pct
        p_pct = tf_pct.paragraphs[0]
        p_pct.font.name = "Consolas"
        p_pct.font.size = Pt(9)
        p_pct.font.bold = True
        p_pct.font.color.rgb = color

        y += Inches(0.5)

# ============================================================================
# SLIDE 20: CLOSING / EOF
# ============================================================================

def slide_20_closing(prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    add_dark_background(slide, CYBER_SPACE)
    add_scan_line_decoration(slide, 0.3, NEON_CYAN)

    # ASCII art title
    ascii_art = "[ EOF.reached ]"
    add_large_text(slide, ascii_art, Inches(0.5), Inches(2.2), Inches(9), Inches(0.6), size=80, color=NEON_CYAN, font="Consolas")

    # Quote
    quote1 = "`v1 was built 14 days before Opus 4.7.`"
    quote2 = "`v1 evolved 24 hours after Opus 4.7.`"

    tb_q1 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.3))
    tf_q1 = tb_q1.text_frame
    tf_q1.text = quote1
    p_q1 = tf_q1.paragraphs[0]
    p_q1.font.name = "Consolas"
    p_q1.font.size = Pt(12)
    p_q1.font.color.rgb = NEON_GREEN
    p_q1.alignment = PP_ALIGN.CENTER

    tb_q2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.3))
    tf_q2 = tb_q2.text_frame
    tf_q2.text = quote2
    p_q2 = tf_q2.paragraphs[0]
    p_q2.font.name = "Consolas"
    p_q2.font.size = Pt(12)
    p_q2.font.color.rgb = NEON_GREEN
    p_q2.alignment = PP_ALIGN.CENTER

    # Metadata table
    add_large_text(slide, "PROJECT METADATA", Inches(0.5), Inches(4.5), Inches(9), Inches(0.25), size=10, color=NEON_CYAN, font="Consolas")

    metadata = [
        ("commit", "e973a0f"),
        ("date", "2026-04-23"),
        ("version", "v1.0 + Phase 1+2"),
        ("repo", "github.com/bernakilljos/orchestration"),
        ("status", "production ready"),
    ]

    y = Inches(4.95)
    for key, value in metadata:
        tb = slide.shapes.add_textbox(Inches(1.5), y, Inches(7), Inches(0.22))
        tf = tb.text_frame
        tf.text = f"{key.ljust(12)} :: {value}"
        p = tf.paragraphs[0]
        p.font.name = "Consolas"
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MID
        y += Inches(0.3)

    # System shutdown log
    add_terminal_log(slide, [
        "[2026-04-23 17:42:15] system.summary.complete",
        "[2026-04-23 17:42:20] cache.flush.ok",
        "[2026-04-23 17:42:25] system.shutdown.initiated",
        "[2026-04-23 17:42:30] goodbye.",
    ], Inches(1.5), Inches(6.5), Inches(7), Inches(1.2))

# ============================================================================
# MAIN
# ============================================================================

def main():
    print("[*] Creating Premium PPT v5 (Cyberpunk Theme)...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define all slides
    slides_funcs = [
        slide_01_cover,
        slide_02_timeline,
        slide_03_plugins_matrix,
        slide_04_exec_orch,
        slide_05_source_of_truth,
        slide_06_roadmap,
        slide_07_premise_changed,
        slide_08_before_after,
        slide_09_upgrade_metrics,
        slide_10_schema,
        slide_11_watchdog,
        slide_12_routing,
        slide_13_ai_matrix,
        slide_14_cache_savings,
        slide_15_pipeline,
        slide_16_cli,
        slide_17_query,
        slide_18_legacy,
        slide_19_phase3,
        slide_20_closing,
    ]

    for i, slide_func in enumerate(slides_funcs, 1):
        print(f"  [+] Slide {i:02d}: {slide_func.__name__}")
        slide_func(prs)

    # Save
    output_path = r"C:\pjt\orchestration_v1\outputs\ppt\orchestration-v1-premium-2026-04-23-v5-cyberpunk.pptx"
    prs.save(output_path)
    file_size = __import__('os').path.getsize(output_path) / 1024
    print(f"\n[OK] PPT saved: {output_path}")
    print(f"[OK] Total slides: {len(prs.slides)}")
    print(f"[OK] File size: {file_size:.1f} KB")

if __name__ == "__main__":
    main()
