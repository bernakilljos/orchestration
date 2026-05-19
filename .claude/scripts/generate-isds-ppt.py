#!/usr/bin/env python3
"""ISDS (정보보호공시) PPT 생성 — generate-final-ppt.py 의 경로만 변경한 sibling."""
import asyncio
import sys
from pathlib import Path

try:
    from playwright.async_api import async_playwright
except ImportError:
    print("ERROR: pip install playwright + playwright install chromium")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERROR: pip install python-pptx")
    sys.exit(1)

ROOT = Path(__file__).resolve().parent.parent.parent
HTML_DIR = ROOT / "outputs/ppt-isds/html-source/slides"
PNG_DIR = ROOT / "outputs/ppt-isds/html-source/png-output"
OUTPUT_PPTX = ROOT / "outputs/ppt-isds/정보보호공시-신규사업.pptx"

SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
DEVICE_SCALE = 2


async def render_html_to_png(html_path: Path, png_path: Path):
    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        ctx = await browser.new_context(
            viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
            device_scale_factor=DEVICE_SCALE,
        )
        page = await ctx.new_page()
        await page.goto(html_path.as_uri(), wait_until="networkidle", timeout=30000)
        await page.wait_for_timeout(800)  # 폰트 로딩
        await page.screenshot(
            path=str(png_path),
            full_page=False,
            clip={"x": 0, "y": 0, "width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
        )
        await browser.close()


async def main():
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    html_files = sorted(HTML_DIR.glob("slide-*.html"))
    if not html_files:
        print(f"ERROR: no slides in {HTML_DIR}")
        return 1

    print(f"=== rendering {len(html_files)} slides ===")
    for i, html in enumerate(html_files, 1):
        png = PNG_DIR / f"slide-{i:02d}.png"
        await render_html_to_png(html, png)
        print(f"  [{i}/{len(html_files)}] {html.name} -> {png.name}")

    print(f"\n=== assembling PPTX ===")
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH * 9525)  # 1 EMU = 1/9525 px @ 96dpi
    prs.slide_height = Emu(SLIDE_HEIGHT * 9525)
    blank = prs.slide_layouts[6]
    for png in sorted(PNG_DIR.glob("slide-*.png")):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f"\n[OK] {OUTPUT_PPTX} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(asyncio.run(main()))
