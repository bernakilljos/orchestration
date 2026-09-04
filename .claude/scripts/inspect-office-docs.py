#!/usr/bin/env python3
"""
inspect-office-docs.py — docx/xlsx/pptx/xls/ppt 일괄 본문 미리보기
사용: python inspect-office-docs.py <dir>
출력: 각 파일의 첫 sheet/슬라이드/paragraph 의 텍스트 일부
"""
import sys, os, subprocess
from pathlib import Path

def _ensure(pkg, mod):
    try:
        __import__(mod)
    except ImportError:
        subprocess.run([sys.executable, '-m', 'pip', 'install', '--quiet', pkg], check=True, timeout=300)
        __import__(mod)

_ensure('openpyxl', 'openpyxl')
_ensure('python-docx', 'docx')
_ensure('python-pptx', 'pptx')
_ensure('xlrd', 'xlrd')

import openpyxl, xlrd
from docx import Document
from pptx import Presentation

def inspect_xlsx(p):
    try:
        wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
        sheets = wb.sheetnames
        first = wb[sheets[0]]
        cells = []
        for row in first.iter_rows(max_row=3, max_col=8, values_only=True):
            cells.append(' | '.join(str(c) if c is not None else '' for c in row))
        wb.close()
        return f"sheets={sheets[:3]}; rows: {' / '.join(cells)}"[:300]
    except Exception as e:
        return f"ERR: {e}"

def inspect_xls(p):
    try:
        wb = xlrd.open_workbook(p)
        sheets = wb.sheet_names()
        s = wb.sheet_by_index(0)
        rows = []
        for r in range(min(3, s.nrows)):
            cells = [str(s.cell_value(r,c)) for c in range(min(8, s.ncols))]
            rows.append(' | '.join(cells))
        return f"sheets={sheets[:3]}; rows: {' / '.join(rows)}"[:300]
    except Exception as e:
        return f"ERR: {e}"

def inspect_docx(p):
    try:
        d = Document(p)
        paras = [p.text.strip() for p in d.paragraphs if p.text.strip()][:5]
        return ' / '.join(paras)[:300] or '(no text)'
    except Exception as e:
        return f"ERR: {e}"

def inspect_pptx(p):
    try:
        prs = Presentation(p)
        out = []
        for i, slide in enumerate(prs.slides):
            if i >= 3: break
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        out.append(f"S{i+1}: {txt[:100]}")
                        break
        return ' / '.join(out)[:400] or '(no text)'
    except Exception as e:
        return f"ERR: {e}"

def inspect_ppt(p):
    try:
        _ensure('pywin32', 'win32com')
        import win32com.client
        pp = win32com.client.Dispatch("PowerPoint.Application")
        pres = pp.Presentations.Open(p, WithWindow=False, ReadOnly=True)
        out = []
        for i in range(1, min(4, pres.Slides.Count + 1)):
            slide = pres.Slides.Item(i)
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    txt = shape.TextFrame.TextRange.Text.strip()
                    if txt:
                        out.append(f"S{i}: {txt[:100]}")
                        break
        pres.Close()
        pp.Quit()
        return ' / '.join(out)[:400] or '(no text)'
    except Exception as e:
        return f"ERR (ppt COM): {e}"

def inspect_hwp(p):
    try:
        _ensure('pywin32', 'win32com')
        import win32com.client
        hwp = win32com.client.Dispatch("HWPFrame.HwpObject")
        try: hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
        except: pass
        hwp.Open(p, "HWP", "forceopen:true")
        text = hwp.GetTextFile("TEXT", "")
        hwp.Quit()
        lines = [l.strip() for l in text.split('\n') if l.strip()][:5]
        return ' / '.join(lines)[:400] or '(no text)'
    except Exception as e:
        return f"ERR (hwp COM): {e}"

INSPECTORS = {
    '.xlsx': inspect_xlsx,
    '.xls':  inspect_xls,
    '.docx': inspect_docx,
    '.pptx': inspect_pptx,
    '.ppt':  inspect_ppt,
    # .hwp: 한글 보안 dialog 가 사용자 액션 요구 -> Zero-touch 위반. SKIP
}

def main():
    if len(sys.argv) < 2:
        print("usage: python inspect-office-docs.py <dir>")
        sys.exit(2)
    root = Path(sys.argv[1])
    for f in sorted(root.rglob('*')):
        if f.is_file() and f.suffix.lower() in INSPECTORS:
            insp = INSPECTORS[f.suffix.lower()]
            print(f"--- {f.relative_to(root)} ---")
            print(insp(str(f)))
            print()

if __name__ == '__main__':
    main()
