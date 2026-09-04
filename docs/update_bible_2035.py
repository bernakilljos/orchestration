"""AI Evolution Bible 2035 pptx delta update.

Adds 2 delta slides at the end noting 2026-06 ~ 2026-08 (7주간) new tech
that emerged after the 2026-07-01 original build. Preserves all original
18 slides. Backup at .bak before overwrite.
"""
from __future__ import annotations
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent
SRC = HERE / "AI_Evolution_Bible_2035.pptx"
BAK = HERE / "AI_Evolution_Bible_2035.pptx.bak"

# 13 organs × 2달 새 사건 (우선순위  급함 /  중요 /  참고)
ORGAN_UPDATES = [
    ("Brain 뇌",         "",   "NVIDIA Rubin CPX · Cerebras Wafer scale 최신"),
    ("Memory 기억",      "", "Mem0g graph-enhanced +29.6점 · Zep GraphRAG sub-second"),
    ("Generation 생성",  "",  "Gemini 3.7 Flash (8/13) · code diffusion · Sora 2 rumor"),
    ("Connectivity 연결","",   "CXL 3.1 · NVLink 6 · Meta MTIA v3"),
    ("Perception 지각",  "",  "Genie 3 (DeepMind world model, 8월) · 멀티모달 확장"),
    ("Reasoning 추론",   "", "OpenAI Erdős 자율 발견 (Fields Medalist 인정) · Gemini 2.5 Deep Think"),
    ("Energy 에너지",    "",   "Anthropic 소형 원전 · Microsoft Three Mile Island · Google 핵융합"),
    ("Trust 신뢰",       "",   "Lakera·HiddenLayer 신제품 · Anthropic MechInterp 논문"),
    ("Agency 행위",      "", "ChatGPT Work (GPT-5.6) · OpenAI Sol/Terra/Luna tier · Managed Agents · Fable 5 RESTORED"),
    ("Embodiment 체화",  "",  "Figure 03 · Optimus Gen 3 · Unitree G1 대량생산"),
    ("Simulation 시뮬",  "",  "NVIDIA Cosmos world foundation · Genie 3 playable worlds"),
    ("Civilization 문명","",   "Managed Agents webhooks (환경·메모리 lifecycle) · agent 경제 논문"),
    ("Evolution 진화",   "", "자율 수학 발견 (Erdős) · Claude self-improvement · AlphaEvolve"),
]

TIMELINE_ADDS = [
    ("2026-07-01",  "AI Evolution Bible 2035 원본 작성"),
    ("2026-07-24",  "Claude Opus 5 launch — 1M context 기본"),
    ("2026-07-01",  "Anthropic Fable 5 RESTORED (US export-control 6/12 → 7/1 복원)"),
    ("2026-07-02",  "Claude Sonnet 5 launch · 이후 8/10 $2/$10 확정"),
    ("2026-07-21",  "Gemini 3.6 Flash · 8/13 Gemini 3.7 Flash"),
    ("2026-08 초",  "OpenAI Sol/Terra/Luna tier 재편 · ChatGPT Work"),
    ("2026-08",     "OpenAI reasoning 자율 수학 발견 (Erdős) — Fields Medalist 인정"),
    ("2026-04~06",  "Mem0g graph-enhanced (temporal +29.6점)"),
    ("2026-06~07",  "Zep GraphRAG · Managed Agents session streams+webhooks"),
]


def add_delta_slide_1(prs: Presentation) -> None:
    """Slide 19: 13 organs × 지난 2달 새 사건."""
    layout = prs.slide_layouts[0]  # DEFAULT (only one)
    slide = prs.slides.add_slide(layout)

    # Add title as textbox (no title placeholder in this template)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.8))
    title_box.text_frame.text = "델타 갱신 · 2026-07 이후 7주 새 기술"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True

    # Table
    rows = len(ORGAN_UPDATES) + 1
    cols = 3
    left, top, width, height = Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.8)
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table

    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(1.0)
    tbl.columns[2].width = Inches(9.3)

    headers = ["Organ (13/13)", "급함", "지난 2달 새 사건"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True

    for i, (organ, prio, event) in enumerate(ORGAN_UPDATES, 1):
        for j, val in enumerate([organ, prio, event]):
            cell = tbl.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)


def add_delta_slide_2(prs: Presentation) -> None:
    """Slide 20: 2026-07~08 진화 타임라인 추가분."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.9))
    title_box.text_frame.text = "2026-07~08 타임라인 (원본 슬라이드 16 후속)"
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True

    left, top, width, height = Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6)
    tbl = slide.shapes.add_table(len(TIMELINE_ADDS) + 1, 2, left, top, width, height).table
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(10.5)

    tbl.cell(0, 0).text = "시점"
    tbl.cell(0, 1).text = "사건"
    for j in range(2):
        tbl.cell(0, j).text_frame.paragraphs[0].font.bold = True
        tbl.cell(0, j).text_frame.paragraphs[0].font.size = Pt(12)

    for i, (when, what) in enumerate(TIMELINE_ADDS, 1):
        tbl.cell(i, 0).text = when
        tbl.cell(i, 1).text = what
        for j in range(2):
            for p in tbl.cell(i, j).text_frame.paragraphs:
                p.font.size = Pt(11)


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"원본 없음: {SRC}")

    # 백업
    shutil.copy2(SRC, BAK)
    print(f"[backup] {BAK.name}")

    prs = Presentation(str(SRC))
    orig_count = len(prs.slides)
    print(f"[open]   원본 {orig_count} 슬라이드")

    add_delta_slide_1(prs)
    add_delta_slide_2(prs)

    prs.save(str(SRC))
    new_count = len(Presentation(str(SRC)).slides)
    print(f"[save]   {new_count} 슬라이드 (+{new_count - orig_count} 델타)")
    print(f"[done]   {SRC.name}")


if __name__ == "__main__":
    main()
