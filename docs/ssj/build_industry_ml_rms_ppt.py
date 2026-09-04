"""
산업 ML → RMS 이식 종합 지도 PPT (누락 없이)
- 16 산업 × ML → RMS O/X 매핑
- 14 AI 기술 카테고리
- 종합 Killer 19개
- Phase 로드맵
- KPMG 대응 전략
- AI ISMS-P Copilot
- 40+ 슬라이드
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── 색 팔레트 ─────────────────────────
NAVY = RGBColor(0x1A, 0x2B, 0x4E)
BLUE = RGBColor(0x1E, 0x6F, 0xB8)
CYAN = RGBColor(0x00, 0xB4, 0xD8)
GREEN = RGBColor(0x1F, 0xB2, 0x5A)
RED = RGBColor(0xE7, 0x3C, 0x3C)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
GRAY_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_DK = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x14, 0x1A, 0x2E)

# ─── PPT 설정 (16:9) ────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def add_bg(slide, color=WHITE):
    left = top = 0
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, SW, SH)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bg=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    if bg is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg; tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "맑은 고딕"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, color=BLUE, line=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = color
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line; box.line.width = Pt(0.5)
    return box

def add_header(slide, title, subtitle=None):
    """상단 헤더 (파랑 바)."""
    add_rect(slide, 0, 0, SW, Inches(0.75), NAVY)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(12), Inches(0.6),
             title, size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.85), Inches(12), Inches(0.4),
                 subtitle, size=13, color=GRAY_DK, anchor=MSO_ANCHOR.TOP)

def add_footer(slide, text):
    """하단 바."""
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), NAVY)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.3),
             text, size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, WHITE)
    return slide

def add_table(slide, x, y, w, h, headers, rows, header_bg=BLUE, header_fg=WHITE,
              row_bg=GRAY_BG, alt_bg=WHITE, first_col_bg=None, col_widths=None,
              font_size=11, header_size=12):
    """수동 표. col_widths = 비율 list. rows = [(cell1, cell2, ...), ...] cell 은 str or (str, color)."""
    n_col = len(headers)
    n_row = len(rows)
    if col_widths is None:
        col_widths = [1] * n_col
    total_w = sum(col_widths)
    col_w = [w * cw / total_w for cw in col_widths]

    row_h = h / (n_row + 1)
    header_h = row_h
    row_h_body = (h - header_h) / max(n_row, 1)

    # header
    cx = x
    for i, hd in enumerate(headers):
        add_rect(slide, cx, y, col_w[i], header_h, header_bg)
        add_text(slide, cx, y, col_w[i], header_h, hd,
                 size=header_size, bold=True, color=header_fg,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]

    # rows
    ry = y + header_h
    for ri, row in enumerate(rows):
        bg = alt_bg if ri % 2 == 0 else row_bg
        cx = x
        for ci, cell in enumerate(row):
            if isinstance(cell, tuple):
                text_val, color = cell
            else:
                text_val, color = cell, BLACK
            cell_bg = first_col_bg if ci == 0 and first_col_bg else bg
            add_rect(slide, cx, ry, col_w[ci], row_h_body, cell_bg)
            align = PP_ALIGN.CENTER if ci == n_col - 1 else PP_ALIGN.LEFT
            add_text(slide, cx, ry, col_w[ci], row_h_body, text_val,
                     size=font_size, color=color,
                     align=align, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[ci]
        ry += row_h_body

def ox_cell(val):
    """O/X 셀 색상."""
    if val == "O":
        return ("O", GREEN)
    if val.startswith("O"):
        return (val, GREEN)
    if val == "X":
        return ("X", RED)
    if val.startswith("X"):
        return (val, RED)
    if "검토" in val:
        return (val, ORANGE)
    return (val, BLACK)

# ─── SLIDE 01 · 타이틀 ────────────────────
s = new_slide()
add_bg(s, NAVY)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.8), Inches(0.9),
         "산업 ML → RMS 이식 종합 지도",
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.4), Inches(11.8), Inches(0.6),
         "Industry ML Transfer to Risk Monitoring Solution",
         size=22, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.5), Inches(11.8), Inches(0.5),
         "16 산업 · 14 AI 기술 카테고리 · 40+ 접목 아이디어",
         size=18, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.8), Inches(0.4),
         "아이티센 코어 ESG 사업부  ·  2026-07-01",
         size=13, color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 02 · 목차 ────────────────────
s = new_slide()
add_header(s, "목차", "이 문서에 담긴 것 (누락 없이)")
toc = [
    ("PART 1", "전체 프레임 · Bible × 산업 × RMS"),
    ("PART 2", "16 산업별 ML → RMS 매핑 (자율주행·해양·항공·위성·게임·의료·반도체·농업·물류·통신·광고·스포츠·바이오·방위·에너지·Physical AI)"),
    ("PART 3", "14 AI 기술 카테고리 (두뇌·기억·에이전트·이상탐지·시계열·그래프·비전·프라이버시·최적화·시뮬·감사·학습·생성·크로스도메인)"),
    ("PART 4", "종합 Killer 19개 · X (안 어울림) 리스트"),
    ("PART 5", "Phase 로드맵 (0~24개월)"),
    ("PART 6", "KPMG 대응 전략 · 시장 분리"),
    ("PART 7", "AI ISMS-P Copilot (사용자 개인 강점)"),
    ("PART 8", "Bible 13 organs 매핑 · 다음 액션"),
]
y = Inches(1.5)
for tag, desc in toc:
    add_rect(s, Inches(0.7), y, Inches(1.4), Inches(0.55), BLUE)
    add_text(s, Inches(0.7), y, Inches(1.4), Inches(0.55), tag,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.3), y, Inches(10.5), Inches(0.55), desc,
             size=13, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)
add_footer(s, "PART 8 로 마무리 · 총 40+ 슬라이드")

# ─── SLIDE 03 · 전체 프레임 ────────────────────
s = new_slide()
add_header(s, "PART 1 · 전체 프레임", "Bible 13 organs × 산업 ML × RMS 접목")
add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
         "🧠 Bible 13 organs (AI 진화 지도)",
         size=15, bold=True, color=NAVY)
bible = [
    ("뇌 Brain", "추론·판단"),
    ("기억 Memory", "장기 이력"),
    ("생성 Generation", "답 만듦"),
    ("연결 Connectivity", "관계망"),
    ("지각 Perception", "보고 듣기"),
    ("추론 Reasoning", "생각 시간"),
    ("에너지 Energy", "GPU 비용"),
    ("신뢰 Trust", "증명"),
    ("행위 Agency", "자율 실행"),
    ("체화 Embodiment", "물리 세계"),
    ("시뮬 Simulation", "미래 예측"),
    ("AI 문명", "AI 협업"),
    ("진화 Evolution", "자기학습"),
]
y = Inches(1.85)
for i, (name, desc) in enumerate(bible):
    col = i % 3
    row = i // 3
    xx = Inches(0.6 + col * 2.15)
    yy = Inches(1.85 + row * 0.55)
    add_rect(s, xx, yy, Inches(2.05), Inches(0.5), GRAY_BG, line=BLUE)
    add_text(s, xx, yy, Inches(0.9), Inches(0.5), name,
             size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, xx + Inches(0.9), yy, Inches(1.15), Inches(0.5), desc,
             size=10, color=GRAY_DK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(7.3), Inches(1.4), Inches(5.8), Inches(0.4),
         "🌍 산업 ML → RMS 접목 원리",
         size=15, bold=True, color=NAVY)
principles = [
    "① 다른 산업의 ML 은 이미 성숙 · 우리는 이식만",
    "② 산업 기술의 mechanism 을 RMS 데이터에 매핑",
    "③ 물리 센서 X · 논리·행동·거래 데이터 O",
    "④ Bible 13 organs 중 미사용 organ 채우기",
    "⑤ 5살 톤 · 비유 · O/X 명확 · 접목 essay 금지",
    "⑥ 한 산업만 X · 여러 산업 breadth · spark 오면 depth",
]
y = Inches(1.85)
for p in principles:
    add_rect(s, Inches(7.3), y, Inches(5.8), Inches(0.5), GRAY_BG)
    add_text(s, Inches(7.3), y, Inches(5.8), Inches(0.5), p,
             size=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)

# 하단 결론
add_rect(s, Inches(0.6), Inches(6.1), Inches(12.5), Inches(0.9), NAVY)
add_text(s, Inches(0.6), Inches(6.1), Inches(12.5), Inches(0.9),
         "결론 · Bible × 산업 = 씨앗 무한 생성. RMS 는 지각·행위·신뢰 3 organ 만 사용 → 10 organ 미개척",
         size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "PART 1 · 3/40")

# ─── PART 2 · 16 산업 매핑 ────────────────────

INDUSTRIES = [
    ("자율주행", "Autonomous Driving", "🚗", [
        ("Edge Case Detection", "훈련 안 한 이상한 상황 감지 (야생동물 튀어나옴)", "O", "처음 보는 결재 패턴 자동 flag"),
        ("GAN Synthetic Data", "가짜 훈련 데이터 만드는 화가", "O ", "부정 사례 소수 → 대량 가짜 부정 생성"),
        ("Sensor Fusion", "눈+귀+코 다 써서 판단 (LiDAR+Camera+RADAR)", "O", "ERP+이메일+CCTV 합쳐서 판단"),
        ("Facial Expression AI", "운전자 졸음·피로 얼굴로 감지", "O", "임원 회의 스트레스·압박 감지"),
        ("V2X (차량 통신)", "차끼리 위험 정보 실시간 공유", "O", "회사 간 부정 정보 공유 (컨소시엄)"),
        ("Path Planning", "목적지까지 경로 계획", "검토", "감사 조사 경로 계획 각도로 가능"),
        ("LiDAR (3D 스캔)", "레이저로 3D 지도", "X", "RMS 는 물리 스캔 안 함"),
        ("HD Map", "cm 단위 도로 지도", "X", "지도 필요 없음"),
        ("Occupancy Network", "3D 공간 점유 예측", "X", "물리 공간 개념 없음"),
        ("End-to-End Driving", "카메라 → 바로 핸들 조작", "X", "RMS 는 별도 아키텍처"),
    ]),
    ("해양", "Maritime", "🚢", [
        ("AIS-LLM 궤적 분석", "배 위치·정상 항로 학습 · dark zone·loitering·jump point 감지", "O ", "결재 궤적 dark zone (승인 우회)·loitering (반복 결재)·jump point (금액 급증)"),
        ("Kalman Filter", "센서 노이즈 낀 배 위치·속도 실시간 최적 추정", "O ", "실시간 위험도 curve 유지 (매 순간 업데이트)"),
        ("SAR 위성 이미지 융합", "레이더 위성으로 배 감지", "O", "다차원 데이터 융합 (외부 시그널)"),
        ("Windward Behavioral Risk", "선박 행동 리스크 2026 상용", "O", "UEBA 다음 단계 (행동 프로필)"),
        ("배 물리 시스템 (LNG 탱크)", "실제 선박 기계", "X", "물리 세계 X"),
    ]),
    ("항공", "Aviation", "✈", [
        ("Predictive Maintenance (RUL)", "부품 15~30일 전 고장 예측 (진동·온도)", "O ", "부정 시도 30일 전 예측 · KPMG 도 없음 · killer"),
        ("Flight Data Recorder (블랙박스)", "모든 센서·조종사 대화 통합 기록", "O", "결재+이메일+통화 통합 블랙박스 (v14 실장)"),
        ("SMS 근본원인 분석", "사고 5-Whys 자동 · 사람·프로세스·기술·환경 분류", "O", "부정 발생 시 근본원인 자동 (재발 방지)"),
        ("AOG 경제 손실 예방", "$10K~$150K/h 손실 예방", "O", "부정 손실 예방 · 경제 근거"),
        ("실제 항공기 물리 부품", "엔진·날개", "X", "물리 세계 X"),
    ]),
    ("위성", "Satellite Telemetry", "🛰", [
        ("Explainable Telemetry Anomaly", "94% precision · 판단 근거 자동 설명 (ESA 벤치마크)", "O ", "RFP § 3.2.2 하이라이팅 요구 정확 대응"),
        ("GCN/TCN 시계열", "그래프 CNN + Temporal CNN", "O", "결재 시계열 실시간 이상"),
        ("Auto-Encoder Reconstruction", "정상 압축·복원 오차", "O", "결재 이상 감지"),
        ("Adversarial Auto-Encoder", "생성 대립 이상탐지", "O", "부정 데이터 강화"),
        ("위성 관제 물리 시스템", "궤도 제어", "X", "물리 세계 X"),
    ]),
    ("게임", "Gaming Anti-Cheat", "🎮", [
        ("LSTM+Transformer sequence", "6개월 게임 sequence 통째 학습 (매치 전체 span)", "O ", "6개월 결재 sequence Transformer 학습"),
        ("Behavioral Anomaly", "움직임·조준·반응시간 통계 이상", "O", "결재 시간·금액·거래처 통계 이상"),
        ("AI 치트 vs AI 안티-치트 진화", "서로 학습 대결", "O", "AI 부정 vs AI RMS 진화 프레임"),
        ("HITL 모드", "AI 감지 후 관리자 최종 결정", "O", "감사인 최종 결정 · RFP § HITL 요구"),
        ("게임 그래픽 렌더링", "3D 렌더", "X", "그래픽 관련 X"),
    ]),
    ("의료", "Medical Radiology", "🏥", [
        ("Few-shot pathology detection", "적은 소수 데이터로 rare 질환 학습", "O ", "소수 부정 사례로도 학습 · 데이터 부족 해결"),
        ("Contrastive Learning", "정상 vs 이상 임베딩 대조 (자폐 생쥐 논문)", "O ", "정상 결재 임베딩 vs 부정 결재 대조"),
        ("Interpretable Generative Anomaly", "판단 근거 자동 설명 이미지 생성", "O", "판단 근거 자동 시각화 · Explainable"),
        ("FDA/EMA Explainable Governance", "규제 대응 표준", "O", "감사인·규제기관 receptive"),
        ("실제 CT/MRI 촬영", "물리 영상 획득", "X", "촬영 관련 X"),
    ]),
    ("반도체", "Semiconductor", "💾", [
        ("Synthetic Defect Generation", "실제 결함 부족 → GAN 으로 합성 결함", "O ", "부정 사례 대량 합성 학습"),
        ("Two-layer AI (hardware+data)", "제조 층 + 데이터 층 분리 감시", "O", "ERP 층 + RMS 층 분리"),
        ("False Positive -60% 필터링", "알림 폭발 해결", "O", "RMS 알림 폭발 해결"),
        ("95% Accuracy classification", "결함 유형 자동 분류", "O", "부정 유형 자동 분류"),
        ("실제 웨이퍼 검사", "물리 스캔", "X", "물리 세계 X"),
    ]),
    ("농업", "Precision Agriculture", "🌾", [
        ("UAV+IoT+위성 multi-source fusion", "드론·센서·위성 다층 데이터 융합", "O ", "ERP+이메일+CCTV+전화 다센서 융합"),
        ("정밀 조기 감지 (질병 3주 전)", "농작물 병해 조기", "O", "부정 조기 신호 감지 (몇 주 전)"),
        ("Spectral Imaging", "농작물 스펙트럼 분석", "O", "다차원 데이터 (참고)"),
        ("실제 농약 살포", "물리 액션", "X", "농약 살포 X"),
        ("UAV 물리 조종", "드론 하드웨어", "X", "드론 조종 X"),
    ]),
    ("물류·공급망", "Logistics", "📦", [
        ("Self-healing Supply Chain", "이상 감지 후 자동 시정", "O ", "이상 감지 후 자동 시정 RMS"),
        ("Predictive → Prescriptive", "detection → 자동 처방", "O", "탐지에서 자동 처방으로 패러다임 이동"),
        ("Time-series 재해 -20~30% 단축", "예측 기반 복구 시간 감소", "O", "부정 사건 복구 시간 급감"),
        ("멀티모달 리스크 신호", "여러 신호 통합", "O", "여러 데이터 소스 통합"),
        ("실제 트럭 GPS", "물리 위치", "X", "물리 세계 X"),
    ]),
    ("통신 5G", "Telecom 5G/6G", "📡", [
        ("Federated Reinforcement Learning", "여러 회사 데이터 안 옮기고 공동 학습", "O ", "여러 증권사 공동 학습 (경쟁 유지)"),
        ("Neuromorphic SNN 초저지연", "밀리초 대응 (뉴로모픽 칩)", "O", "실시간 대응 (밀리초)"),
        ("Transformer traffic anomaly", "네트워크 트래픽 Transformer", "O", "결재·이체 트래픽 감시"),
        ("Agentic Network 자율 관제", "자율 네트워크", "O", "자율 RMS 관제"),
        ("실제 5G 안테나 하드웨어", "물리 안테나", "X", "물리 세계 X"),
    ]),
    ("광고", "AdTech", "📢", [
        ("18-signal real-time feature engine", "결재당 18+ 지표 병렬 스코어링", "O ", "결재당 18+ 지표 병렬 스코어링"),
        ("Adversarial AI Fraud 방어", "AI 로 진화하는 사기 대응", "O", "AI 로 진화하는 부정 대응"),
        ("Stacking Ensemble + 9 ML", "여러 모델 조합 정확도", "O", "여러 모델 조합 정확도"),
        ("Agentic AI Chatbot 실시간 대응", "자율 대응", "O", "부정 실시간 대응"),
        ("실제 광고 CTR 예측", "광고 도메인 특화", "X", "광고 관련 X"),
    ]),
    ("스포츠 VAR", "Sports Officiating", "⚽", [
        ("3D 카메라 트래킹 (Sony Hawk-Eye)", "3D 공간 물체 추적", "O", "사무실 CCTV 3D 재구성"),
        ("Vision Foundation Model (SoccerMaster)", "행동 파운데이션 모델", "O", "행동 파운데이션 모델"),
        ("반자동 판정 (AI + 사람)", "AI 계산 후 사람 검토", "O", "HITL 표준 · RFP 요구"),
        ("실시간 3D 재구성", "다각도 카메라 융합", "O", "다센서 융합 참고"),
        ("실제 볼 궤적 물리", "물리 궤적", "X", "물리 계산 X"),
    ]),
    ("바이오·제약", "Pharmacovigilance", "💊", [
        ("Pharmacovigilance Signal Detection", "소셜·문헌에서 부작용 사전 추출", "O ", "소셜·내부 대화에서 부정 신호 사전 추출"),
        ("15일 → 5일 자동화", "케이스 처리 시간 급감", "O", "신고서·감사 대응 시간 급감"),
        ("FDA/EMA Explainable+Inspection-ready", "규제 대응 표준", "O", "감사인 조사 대응 표준"),
        ("WHO VigiBase 3500만 케이스", "대량 데이터 활용", "O", "K-DART 등 공개 데이터 활용"),
        ("실제 약물 시험", "물리 실험", "X", "약물 관련 X"),
    ]),
    ("방위 드론", "Defense Drone Swarm", "🛸", [
        ("Multi-sensor swarm counter", "RF+radar+EO+acoustic 다센서 융합", "O ", "다센서 부정 스웜 감지"),
        ("Behavior baseline + trajectory anomaly", "정상 프로필 대비 실시간 이탈", "O", "정상 프로필 대비 실시간 이탈"),
        ("AI swarm coordination", "여러 드론 자율 협업", "O", "여러 RMS 에이전트 자율 협업"),
        ("Autonomous intercept", "자율 차단", "O", "자율 대응 (Computer Use 관점)"),
        ("실제 무기 통제", "물리 무기", "X", "무기 X"),
    ]),
    ("에너지 스마트그리드", "Smart Grid", "⚡", [
        ("GAN-LSTM Electricity Theft", "전기 도용 감지 (GAN+LSTM)", "O ", "자금 도용 (같은 원리)"),
        ("Spatio-Temporal Graph", "공간·시간 그래프", "O", "자금 이동 그래프 감시"),
        ("Smart Meter Anomaly", "스마트미터 이상탐지", "O", "결재·계좌 이상탐지 참고"),
        ("Neural Structured Prediction", "구조화 예측", "O", "구조 데이터 예측"),
        ("실제 발전소 물리", "물리 설비", "X", "설비 X"),
    ]),
    ("Physical AI", "Physical AI / 산업안전", "🤖", [
        ("PPE Detection + Near-miss capture", "안전장구 감지·아슬아슬한 순간 캡처", "O ", "현장 부정 준비 행동 감지"),
        ("Motion Capture AI", "자세·행동 자동 분석", "O", "사무실 이상 행동 (문서 훔침)"),
        ("Voxel/Cority/Protex 상용", "기존 CCTV 재활용", "O", "v14 Robot Blackbox 실장"),
        ("Ergonomic Risk Assessment", "자세 위험 자동", "O", "직원 스트레스 자세 지표"),
        ("실제 로봇 하드웨어", "물리 로봇", "X", "물리 로봇 X"),
    ]),
]

for idx, (ko, en, emoji, techs) in enumerate(INDUSTRIES, 1):
    s = new_slide()
    add_header(s, f"{emoji}  {ko} → RMS 이식", f"{en}  ·  기술 {len(techs)}개  ·  O/X 판정")
    rows = []
    for tech, desc, ox, rms in techs:
        # 이름 + 5살 설명 결합
        combined = f"{tech}\n  ↳ {desc}"
        rows.append((combined, ox_cell(ox), rms))
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
              ["기술 (5살 설명)", "RMS", "접목 시 효과"],
              rows,
              header_bg=NAVY, header_fg=WHITE,
              first_col_bg=None, col_widths=[4.5, 1.2, 6.6],
              font_size=10, header_size=12)
    # 하단 요약
    o_count = sum(1 for t in techs if t[2].startswith("O"))
    x_count = sum(1 for t in techs if t[2].startswith("X"))
    review_count = sum(1 for t in techs if "검토" in t[2])
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), GRAY_BG, line=BLUE)
    summary = f"요약  ·  O (접목 가능): {o_count}  ·  X (안 어울림): {x_count}  ·  검토: {review_count}"
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), summary,
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, f"PART 2 · 산업 {idx}/16")

# ─── PART 3 · 14 AI 기술 카테고리 ────────────────────

CATEGORIES = [
    ("🧠 AI 두뇌 (Reasoning)", "AI 가 판단·추론하는 기술", [
        ("Claude Opus 5 🆕", "1M context 기본+최대 · thinking on-by-default · effort ladder (2026-07-24)", "O ", "복잡 판정 default · Fable 5 대비 절반 비용"),
        ("자율 수학 발견 🆕", "OpenAI reasoning 이 Erdős 추측 반증 (1946년 open · Fields Medalist 인정)", "O ", "agent 스스로 새 규정·룰 발견 가능성 시사"),
        ("Gemini 2.5 Pro Deep Think 🆕", "GPQA Diamond 82.4% · MMLU-Pro 89.8%", "O ", "검증 fallback 후보"),
        ("LLM (GPT·Claude·Gemini)", "큰 언어 모델", "O", "이미 씀"),
        ("SLM (Phi·Gemma·Qwen)", "작은 온프레미스 LLM", "O", "RFP 대응"),
        ("Test-Time Compute (o1·o3·R1)", "어려운 문제 오래 생각", "O", "고위험 결재만 deep · 비용 -90%"),
        ("Chain-of-Thought", "단계별 추론", "O", "판단 근거 명시"),
        ("Tree-of-Thought", "여러 경로 병렬 추론", "O", "여러 시나리오 병렬"),
        ("Multi-Agent Debate", "여러 AI 합의", "O", "환각 방지 · False positive 급감"),
        ("Constitutional AI", "AI 자기 규정 검열", "O", "K-SOX 규정 내장"),
        ("MoE (Mixture of Experts)", "전문가 활성화", "O", "도메인별 sLLM 전환"),
        ("Mamba (State Space)", "긴 시퀀스 · Transformer 대안", "검토", "성능 대비 필요"),
        ("World Model", "세상 모델링 (Physical Intelligence)", "X", "실험 단계"),
    ]),
    ("💾 AI 기억·검색 (Memory·Retrieval)", "AI 가 이력·문서 기억·검색", [
        ("Mem0g (graph-enhanced) 🆕", "mem0 진화 · temporal +29.6점 · time-sensitive 58% vs OpenAI 21%", "O ", "결재 이력 시계열 정확도 2.7배 (2026-06 신규)"),
        ("mem0", "LLM 장기 기억", "O ", "6개월 결재 이력 자동 로드"),
        ("MemGraphRAG", "mem0 + Graph · KDD 2026", "O ", "관계망 통합 기억"),
        ("Zep GraphRAG 🆕", "graph 순회 + vector 하이브리드 · sub-second", "O ", "Microsoft GraphRAG 대안 · 2026-07 등장"),
        ("Vector DB (Milvus·Weaviate·Qdrant)", "벡터 저장·검색 · 온프레미스", "O", "필수 인프라 · RFP 대응"),
        ("Corrective RAG", "자기 검증 RAG", "O", "환각 자동 정정"),
        ("Adaptive RAG", "상황별 RAG 전략", "O", "쿼리별 최적 검색"),
        ("Self-RAG", "스스로 검색 여부 판단", "O", "필요 시만 검색"),
        ("GraphRAG (Microsoft)", "관계망 기반 검색", "O", "부정 관계망 감지"),
        ("HyDE (가상 답 생성)", "가상 답 생성 후 검색", "검토", "특수 케이스만"),
    ]),
    ("🤖 AI 에이전트·오케스트레이션", "AI 자율 실행·협업", [
        ("ChatGPT Work 🆕", "GPT-5.6 기반 multi-hour agentic (2026-08)", "O ", "장시간 자율 감사 프로젝트"),
        ("Managed Agents (Anthropic) 🆕", "session thread streams + webhooks · effort 지정 (2026-07)", "O ", "우리 kit 이미 사용 · route.py 확장"),
        ("Sol/Terra/Luna tier (OpenAI) 🆕", "durable capability tiers (2026-08 재편)", "O", "라우팅 매트릭스 반영"),
        ("LangChain", "LLM 체인 프레임", "O", "기본 인프라"),
        ("LangGraph", "그래프 기반 에이전트 조율", "O ", "Multi-Agent 필수"),
        ("AutoGen (Microsoft)", "Multi-Agent 대화", "O", "협업 프레임"),
        ("CrewAI", "역할별 에이전트 팀", "O", "역할 분업 (감사·조사·대응)"),
        ("Computer Use (Claude)", "AI 가 마우스·키보드", "O ", "자동 대응 실행"),
        ("Browser-Use", "웹 자율 조작", "O", "감사 시스템 자동 조작"),
        ("Devin (Cognition)", "자율 개발자", "X", "RMS 아닌 개발팀 도구"),
        ("ReAct", "Reasoning + Acting", "O", "추론 + 실행 통합"),
    ]),
    ("🔍 이상탐지 (Anomaly Detection)", "정상 vs 이상 자동 감지", [
        ("Autoencoder", "압축·복원 오차로 이상 감지", "O", "결재 이상탐지"),
        ("Isolation Forest", "이상치 격리 나무", "O", "고전 안정 알고리즘"),
        ("LOF (Local Outlier Factor)", "근접도 기반", "O", "지역 이상"),
        ("One-Class SVM", "정상만 학습", "O", "부정 데이터 없어도 학습"),
        ("Contrastive Learning", "정상 vs 이상 대조 · 자폐 생쥐", "O ", "임베딩 대조 학습"),
        ("GAN-based Anomaly", "생성 오차로 이상", "O", "복잡 패턴 감지"),
        ("Diffusion-based", "노이즈 복원 오차", "검토", "성능 대비 비용"),
        ("Deep SVDD", "딥러닝 SVDD", "검토", "특수 케이스"),
    ]),
    ("⏱ 시계열 (Time Series)", "시간 흐름 데이터 이해", [
        ("LSTM", "순차 신경망 (성숙)", "O", "결재 시계열 이해"),
        ("Transformer sequence", "Attention 기반 시퀀스", "O ", "6개월 sequence 통째"),
        ("TimesFM (Google 2024)", "시계열 파운데이션 모델", "O", "사전학습 활용"),
        ("Chronos (Amazon)", "시계열 LLM", "O", "LLM 방식 시계열"),
        ("Prophet (Meta)", "트렌드 예측", "O", "간단한 트렌드"),
        ("N-BEATS", "순수 딥러닝 forecast", "검토", "특정 케이스"),
        ("Kalman Filter", "실시간 상태 추정", "O ", "실시간 위험도 curve"),
        ("Particle Filter", "비선형 상태 추정", "검토", "복잡 시스템"),
    ]),
    ("🕸 그래프 (Graph AI)", "관계망 데이터 이해", [
        ("GNN (Graph Neural Network)", "그래프 신경망", "O ", "부정 관계망 학습"),
        ("GAT (Graph Attention)", "그래프 attention", "O", "중요 관계 자동"),
        ("Node2Vec", "노드 임베딩", "O", "그래프 → 벡터"),
        ("GraphSAGE", "그래프 sampling", "O", "대규모 그래프"),
        ("Graph Transformer", "Transformer + Graph", "O", "복합 관계망"),
        ("Entity Resolution", "개체 자동 매칭", "O ", "특수관계자 자동 감지"),
    ]),
    ("👁 컴퓨터비전 (Vision)", "이미지·영상 이해", [
        ("CNN", "이미지 신경망 (기본)", "O", "이미지 감사"),
        ("YOLO", "실시간 물체 탐지", "O", "Physical AI · CCTV"),
        ("Vision Transformer (ViT)", "이미지 Transformer", "O", "고해상도 이미지"),
        ("SAM (Segment Anything)", "세그멘테이션", "O", "문서·이미지 분할"),
        ("DeepLabCut·SLEAP", "자세 추정 (동물 논문)", "O", "행동 자세 분석"),
        ("CLIP", "이미지·텍스트 매칭", "O", "문서 이해"),
        ("VLM (Vision-Language)", "이미지 이해 + 언어", "O", "이미지+텍스트 통합"),
    ]),
    ("🔐 프라이버시·암호 (Privacy·Crypto)", "데이터 보호 기술", [
        ("Homomorphic Encryption", "암호 상태 계산 (Microsoft SEAL·OpenFHE)", "O ", "감사인·규제 대응 · 데이터 안 열고 계산"),
        ("Differential Privacy", "통계에 노이즈 추가", "O", "개인정보 보호 통계"),
        ("Federated Learning", "데이터 안 옮기고 공동 학습", "O", "여러 증권사 공동 학습"),
        ("Zero-Knowledge Proof", "내용 안 보여주고 증명 (zk-SNARK)", "O ", "감사 근거 증명"),
        ("Secure Multi-Party (MPC)", "여러 파티 공동 계산", "검토", "Federated 로 충분"),
    ]),
    (" 최적화 (Optimization)", "최적해 자동 찾기", [
        ("QUBO", "이진 최적화 (양자영감)", "O ", "최적 조사 대상 선정 (v14 S16)"),
        ("Simulated Annealing", "담금질 최적화", "O", "QUBO 백엔드"),
        ("Genetic Algorithm", "유전 알고리즘", "O", "탐색 문제"),
        ("Bayesian Optimization", "확률적 탐색", "O", "하이퍼파라미터 튜닝"),
        ("Reinforcement Learning (RL)", "강화학습", "O", "자율 학습"),
        ("PPO·DPO", "RL 알고리즘", "검토", "특수 케이스"),
    ]),
    ("🌐 시뮬레이션 (Simulation)", "미래·What-if 예측", [
        ("Digital Twin", "시스템 컴퓨터 쌍둥이", "O ", "v14 Risk Twin 실장"),
        ("Monte Carlo", "확률 시뮬", "O", "위험 확률 계산"),
        ("Agent-Based Model", "개체 상호작용", "O", "조직 시뮬"),
        ("Physics-Informed NN", "물리 규칙 학습", "검토", "특수 케이스"),
    ]),
    ("📋 감사·컴플라이언스 특화", "감사·규제 대응 특화", [
        ("Process Mining (Celonis)", "ERP 로그 → 실제 프로세스", "O ", "규정 프로세스 vs 실제 이탈"),
        ("Business Rules Engine", "룰 실행 엔진", "O", "규정 자동 실행"),
        ("Explainable AI (SHAP·LIME)", "AI 판단 설명", "O ", "RFP 필수 요구"),
        ("Model Governance (Fiddler·Arize)", "AI 모델 감시", "O", "우리 AI 자체 감시"),
        ("Data Lineage", "데이터 계보 추적", "O", "K-SOX audit trail"),
        ("CBR (Case-Based Reasoning)", "과거 사례 매칭", "O", "유사 부정 사례 자동"),
        ("Benford's Law", "자연 숫자 분포 이상", "O", "회계 조작 감지"),
    ]),
    ("🔄 학습 패러다임 (Learning Paradigms)", "학습 전략·기법", [
        ("Fine-tuning", "사전학습 모델 재학습", "O", "도메인 특화"),
        ("LoRA·QLoRA", "저비용 fine-tuning", "O ", "온프레미스 저비용"),
        ("Quantization", "모델 압축", "O", "저사양 배포"),
        ("Distillation", "큰 모델 → 작은 모델", "O", "SLM 만들기"),
        ("Transfer Learning", "다른 도메인 지식 활용", "O", "산업 이식"),
        ("Few-shot / Zero-shot", "소수 데이터로 학습", "O ", "부정 사례 극소량 학습"),
        ("Active Learning", "애매한 것만 사람에 질문", "O ", "사람 개입 최소"),
        ("Semi-Supervised", "라벨 일부만", "O", "부분 라벨 활용"),
        ("Meta-Learning", "학습법 학습", "검토", "장기 R&D"),
    ]),
    ("🎨 데이터 생성 (Data Generation)", "가짜 데이터 자동 만들기", [
        ("GAN", "대립 생성", "O", "Synthetic 부정"),
        ("Diffusion Model", "노이즈 → 진짜", "O", "고품질 생성"),
        ("VAE (Variational Autoencoder)", "잠재 공간 생성", "O", "부드러운 생성"),
        ("Synthetic Data (SDV·Gretel)", "개인정보 없는 가짜 데이터", "O ", "부정 데이터 부족 해결"),
    ]),
    ("🌍 크로스도메인 (Cross-Industry)", "다른 산업에서 이식된 것", [
        ("Predictive Maintenance (RUL)", "항공·제조 원 · 부품 고장 예측", "O ", "부정 30일 전 예측"),
        ("Sensor Fusion", "자율주행 원 · 다센서 융합", "O", "ERP+이메일+CCTV 통합"),
        ("Edge Case Detection", "자율주행 원 · 훈련 밖 감지", "O", "정상 훈련 밖 결재"),
        ("AIS-LLM 궤적 분석", "해양 원 · Windward.ai 2026", "O", "결재 궤적 이상"),
        ("Flight Data Recorder", "항공 원 · 통합 기록", "O", "v14 Robot Blackbox"),
        ("Satellite Telemetry Anomaly", "위성 원 · ESA 벤치마크", "O", "결재 텔레메트리"),
    ]),
]

for idx, (title, subtitle, techs) in enumerate(CATEGORIES, 1):
    s = new_slide()
    add_header(s, title, subtitle)
    rows = [(name, desc, ox_cell(ox), effect) for name, desc, ox, effect in techs]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
              ["기술", "뭐 하는 것", "RMS", "RMS 효과"],
              rows,
              header_bg=BLUE, header_fg=WHITE,
              col_widths=[2.8, 4.0, 1.2, 4.3],
              font_size=10, header_size=11)
    o_count = sum(1 for t in techs if t[2].startswith("O"))
    x_count = sum(1 for t in techs if t[2].startswith("X"))
    review_count = sum(1 for t in techs if "검토" in t[2])
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), GRAY_BG, line=BLUE)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5),
             f"요약  ·  O: {o_count}  ·  X: {x_count}  ·  검토: {review_count}",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, f"PART 3 · 카테고리 {idx}/14")

# ─── PART 4 · Killer 19 종합 ────────────────────
s = new_slide()
add_header(s, "PART 4 · 종합 Killer 19", " RMS 에 진짜 있으면 좋은 것")
killers = [
    "1. mem0 (기억)",
    "2. MemGraphRAG (관계망 기억)",
    "3. LangGraph (에이전트 조율)",
    "4. Computer Use (자동 대응)",
    "5. Contrastive Learning (정상 vs 이상)",
    "6. Transformer sequence (긴 이력 이해)",
    "7. Kalman Filter (실시간 상태)",
    "8. GNN + Entity Resolution (특수관계자·자금세탁)",
    "9. Homomorphic Encryption (프라이버시)",
    "10. Zero-Knowledge Proof (증명)",
    "11. QUBO (최적 조사대상)",
    "12. Digital Twin (시뮬)",
    "13. Process Mining (실제 vs 규정 프로세스)",
    "14. Explainable AI (SHAP·LIME) (RFP 필수)",
    "15. LoRA·QLoRA (저비용 fine-tuning)",
    "16. Few-shot Learning (소수 부정 사례)",
    "17. Active Learning (사람 개입 최소)",
    "18. Synthetic Data (GAN) (부정 데이터 부족 해결)",
    "19. Predictive Maintenance / RUL (사전 예측)",
]
y = Inches(1.4)
for i, k in enumerate(killers):
    col = i % 3
    row = i // 3
    xx = Inches(0.5 + col * 4.25)
    yy = Inches(1.4 + row * 0.75)
    add_rect(s, xx, yy, Inches(4.15), Inches(0.65), GRAY_BG, line=BLUE)
    add_text(s, xx, yy, Inches(4.15), Inches(0.65), k,
             size=12, bold=True, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "PART 4 · Killer 19개")

# ─── SLIDE · X 안 어울림 리스트 ────────────────────
s = new_slide()
add_header(s, "PART 4 · X (RMS 안 어울림)", "왜 안 쓰나 이유 포함")
x_list = [
    ("Devin (자율 개발자)", "RMS 아님 · 우리 개발팀 자체 도구엔 O"),
    ("World Model (Physical Intelligence·Genie)", "아직 실험 단계 · 물리 세계 대상"),
    ("Mamba (State Space Model)", "LSTM/Transformer 로 충분"),
    ("MPC (Secure Multi-Party Computation)", "Federated Learning 으로 충분"),
    ("Physics-Informed NN", "물리 규칙 기반 · RMS 물리 X"),
    ("Diffusion Anomaly", "Autoencoder 로 충분"),
    ("LiDAR / 3D 스캔", "물리 센서 · RMS 물리 X"),
    ("HD Map", "지도 기반 · RMS 관계없음"),
    ("Occupancy Network", "3D 공간 점유 · 물리 X"),
    ("End-to-End Driving (Tesla FSD)", "전용 아키텍처 · RMS 별도"),
    ("실제 항공기·선박·발전소 물리", "물리 세계 · RMS X"),
    ("Speech Synthesis (TTS)", "생성만 · 감지·판단 아님"),
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
          ["X 기술", "이유"],
          [(name, reason) for name, reason in x_list],
          header_bg=RED, header_fg=WHITE,
          col_widths=[4.5, 7.8], font_size=12, header_size=13)
add_footer(s, "PART 4 · X 12개")

# ─── PART 5 · Phase 로드맵 ────────────────────
s = new_slide()
add_header(s, "PART 5 · Phase 로드맵", "0 → 24개월 실행 순서")
phases = [
    ("Phase 1 · 0~3개월 · 즉시 차별화", GREEN, [
        "① mem0 붙이기 · 6개월 이력 축적 시작",
        "② MemGraphRAG · Cen's TRM 확장",
        "③ 위성 텔레메트리 이상탐지 (94% precision)",
        "④ LangGraph 인프라 세팅",
    ]),
    ("Phase 2 · 3~6개월 · 카테고리 신설", CYAN, [
        "⑤ Agent Behavior Analytics (Exabeam 방식)",
        "⑥ AI-RPA / Computer Use (RPA 불요)",
        "⑦ 해양 AIS-LLM 궤적 분석",
        "⑧ Contrastive Learning (자폐 생쥐 방식)",
    ]),
    ("Phase 3 · 6~12개월 · 예측 패러다임", ORANGE, [
        "⑨ Predictive Maintenance (RUL) · 부정 30일 전 예측",
        "⑩ Synthetic Fraud Data (GAN)",
        "⑪ Physical AI (Voxel/Protex)",
        "⑫ QUBO (양자영감 최적 조사대상)",
        "⑬ Zero-Knowledge Proof PoC (유안타 크로스보더)",
    ]),
    ("Phase 4 · 12개월+ · 최상위 IP", NAVY, [
        "⑭ 게임 방식 Transformer sequence (6개월 통째)",
        "⑮ Federated Learning (증권사 컨소시엄)",
        "⑯ Homomorphic Encryption (감사 데이터 극한 privacy)",
        "⑰ 의료 Few-shot pathology 방식",
        "⑱ 광고 Adversarial AI fraud 방어",
        "⑲ ZK Proof audit trail (규제 강제 오면 즉시)",
    ]),
]
y = Inches(1.4)
for phase_title, color, items in phases:
    add_rect(s, Inches(0.5), y, Inches(3.0), Inches(1.4), color)
    add_text(s, Inches(0.5), y, Inches(3.0), Inches(1.4), phase_title,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(3.6), y, Inches(9.5), Inches(1.4), GRAY_BG, line=color)
    text = "\n".join(items)
    tb = s.shapes.add_textbox(Inches(3.7), y, Inches(9.3), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.05)
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.name = "맑은 고딕"; r.font.size = Pt(11); r.font.color.rgb = BLACK
    y += Inches(1.5)
add_footer(s, "PART 5 · 로드맵")

# ─── PART 6 · KPMG 대응 ────────────────────
s = new_slide()
add_header(s, "PART 6 · KPMG 대응 전략", "직접 대결 X · 시장 분리")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
         "나쁜 소식 · KPMG 5년 $2B · $12B 매출 목표 · 감사업무 완전 자동화 중",
         size=14, bold=True, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.5),
         "좋은 소식 · KPMG = 감사법인 (서비스) · 우리 = 감사 대상 회사에 파는 SaaS (다른 시장)",
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

matrix = [
    ("온프레미스", "클라우드 중심", "우리 kit exec_offline · RFP 요구"),
    ("실시간·상시", "감사 기간만", "24/7 상시 · Bible 행위 organ"),
    ("한국 규제 특화", "글로벌 표준", "K-SOX·자금부정공시·중대재해법"),
    ("Physical AI 통합", "없음", "v14 Robot Blackbox"),
    ("감사인 협업 툴", "자기 내부", "회사 → 감사인 준비 도구 (B2B)"),
    ("금융권 온프레미스", "클라우드만", "증권사·은행 RFP 대응"),
    ("한국어 OCR + HWP", "지원 미미", "국내 정부·공기업 필수"),
    ("ISMS-P 통합", "없음", "사용자 개인 3년 강점"),
]
add_table(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2),
          ["축", "KPMG", "우리 진입 지점"],
          matrix,
          header_bg=NAVY, header_fg=WHITE,
          col_widths=[2.5, 3.5, 6.3], font_size=11, header_size=12)

add_rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35), NAVY)
add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
         "결론 · KPMG 대체 X · KPMG 감사 대비 SaaS · 시장 분리 + 파트너 관계",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "PART 6 · KPMG 대응")

# ─── PART 7 · AI ISMS-P Copilot ────────────────────
s = new_slide()
add_header(s, "PART 7 · AI ISMS-P Copilot", "사용자 개인 강점 3년 도메인 × AI = 국내 유일")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
         " 핵심 인사이트 · K-SOX 감사 시장 열림 · ISMS-P 정보보호 인증 시장 열림 · 둘 다 매년 심사 · 매년 사람 손 · 매년 지침",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
data = [
    ("ISMS-P 인증 유지 회사", "4000+ (상장사·금융·공공 대부분)"),
    ("매년 심사 대응 비용", "회사당 5천만 ~ 3억"),
    ("지금 대응 방법", "사람 · Excel · 반복 노동"),
    ("AI 자동화 시장 규모", "2000억+ (국내 아직 없음)"),
    ("우리 kit 재활용", "plugins/ai_rag/ + plugins/exec_offline/ + approval-gate 즉시 활용"),
    ("MicroICM 재활용", "K-SOX 논리 유사 · 그대로 확장"),
    ("KPMG 못 하는 이유", "감사법인 = 정보보호 인증 없음 · 다른 시장 (KISA·개인정보위)"),
    ("사용자 유일 강점", "ISMS-P 3년 실무 · 통제·감사·심사 뉘앙스 국내 몇 안 됨"),
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.5),
          ["항목", "내용"],
          data,
          header_bg=BLUE, header_fg=WHITE,
          col_widths=[3.5, 8.8], font_size=12, header_size=13)

add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.55), GREEN)
add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.55),
         "이번 주 액션 · ① 오늘 밤 종이 1장 3줄 정리 · ② 내일 ISMS-P 심사원 1명에게 문자 · ③ 답 오면 시작",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "PART 7 · ISMS-P Copilot")

# ─── PART 8 · Bible 매핑 + 다음 액션 ────────────────────
s = new_slide()
add_header(s, "PART 8 · Bible 13 organs 매핑", "지금 RMS 사용 organ · 미개척 organ")
mapping = [
    ("뇌 (Reasoning)", "부분", "Multi-Agent 붙이면 강화"),
    ("기억 (Memory)", "", "mem0 넣으면 실현 (Phase 1)"),
    ("생성 (Generation)", "부분", "감사 조서 자동 생성"),
    ("연결 (Connectivity)", "", "MemGraphRAG + GNN (Phase 1)"),
    ("지각 (Perception)", "", "AI CCTV · 부서 IP 사용"),
    ("추론 (Reasoning)", "부분", "Test-Time Compute 붙이면"),
    ("에너지 (Energy)", "", "SLM + vLLM (Phase 2)"),
    ("신뢰 (Trust)", "", "Evidence Ledger · Audit Trail"),
    ("행위 (Agency)", "", "Agentic AI · Computer Use"),
    ("체화 (Embodiment)", "부분", "Physical AI Safety RMS (Phase 3)"),
    ("시뮬 (Simulation)", "", "Digital Twin · Risk Twin (Phase 3)"),
    ("AI 문명", "", "Federated Learning (Phase 4)"),
    ("진화 (Evolution)", "", "Self-Improving AutoCheck (Phase 4)"),
]
rows = []
for name, usage, action in mapping:
    if usage == "":
        u = ("사용 중", GREEN)
    elif usage == "부분":
        u = ("부분", ORANGE)
    else:
        u = ("미사용", RED)
    rows.append((name, u, action))

add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.3),
          ["Bible organ", "지금 RMS", "다음 액션"],
          rows,
          header_bg=NAVY, header_fg=WHITE,
          col_widths=[3.5, 2.5, 6.3], font_size=11, header_size=12)
add_rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35), NAVY)
add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.35),
         "현재 사용: 3 organ (지각·신뢰·행위) · 미개척: 6 organ · 부분: 4 organ · Phase 1~4 로 다 채움",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, "PART 8 · Bible 매핑")

# ─── 마지막 · 다음 액션 요약 ────────────────────
s = new_slide()
add_bg(s, NAVY)
add_text(s, Inches(0.8), Inches(0.8), Inches(11.8), Inches(0.9),
         "다음 액션 (사용자 선택)",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.9), Inches(11.8), Inches(0.5),
         "이 40+ 슬라이드 중 spark 되는 것 짚어주세요",
         size=16, color=CYAN, align=PP_ALIGN.CENTER)

actions = [
    ("A", "산업 하나 심층", "16 산업 중 하나 spark → 3배 depth (구현 코드·아키텍처·비용)"),
    ("B", "기술 카테고리 심층", "14 카테고리 중 하나 spark → mechanism 상세"),
    ("C", "Phase 1 실행", "mem0 · MemGraphRAG · LangGraph 즉시 시작 (3개월)"),
    ("D", "AI ISMS-P Copilot", "사용자 개인 강점 · 이번 주 검증 (문자 3명)"),
    ("E", "유안타 RFP 응답", "K-SOX + ISMS-P + Physical AI 종합 대응"),
    ("F", "새 산업 스캔", "부동산·헬스케어·법률·게임 등 새 도메인 요청"),
]
y = Inches(2.6)
for tag, title, desc in actions:
    add_rect(s, Inches(0.8), y, Inches(0.7), Inches(0.55), YELLOW)
    add_text(s, Inches(0.8), y, Inches(0.7), Inches(0.55), tag,
             size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.6), y, Inches(2.6), Inches(0.55), CYAN)
    add_text(s, Inches(1.6), y, Inches(2.6), Inches(0.55), title,
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(4.3), y, Inches(8.4), Inches(0.55), WHITE)
    add_text(s, Inches(4.3), y, Inches(8.4), Inches(0.55), desc,
             size=11, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)

add_text(s, Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5),
         "아이티센 코어 ESG 사업부  ·  2026-07-01  ·  MicroRisk-X",
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ─── 저장 ────────────────────
out = Path(__file__).parent / "산업_ML_RMS_이식_종합지도.pptx"
if out.exists():
    bak = out.with_suffix(out.suffix + ".bak")
    if bak.exists():
        bak.unlink()
    out.rename(bak)
prs.save(str(out))
print(f" 생성 완료: {out}")
print(f"   슬라이드: {len(prs.slides)}장")
