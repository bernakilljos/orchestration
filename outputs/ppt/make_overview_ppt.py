#!/usr/bin/env python3
"""
Orchestration Kit v3 — 전체 개요 PPT 생성
(기존 ai-system-stages-2026-04-16.pptx 대체)

내용:
  1. 커버
  2. 프로젝트 정체성 (킷 vs 플랫폼)
  3. 아키텍처 — exec_orch 코어 + 플러그인 응용 레이어
  4. 21개 플러그인 카탈로그 (카테고리별)
  5. 핵심 플러그인 5개 상세
  6. 멀티AI 라우팅 흐름
  7. 로드맵 Phase 1~3
  8. 사용 시나리오

생성물:
  docs/screens/*.jpg (5개 다이어그램)
  outputs/ppt/orchestration-overview-2026-04-19.pptx
"""
import sys, os
from pathlib import Path

try: sys.stdout.reconfigure(encoding='utf-8')
except: pass

# 한글 폰트
import matplotlib
import matplotlib.pyplot as plt
from matplotlib import patches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import matplotlib.font_manager as fm

# Windows 한글 폰트
for font in ['Malgun Gothic', 'NanumGothic', 'Apple SD Gothic Neo', 'Noto Sans CJK KR']:
    if any(font.lower() in f.name.lower() for f in fm.fontManager.ttflist):
        matplotlib.rcParams['font.family'] = font
        break
matplotlib.rcParams['axes.unicode_minus'] = False

SCREENS = Path("docs/screens")
SCREENS.mkdir(parents=True, exist_ok=True)

# ─────────────────────────────────────────────────────
# 다이어그램 1: 킷 vs 플랫폼 구조
# ─────────────────────────────────────────────────────
def make_kit_vs_platform():
    fig, ax = plt.subplots(figsize=(14, 8), dpi=120)
    ax.set_xlim(0, 14); ax.set_ylim(0, 8)
    ax.axis('off')

    # 킷 영역 (왼쪽)
    ax.add_patch(FancyBboxPatch((0.3, 0.5), 6.2, 7, boxstyle="round,pad=0.1",
                                 edgecolor='#2E86AB', facecolor='#EAF4FB', linewidth=2.5))
    ax.text(3.4, 7.0, "orchestration_v1 (KIT)", fontsize=18, fontweight='bold', ha='center', color='#2E86AB')
    ax.text(3.4, 6.5, "— 이 저장소 —", fontsize=11, ha='center', color='#666')

    kit_items = [
        "exec_orch 코어 엔진",
        "플러그인 스펙·템플릿",
        "sync 시스템",
        "설계 문서·로드맵",
        "공유 규칙 (.claude/rules/)",
        "인프라 스크립트",
    ]
    for i, item in enumerate(kit_items):
        y = 5.8 - i * 0.7
        ax.add_patch(FancyBboxPatch((0.8, y), 5.2, 0.55, boxstyle="round,pad=0.05",
                                     edgecolor='#2E86AB', facecolor='white', linewidth=1))
        ax.text(3.4, y + 0.28, item, fontsize=11, ha='center')

    # 플랫폼 영역 (오른쪽)
    ax.add_patch(FancyBboxPatch((7.5, 0.5), 6.2, 7, boxstyle="round,pad=0.1",
                                 edgecolor='#A23B72', facecolor='#F9ECEF', linewidth=2.5))
    ax.text(10.6, 7.0, "설치 후 플랫폼들", fontsize=18, fontweight='bold', ha='center', color='#A23B72')
    ax.text(10.6, 6.5, "— 각자 프로젝트 —", fontsize=11, ha='center', color='#666')

    platform_items = [
        "수익화 프로젝트 (YouTube·IG)",
        "본업 (Micropolis AWS)",
        "실제 API 호출 구현",
        "비즈니스 로직",
        "프로덕션 트래픽 처리",
        "도메인별 커스터마이징",
    ]
    for i, item in enumerate(platform_items):
        y = 5.8 - i * 0.7
        ax.add_patch(FancyBboxPatch((8.0, y), 5.2, 0.55, boxstyle="round,pad=0.05",
                                     edgecolor='#A23B72', facecolor='white', linewidth=1))
        ax.text(10.6, y + 0.28, item, fontsize=11, ha='center')

    # install 화살표
    arrow = FancyArrowPatch((6.6, 4), (7.4, 4), arrowstyle='->', mutation_scale=30,
                            color='#F18F01', linewidth=3)
    ax.add_patch(arrow)
    ax.text(7.0, 4.5, "install", fontsize=12, fontweight='bold', ha='center', color='#F18F01')

    plt.title("프로젝트 정체성 — 킷 vs 플랫폼",
              fontsize=20, fontweight='bold', pad=20)
    plt.savefig(SCREENS / "01-kit-vs-platform.jpg", dpi=120, bbox_inches='tight', facecolor='white')
    plt.close()
    print("OK: 01-kit-vs-platform.jpg")

# ─────────────────────────────────────────────────────
# 다이어그램 2: exec_orch 코어 + 응용 레이어
# ─────────────────────────────────────────────────────
def make_architecture():
    fig, ax = plt.subplots(figsize=(14, 9), dpi=120)
    ax.set_xlim(0, 14); ax.set_ylim(0, 9)
    ax.axis('off')

    # 코어 (가운데)
    ax.add_patch(FancyBboxPatch((5, 3.5), 4, 2, boxstyle="round,pad=0.15",
                                 edgecolor='#C73E1D', facecolor='#FFE5E0', linewidth=3))
    ax.text(7, 5.0, "exec_orch", fontsize=22, fontweight='bold', ha='center', color='#C73E1D')
    ax.text(7, 4.4, "멀티AI 오케스트레이션 엔진", fontsize=12, ha='center')
    ax.text(7, 3.9, "(Claude + Codex + Gemini)", fontsize=10, ha='center', color='#666')

    # 응용 플러그인 카테고리
    categories = [
        ("exec_", 4, 7.5, '#2E86AB', "실행·학습·세션·음성"),
        ("mcp_",  10, 7.5, '#A23B72', "MCP 서버 설치 허브"),
        ("design_", 10, 1.5, '#F18F01', "Excel·PPT·Word·Web·PDF·Video"),
        ("review_", 4, 1.5, '#57A773', "코드·보안·성능·테스트"),
    ]
    for prefix, x, y, color, desc in categories:
        ax.add_patch(FancyBboxPatch((x-1.5, y-0.6), 3, 1.2, boxstyle="round,pad=0.08",
                                     edgecolor=color, facecolor='white', linewidth=2))
        ax.text(x, y+0.15, prefix + "*", fontsize=14, fontweight='bold', ha='center', color=color)
        ax.text(x, y-0.3, desc, fontsize=9, ha='center')

        # 코어로 화살표
        dx = 7 - x
        dy = 4.5 - y
        arrow = FancyArrowPatch((x, y-0.6 if y > 4 else y+0.6),
                                (7 + (1.5 if dx > 0 else -1.5) * 0.3,
                                 4.5 + (1 if dy > 0 else -1) * 0.5),
                                arrowstyle='->', mutation_scale=15,
                                color=color, linewidth=1.5, linestyle='--')
        ax.add_patch(arrow)

    # 하단 인프라
    ax.add_patch(FancyBboxPatch((1, 0.2), 12, 0.8, boxstyle="round,pad=0.05",
                                 edgecolor='#333', facecolor='#F5F5F5', linewidth=1.5))
    ax.text(7, 0.6, "sync-plugins.sh  |  validate-schema.py  |  install.sh  |  orca-status.sh  |  worker-health.sh",
            fontsize=10, ha='center', family='monospace')

    # 상단 문서
    ax.add_patch(FancyBboxPatch((1, 8.0), 12, 0.7, boxstyle="round,pad=0.05",
                                 edgecolor='#333', facecolor='#EFEFEF', linewidth=1.5))
    ax.text(7, 8.35, "CLAUDE.md  |  guide.txt  |  .claude/rules/  |  docs/architecture-patterns.md",
            fontsize=10, ha='center', family='monospace')

    plt.title("아키텍처 — exec_orch 코어 + 4개 응용 카테고리",
              fontsize=20, fontweight='bold', pad=20)
    plt.savefig(SCREENS / "02-architecture.jpg", dpi=120, bbox_inches='tight', facecolor='white')
    plt.close()
    print("OK: 02-architecture.jpg")

# ─────────────────────────────────────────────────────
# 다이어그램 3: 21개 플러그인 카탈로그
# ─────────────────────────────────────────────────────
def make_catalog():
    fig, ax = plt.subplots(figsize=(16, 10), dpi=120)
    ax.set_xlim(0, 16); ax.set_ylim(0, 10)
    ax.axis('off')

    groups = {
        "exec_ (5)": {
            "color": "#2E86AB",
            "items": ["exec_orch ⭐", "exec_learning", "exec_session_guard", "exec_voice", "exec_scheduler 📝"],
            "x": 0.5, "y_top": 8.5,
        },
        "mcp_ (8)": {
            "color": "#A23B72",
            "items": ["mcp_dev", "mcp_collab", "mcp_data", "mcp_docs",
                      "mcp_media", "mcp_web", "mcp_social 📝", "mcp_queue 📝"],
            "x": 4.3, "y_top": 8.5,
        },
        "design_ (6)": {
            "color": "#F18F01",
            "items": ["design_excel", "design_ppt", "design_word",
                      "design_web 📝", "design_pdf 📝", "design_video 📝"],
            "x": 8.1, "y_top": 8.5,
        },
        "review_ (1)": {
            "color": "#57A773",
            "items": ["review_qa"],
            "x": 11.9, "y_top": 8.5,
        },
        "cost_ (1 spec)": {
            "color": "#C73E1D",
            "items": ["cost_youtube 📝"],
            "x": 11.9, "y_top": 5.0,
        },
    }

    for gname, g in groups.items():
        ax.add_patch(FancyBboxPatch((g["x"], g["y_top"] - len(g["items"]) * 0.45 - 0.5),
                                    3.5, len(g["items"]) * 0.45 + 0.7,
                                    boxstyle="round,pad=0.1",
                                    edgecolor=g["color"], facecolor='white', linewidth=2))
        ax.text(g["x"] + 1.75, g["y_top"] + 0.1, gname, fontsize=14, fontweight='bold',
                ha='center', color=g["color"])
        for i, item in enumerate(g["items"]):
            y = g["y_top"] - 0.3 - i * 0.45
            ax.text(g["x"] + 0.2, y, "•  " + item, fontsize=10, ha='left')

    ax.text(8, 0.8, "Total: 21개 (14 stable + 7 spec-only 📝)   ⭐ = 코어",
            fontsize=14, fontweight='bold', ha='center', color='#333',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#FFFFE0', edgecolor='#888'))

    plt.title("21개 플러그인 카탈로그 (카테고리별)",
              fontsize=22, fontweight='bold', pad=20)
    plt.savefig(SCREENS / "03-catalog.jpg", dpi=120, bbox_inches='tight', facecolor='white')
    plt.close()
    print("OK: 03-catalog.jpg")

# ─────────────────────────────────────────────────────
# 다이어그램 4: 멀티AI 라우팅 흐름
# ─────────────────────────────────────────────────────
def make_routing():
    fig, ax = plt.subplots(figsize=(14, 9), dpi=120)
    ax.set_xlim(0, 14); ax.set_ylim(0, 9)
    ax.axis('off')

    # 사용자 요청 (상단)
    ax.add_patch(FancyBboxPatch((5.5, 8), 3, 0.8, boxstyle="round,pad=0.1",
                                 edgecolor='#333', facecolor='#FFF2CC', linewidth=2))
    ax.text(7, 8.4, "사용자 요청", fontsize=14, fontweight='bold', ha='center')

    # route_dispatch (판단)
    ax.add_patch(FancyBboxPatch((5, 6.3), 4, 1.2, boxstyle="round,pad=0.1",
                                 edgecolor='#C73E1D', facecolor='#FFE5E0', linewidth=2.5))
    ax.text(7, 7.1, "route_dispatch v2", fontsize=13, fontweight='bold', ha='center', color='#C73E1D')
    ax.text(7, 6.6, "규모 판단 + AI 매트릭스 + 단가", fontsize=10, ha='center')

    arrow = FancyArrowPatch((7, 8), (7, 7.5), arrowstyle='->', mutation_scale=20, color='#333')
    ax.add_patch(arrow)

    # 3 경로
    paths = [
        (2, 4.5, "Claude", "설계·판단·승인\n복잡 추론", "#1F77B4"),
        (7, 4.5, "Codex × 4", "대량 구현 (500줄+)\n병렬 워커", "#FF7F0E"),
        (12, 4.5, "Gemini × 2", "검증·문서·리서치\n저단가", "#2CA02C"),
    ]
    for x, y, name, desc, color in paths:
        ax.add_patch(FancyBboxPatch((x-1.5, y-0.7), 3, 1.4, boxstyle="round,pad=0.1",
                                     edgecolor=color, facecolor='white', linewidth=2.5))
        ax.text(x, y+0.35, name, fontsize=14, fontweight='bold', ha='center', color=color)
        ax.text(x, y-0.3, desc, fontsize=10, ha='center')

        # 화살표
        arrow = FancyArrowPatch((7, 6.3), (x, y+0.7), arrowstyle='->',
                                 mutation_scale=15, color=color, linewidth=1.5)
        ax.add_patch(arrow)

    # 통합 결과
    ax.add_patch(FancyBboxPatch((5, 2), 4, 1.2, boxstyle="round,pad=0.1",
                                 edgecolor='#333', facecolor='#EAF4FB', linewidth=2))
    ax.text(7, 2.8, "Claude 최종 통합·보완", fontsize=13, fontweight='bold', ha='center')
    ax.text(7, 2.3, "codex 결과 + gemini 검증 → 판단", fontsize=10, ha='center')

    for x in [2, 7, 12]:
        arrow = FancyArrowPatch((x, 3.8), (7, 3.3), arrowstyle='->',
                                 mutation_scale=15, color='#666', linewidth=1, linestyle='--')
        ax.add_patch(arrow)

    # Quota Fallback
    ax.add_patch(FancyBboxPatch((1, 0.3), 12, 0.8, boxstyle="round,pad=0.05",
                                 edgecolor='#F18F01', facecolor='#FFF4E6', linewidth=1.5))
    ax.text(7, 0.7, "⚠ API 한도 초과 시 → codex-quota-check.sh → Claude 직접 모드 자동 fallback (TTL 3h)",
            fontsize=11, ha='center', fontweight='bold', color='#F18F01')

    plt.title("멀티AI 라우팅 — route_dispatch v2",
              fontsize=22, fontweight='bold', pad=20)
    plt.savefig(SCREENS / "04-routing.jpg", dpi=120, bbox_inches='tight', facecolor='white')
    plt.close()
    print("OK: 04-routing.jpg")

# ─────────────────────────────────────────────────────
# 다이어그램 5: 로드맵 Phase 1~3
# ─────────────────────────────────────────────────────
def make_roadmap():
    fig, ax = plt.subplots(figsize=(16, 9), dpi=120)
    ax.set_xlim(0, 16); ax.set_ylim(0, 9)
    ax.axis('off')

    phases = [
        {
            "name": "Phase 0 — 완성",
            "period": "현재",
            "count": "14 stable",
            "color": "#57A773",
            "items": ["exec_orch, exec_learning, exec_session_guard, exec_voice",
                      "mcp_dev, mcp_collab, mcp_data, mcp_docs, mcp_media, mcp_web",
                      "design_excel, design_ppt, design_word",
                      "review_qa"],
            "x": 0.3,
        },
        {
            "name": "Phase 1 — 스펙",
            "period": "1-2개월",
            "count": "4 spec",
            "color": "#F18F01",
            "items": ["exec_scheduler ⭐⭐⭐ (크론·DAG)",
                      "mcp_social ⭐⭐⭐ (YouTube·IG API)",
                      "cost_youtube ⭐⭐⭐ (research·upload·analytics)",
                      "design_web ⭐⭐⭐ (landing·blog·portfolio)"],
            "x": 4.3,
        },
        {
            "name": "Phase 2 — 확장",
            "period": "3-6개월",
            "count": "3 spec + 12 계획",
            "color": "#A23B72",
            "items": ["design_pdf, design_video, mcp_queue (spec 생성)",
                      "cost_instagram, cost_blog, cost_affiliate, cost_saas, cost_dashboard",
                      "perf_monitor, perf_loadtest, perf_optimize",
                      "mcp_payment, mcp_analytics, mcp_ai",
                      "growth_ab, growth_funnel"],
            "x": 8.3,
        },
        {
            "name": "Phase 3 — 본업",
            "period": "6-12개월",
            "count": "10 계획",
            "color": "#C73E1D",
            "items": ["ai_* (RAG·벡터DB·파인튜닝)",
                      "sec_* (ISMS-P·시크릿·감사)",
                      "infra_* (AWS·Terraform·IAM)",
                      "data_* (ETL·레이크·크롤링)",
                      "edu_·test_·db_·life_·biz_·i18n_"],
            "x": 12.3,
        },
    ]

    for ph in phases:
        ax.add_patch(FancyBboxPatch((ph["x"], 0.8), 3.5, 7.5,
                                    boxstyle="round,pad=0.1",
                                    edgecolor=ph["color"], facecolor='white', linewidth=2.5))
        ax.text(ph["x"] + 1.75, 8.0, ph["name"], fontsize=14, fontweight='bold',
                ha='center', color=ph["color"])
        ax.text(ph["x"] + 1.75, 7.5, ph["period"], fontsize=10, ha='center', color='#666')
        ax.text(ph["x"] + 1.75, 7.1, ph["count"], fontsize=11, fontweight='bold',
                ha='center', color=ph["color"])

        for i, item in enumerate(ph["items"]):
            y = 6.4 - i * 0.9
            # 긴 텍스트 줄바꿈
            wrapped = []
            words = item.split(", ")
            cur = ""
            for w in words:
                if len(cur + w) < 35:
                    cur = cur + ", " + w if cur else w
                else:
                    wrapped.append(cur)
                    cur = w
            if cur:
                wrapped.append(cur)
            for j, line in enumerate(wrapped):
                ax.text(ph["x"] + 0.15, y - j * 0.3, "• " + line, fontsize=8, ha='left')

    # 누적 표시
    ax.text(8, 0.3, "누적: 14 → 18 → 30 → 40 플러그인",
            fontsize=13, fontweight='bold', ha='center', color='#333',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#FFFFE0', edgecolor='#888'))

    plt.title("로드맵 — Phase 0~3 (40개 플러그인 목표)",
              fontsize=22, fontweight='bold', pad=20)
    plt.savefig(SCREENS / "05-roadmap.jpg", dpi=120, bbox_inches='tight', facecolor='white')
    plt.close()
    print("OK: 05-roadmap.jpg")

# ─────────────────────────────────────────────────────
# PPT 생성
# ─────────────────────────────────────────────────────
def make_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    tf = title.text_frame
    tf.text = "Multi-AI Orchestration Kit v3"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(54)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xAB)

    sub = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1.5))
    tf = sub.text_frame
    tf.text = "Claude + Codex + Gemini 오케스트레이션 킷\n21개 플러그인 · SoT sync · Quota Fallback"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(22)
            r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    date = s.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
    tf = date.text_frame
    tf.text = "2026-04-19 · bernakilljos"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(16)

    # Slide 2: 정체성
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "1. 프로젝트 정체성"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    s.shapes.add_picture(str(SCREENS / "01-kit-vs-platform.jpg"),
                          Inches(1.5), Inches(1.3), height=Inches(7))

    # Slide 3: 아키텍처
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "2. 아키텍처 — 코어 + 응용 레이어"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    s.shapes.add_picture(str(SCREENS / "02-architecture.jpg"),
                          Inches(0.5), Inches(1.3), height=Inches(7.5))

    # Slide 4: 카탈로그
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "3. 21개 플러그인 카탈로그"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    s.shapes.add_picture(str(SCREENS / "03-catalog.jpg"),
                          Inches(0.5), Inches(1.3), height=Inches(7.5))

    # Slide 5: 핵심 플러그인 상세
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "4. 핵심 플러그인 5개 — 용도 & 사용법"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True

    details = [
        ("exec_orch", "#C73E1D", "멀티AI 오케스트레이션 엔진 (코어)",
         "/check-agents  워커 상태\n/status  통합 대시보드\n/godmode  공격 실행"),
        ("review_qa", "#57A773", "코드 리뷰·보안·품질·테스트",
         "/check  종합\n/security  OWASP\n/performance  CWV\n/screenshot  UI 회귀"),
        ("design_ppt", "#F18F01", "PPT 자동 생성 (Canva·Figma·Gamma)",
         "/make-ppt \"주제\" 20\n/ai-system-stages\n/ppt-install"),
        ("exec_session_guard", "#2E86AB", "토큰·세션 안전 (스냅샷 자동)",
         "/guard-save  즉시 저장\n/token-stats --today\nStop/PreCompact 훅 자동"),
        ("exec_voice", "#A23B72", "음성 STT·TTS·회의록",
         "/transcribe file\n/speak text\n/meeting record.wav\n/voice-task"),
    ]
    y = 1.3
    for name, color, desc, usage in details:
        # name box
        box = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3), Inches(1.4))
        tf = box.text_frame
        tf.text = name
        r = tf.paragraphs[0].runs[0]
        r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(color.lstrip('#'))
        tf.add_paragraph().text = desc
        tf.paragraphs[-1].runs[0].font.size = Pt(10)
        tf.paragraphs[-1].runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # usage box
        box2 = s.shapes.add_textbox(Inches(3.8), Inches(y), Inches(11.5), Inches(1.4))
        tf2 = box2.text_frame
        for i, line in enumerate(usage.split("\n")):
            if i == 0:
                tf2.text = line
            else:
                tf2.add_paragraph().text = line
            tf2.paragraphs[-1].runs[0].font.size = Pt(12)
            tf2.paragraphs[-1].runs[0].font.name = "Consolas"
        y += 1.45

    # Slide 6: 라우팅
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "5. 멀티AI 라우팅 v2"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    s.shapes.add_picture(str(SCREENS / "04-routing.jpg"),
                          Inches(1), Inches(1.3), height=Inches(7.5))

    # Slide 7: 로드맵
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "6. 로드맵 — 40개 플러그인 목표"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True
    s.shapes.add_picture(str(SCREENS / "05-roadmap.jpg"),
                          Inches(0.5), Inches(1.3), height=Inches(7.5))

    # Slide 8: 사용 시나리오
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    tb.text_frame.text = "7. 사용 시나리오 — 실전 예시"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True

    scenarios = [
        ("새 기능 개발 (LARGE)",
         "1. Claude 가 task-instruction.md 작성\n2. Codex 4대 병렬 구현\n3. Claude 보완\n4. Gemini 2대 검증\n5. 배포"),
        ("빠른 버그 수정 (SMALL)",
         "1. Claude 직접 구현 (task 파일 불필요)\n2. /check 로 빠른 검증"),
        ("문서·PPT 생성",
         "1. /make-ppt \"제목\" 슬라이드수\n2. Canva + Mermaid + Figma 자동 삽입\n3. outputs/ppt/ 저장"),
        ("API 한도 초과",
         "1. codex-quota-check.sh 자동 감지\n2. .claude/state/codex-quota-exceeded 플래그\n3. Claude 직접 fallback (TTL 3h)\n4. 수동 해제: --clear"),
    ]
    y = 1.5
    for title, steps in scenarios:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4.5), Inches(1.7))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        tb.text_frame.paragraphs[0].runs[0].font.bold = True
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xAB)

        tb2 = s.shapes.add_textbox(Inches(5.2), Inches(y), Inches(10), Inches(1.7))
        tf = tb2.text_frame
        for i, line in enumerate(steps.split("\n")):
            if i == 0:
                tf.text = line
            else:
                tf.add_paragraph().text = line
            tf.paragraphs[-1].runs[0].font.size = Pt(12)
        y += 1.85

    # Slide 9: 마무리
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title = s.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    tf = title.text_frame
    tf.text = "GitHub · github.com/bernakilljos/orchestration"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True

    meta = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(3))
    tf = meta.text_frame
    for line in [
        "현재 커밋: 38b1a4d (main)",
        "커밋 수 (오늘): 8개",
        "플러그인: 21개 (14 stable + 7 spec)",
        "스크립트: 17개 (sync·validate·install·worker-health 등)",
        "공유 규칙: 5개 (.claude/rules/)",
        "설계 문서: 4개 (로드맵·플러그인·architecture-patterns·guide)",
        "",
        "실제 구현은 install 후 플랫폼에서",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(16)

    out = Path("outputs/ppt/orchestration-overview-2026-04-19.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"\n✓ PPT 저장: {out}")
    print(f"  슬라이드: {len(prs.slides)} 장")

if __name__ == "__main__":
    make_kit_vs_platform()
    make_architecture()
    make_catalog()
    make_routing()
    make_roadmap()
    make_ppt()
