"""
AI 시스템 진화 6단계 PPT 생성 스크립트
출력: outputs/ppt/ai-system-stages-2026-04-16.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── 색상 팔레트 (Stratos 스타일) ──────────────────────
BG_DARK    = RGBColor(0x0A, 0x0E, 0x27)   # 배경 (딥 네이비)
BG_CARD    = RGBColor(0x13, 0x1A, 0x3A)   # 카드 배경
ACCENT1    = RGBColor(0x4A, 0x9F, 0xFF)   # 파란 강조
ACCENT2    = RGBColor(0xFF, 0x6B, 0x35)   # 오렌지 강조 (5단계)
ACCENT3    = RGBColor(0x00, 0xE5, 0xFF)   # 시안 강조 (6단계)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY  = RGBColor(0xB0, 0xBE, 0xD1)
TEXT_LIGHT = RGBColor(0xE8, 0xF0, 0xFF)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # 완전 빈 레이아웃


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=BG_DARK):
    """배경 채우기"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        bg_color=BG_CARD, border=None):
    """사각형 박스"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    """텍스트 박스"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def heading(slide, text, sub=None, accent=ACCENT1):
    """슬라이드 상단 헤딩"""
    # 상단 라인
    line = slide.shapes.add_shape(1,
        Inches(0), Inches(0), W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    txt(slide, text, 0.5, 0.2, 12, 0.8,
        size=32, bold=True, color=TEXT_WHITE)
    if sub:
        txt(slide, sub, 0.5, 0.95, 12, 0.4,
            size=14, color=TEXT_GRAY)


def badge(slide, text, left, top, color=ACCENT1):
    """작은 뱃지"""
    b = box(slide, left, top, len(text)*0.13+0.3, 0.32,
            bg_color=color)
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = TEXT_WHITE


# ═══════════════════════════════════════════════
# 슬라이드 1 — 표지
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)

# 중앙 글로우 효과 (원)
for i, (size, alpha) in enumerate([(6,0.05),(4.5,0.08),(3,0.12)]):
    c = slide = s
    shape = s.shapes.add_shape(9,  # OVAL
        Inches(6.665-size/2), Inches(3.75-size/2),
        Inches(size), Inches(size))
    shape.fill.solid()
    shade = 0x0A + i*0x08
    shape.fill.fore_color.rgb = RGBColor(shade, shade+0x10, shade+0x30)
    shape.line.fill.background()

txt(s, "AI 시스템 진화", 1, 1.8, 11.33, 1.2,
    size=54, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txt(s, "6단계", 1, 2.9, 11.33, 1.0,
    size=72, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
txt(s, "Prompt에서 Platform까지  —  차원이 바뀌는 순간들", 1, 4.1, 11.33, 0.6,
    size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# 하단 라인
line = s.shapes.add_shape(1, Inches(3), Inches(5.0), Inches(7.33), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT1; line.line.fill.background()

txt(s, "2026", 1, 5.4, 11.33, 0.5,
    size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 2 — 전체 로드맵
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "전체 로드맵", "1단계에서 6단계로 — 각 단계는 이전을 포함하며 확장된다")

stages = [
    ("1️⃣", "Prompt",        "질문",     ACCENT1,  1.0),
    ("2️⃣", "Agent",         "역할",     ACCENT1,  3.0),
    ("3️⃣", "Orchestration", "연결",     ACCENT1,  5.0),
    ("4️⃣", "Automation",    "자동화",   ACCENT1,  7.0),
    ("🔥", "Autonomous",    "자율",     ACCENT2,  9.0),
    ("🔥", "Platform",      "생태계",   ACCENT3, 11.0),
]
for emoji, name, keyword, color, left in stages:
    b = box(s, left, 1.6, 2.1, 3.5, bg_color=BG_CARD, border=color)
    txt(s, emoji,   left+0.1, 1.7,  1.9, 0.6, size=28, align=PP_ALIGN.CENTER)
    txt(s, name,    left+0.1, 2.35, 1.9, 0.6, size=15, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    txt(s, keyword, left+0.1, 2.95, 1.9, 0.4, size=13, color=color, align=PP_ALIGN.CENTER)

# 화살표
for i, (_, _, _, _, left) in enumerate(stages[:-1]):
    txt(s, "→", left+2.15, 2.95, 0.7, 0.4, size=20, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "★ 현재: 5단계 완성 단계  |  목표: 6단계 Platform", 0.5, 5.5, 12.33, 0.4,
    size=14, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 3 — 1단계 Prompt
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "1단계 — Prompt (프롬프트)", "그냥 질문하는 단계")
badge(s, "1단계", 0.5, 0.15, ACCENT1)

items = ["단일 입력  →  단일 출력", "기억 없음", "역할 없음", "자동화 없음"]
for i, item in enumerate(items):
    b = box(s, 0.8, 1.5+i*1.2, 5.5, 1.0, bg_color=BG_CARD, border=ACCENT1)
    txt(s, f"✗  {item}", 0.95, 1.55+i*1.2, 5.2, 0.6, size=18, color=TEXT_LIGHT)

txt(s, "예시", 7.5, 1.4, 5.0, 0.4, size=16, bold=True, color=TEXT_GRAY)
examples = ['"코드 짜줘"', '"PPT 만들어줘"', '"이거 번역해줘"', '"오류 찾아줘"']
for i, ex in enumerate(examples):
    b = box(s, 7.5, 1.9+i*1.0, 4.8, 0.8, bg_color=RGBColor(0x1A,0x24,0x45))
    txt(s, ex, 7.65, 1.95+i*1.0, 4.5, 0.5, size=16, color=TEXT_WHITE)


# ═══════════════════════════════════════════════
# 슬라이드 4 — 2단계 Agent
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "2단계 — Agent (에이전트)", "역할을 부여하는 단계")
badge(s, "2단계", 0.5, 0.15, ACCENT1)

items = ["역할 기반 동작", "Context 유지", "특정 목적 수행", "일관된 페르소나"]
for i, item in enumerate(items):
    b = box(s, 0.8, 1.5+i*1.2, 5.5, 1.0, bg_color=BG_CARD, border=ACCENT1)
    txt(s, f"✓  {item}", 0.95, 1.55+i*1.2, 5.2, 0.6, size=18, color=ACCENT1)

txt(s, "예시", 7.5, 1.4, 5.0, 0.4, size=16, bold=True, color=TEXT_GRAY)
examples = ['"너는 백엔드 개발자야"', '"너는 PPT 전문가야"', '"너는 데이터 분석가야"']
for i, ex in enumerate(examples):
    b = box(s, 7.5, 1.9+i*1.2, 4.8, 1.0, bg_color=RGBColor(0x1A,0x24,0x45))
    txt(s, ex, 7.65, 1.98+i*1.2, 4.5, 0.6, size=16, color=TEXT_WHITE)


# ═══════════════════════════════════════════════
# 슬라이드 5 — 3단계 Orchestration
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "3단계 — Orchestration (오케스트레이션)", "여러 AI를 연결하는 단계")
badge(s, "3단계", 0.5, 0.15, ACCENT1)

agents = [
    ("Claude", "설계·판단·승인", ACCENT1, 1.0),
    ("Codex",  "코드 구현 (병렬)", RGBColor(0x00,0xC4,0x80), 4.8),
    ("Gemini", "검증·리뷰·문서화", RGBColor(0xFF,0xA5,0x00), 8.6),
]
for name, role, color, left in agents:
    b = box(s, left, 1.5, 3.5, 2.2, bg_color=BG_CARD, border=color)
    txt(s, name, left+0.1, 1.6, 3.3, 0.7, size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, role, left+0.1, 2.3, 3.3, 0.6, size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "→",  4.55, 2.2, 0.7, 0.5, size=24, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
txt(s, "→",  8.35, 2.2, 0.7, 0.5, size=24, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "A2A (Agent to Agent) 핸드오프 — 결과물을 다음 AI에게 전달", 0.5, 4.1, 12.33, 0.5,
    size=16, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# 핸드오프 설명
handoffs = [
    "Claude  →  task-instruction.md  →  Codex",
    "Codex   →  implementation-report.md  →  Claude",
    "Claude  →  verify-request.md  →  Gemini",
    "Gemini  →  review-result.md  →  Claude (채택 결정)",
]
for i, h in enumerate(handoffs):
    txt(s, h, 0.8, 4.8+i*0.45, 11.73, 0.4, size=13, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════
# 슬라이드 6 — 3단계 파이프라인
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "3단계 파이프라인", "Hook → Planner → Executor → Validator → State → Retry")

pipeline = [
    ("Hook",      "사전 확인\n파일 잠금·충돌 체크",     ACCENT1,  0.3),
    ("Planner",   "Claude 설계\ntask-instruction.md",  ACCENT1,  2.35),
    ("Executor",  "Codex 4개\n병렬 구현",               RGBColor(0x00,0xC4,0x80), 4.4),
    ("Validator", "Gemini 2개\n검증·보안·품질",         RGBColor(0xFF,0xA5,0x00), 6.45),
    ("State",     "스냅샷 저장\n세션 복구",             RGBColor(0xA0,0x78,0xFF), 8.5),
    ("Retry",     "3회 재시도\n→ 에스컬레이션",         ACCENT2, 10.55),
]
for name, desc, color, left in pipeline:
    b = box(s, left, 1.5, 2.6, 3.5, bg_color=BG_CARD, border=color)
    txt(s, name, left+0.05, 1.6, 2.5, 0.7, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, left+0.1, 2.4, 2.4, 1.5, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

for i in range(5):
    txt(s, "→", 2.9+i*2.05, 2.9, 0.5, 0.5, size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "실패 시: Retry → 3회 초과 → Claude 에스컬레이션", 0.5, 5.5, 12.33, 0.4,
    size=14, color=ACCENT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 7 — 4단계 Automation
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "4단계 — Automation (자동화)", "사람이 없어도 돌아가는 단계")
badge(s, "4단계", 0.5, 0.15, ACCENT1)

cols = [
    ("Plugin\n13개", "exec_orch / exec_voice\nexec_learning\ndesign_* / review_qa\nmcp_* × 7", 0.5),
    ("Command\n21개", "vibe-loop / loop-stop\nplug_* × 8\ncheck-agents\ngodmode / 10x", 3.5),
    ("Hook\n9개", "pre-task\npost-impl\npost-review\nnotify / ai-handoff", 6.5),
    ("MCP\n연동", "context7 / playwright\nfigma / canva\ngithub / postgres\nmongodb 외", 9.5),
]
for title, desc, left in cols:
    b = box(s, left, 1.4, 2.8, 4.5, bg_color=BG_CARD, border=ACCENT1)
    txt(s, title, left+0.1, 1.5, 2.6, 0.9, size=18, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    txt(s, desc,  left+0.1, 2.5, 2.6, 3.0, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

txt(s, "트리거 기반 자동 실행 — 이벤트가 발생하면 스스로 동작", 0.5, 6.3, 12.33, 0.4,
    size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 8 — 🔥 5단계 Autonomous System
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)

# 강조 배경
glow = s.shapes.add_shape(9, Inches(4), Inches(1.5), Inches(5.33), Inches(4.5))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0xFF,0x4A,0x00)
glow.line.fill.background()
# 겹쳐서 대부분 가리기
cover = s.shapes.add_shape(1, Inches(0), Inches(0), W, H)
cover.fill.solid(); cover.fill.fore_color.rgb = BG_DARK; cover.line.fill.background()

heading(s, "🔥 5단계 — Autonomous System", "여기서 차원이 바뀐다", ACCENT2)

txt(s, "스스로 판단  ·  스스로 선택  ·  스스로 복구",
    0.5, 1.1, 12.33, 0.6, size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

elements = [
    ("라우팅", "누가 할지 AI가 결정\nroute_dispatch", ACCENT2, 0.5, 1.8),
    ("실패 판단", "성공/실패 스스로 인식\nretry-count.json", ACCENT2, 3.5, 1.8),
    ("재시도", "자동 3회\n→ 에스컬레이션", ACCENT2, 6.5, 1.8),
    ("상태 저장", "기억하고 학습\nexec_learning", ACCENT2, 9.5, 1.8),
]
for name, desc, color, l, t in elements:
    b = box(s, l, t, 2.8, 2.8, bg_color=RGBColor(0x2A,0x10,0x05), border=color)
    txt(s, name, l+0.1, t+0.1, 2.6, 0.7, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, l+0.1, t+0.9, 2.6, 1.6, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

txt(s, "✅ route_dispatch  ✅ retry-count.json  ✅ state_session  ✅ exec_learning",
    0.5, 4.9, 12.33, 0.5, size=15, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "→  5단계 핵심 요소 구현 완료",
    0.5, 5.5, 12.33, 0.4, size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 9 — 🔥 6단계 Platform
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "🔥 6단계 — Platform (플랫폼)", "시스템이 아니라 생태계", ACCENT3)

txt(s, "여러 Workflow  ·  여러 사용자  ·  API · UI  ·  서비스화",
    0.5, 1.1, 12.33, 0.6, size=20, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

items = [
    ("다중 워크플로우",  "여러 프로젝트 동시 운영"),
    ("다중 사용자",      "사용자별 권한 관리"),
    ("비용·로그 관리",   "토큰 추적 + 사용량 분석"),
    ("REST API 제공",    "외부 서비스 연동"),
    ("대시보드 UI",      "dashboard.py 확장"),
    ("서비스화",         "SaaS / 내부 AI 플랫폼"),
]
for i, (title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    l = 0.5 + col * 6.5
    t = 2.0 + row * 1.5
    b = box(s, l, t, 6.0, 1.2, bg_color=RGBColor(0x03,0x1A,0x2A), border=ACCENT3)
    txt(s, title, l+0.15, t+0.1, 3.0, 0.5, size=15, bold=True, color=ACCENT3)
    txt(s, desc,  l+0.15, t+0.6, 5.7, 0.5, size=13, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════
# 슬라이드 10 — 전체 비교표
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "1→6단계 전체 비교", "단계별 능력 매트릭스")

headers = ["단계", "핵심 키워드", "자동화", "자율성", "확장성"]
rows = [
    ["1️⃣ Prompt",        "질문",     "❌", "❌", "❌"],
    ["2️⃣ Agent",         "역할",     "△",  "❌", "△"],
    ["3️⃣ Orchestration", "AI 연결",  "△",  "❌", "○"],
    ["4️⃣ Automation",    "자동 실행","✅", "△",  "○"],
    ["🔥 5 Autonomous",  "자율 판단","✅", "✅", "○"],
    ["🔥 6 Platform",    "생태계",   "✅", "✅", "✅"],
]
col_w = [3.2, 2.5, 1.5, 1.5, 1.5]
col_x = [0.3, 3.55, 6.1, 7.65, 9.2]
row_h = 0.75

# 헤더
for j, (header, x, w) in enumerate(zip(headers, col_x, col_w)):
    b = box(s, x, 1.3, w, 0.55, bg_color=ACCENT1)
    txt(s, header, x+0.05, 1.33, w-0.1, 0.45,
        size=14, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# 데이터
for i, row in enumerate(rows):
    y = 1.95 + i * row_h
    highlight = i >= 4  # 5, 6단계 강조
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        bg_c = RGBColor(0x2A,0x10,0x05) if highlight else BG_CARD
        bc   = ACCENT2 if highlight else None
        b = box(s, x, y, w, row_h-0.05, bg_color=bg_c, border=bc)
        color = GOLD if highlight else TEXT_LIGHT
        txt(s, cell, x+0.05, y+0.1, w-0.1, row_h-0.2,
            size=13, color=color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 슬라이드 11 — 현재 위치와 로드맵
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)
heading(s, "현재 위치와 로드맵", "우리는 어디에 있고, 어디로 가는가")

# 진행 바
stages_bar = ["Prompt","Agent","Orch.","Automation","Autonomous★","Platform"]
colors_bar  = [ACCENT1,ACCENT1,ACCENT1,ACCENT1,GOLD,TEXT_GRAY]
done        = [True, True, True, True, True, False]
for i, (name, color, is_done) in enumerate(zip(stages_bar, colors_bar, done)):
    x = 0.5 + i * 2.1
    b = box(s, x, 1.4, 2.0, 1.2,
            bg_color=(color if is_done else BG_CARD),
            border=color)
    txt(s, name, x+0.05, 1.5, 1.9, 0.7,
        size=13, bold=is_done,
        color=(TEXT_WHITE if is_done else TEXT_GRAY),
        align=PP_ALIGN.CENTER)
    if is_done and i < 5:
        txt(s, "✅", x+0.6, 2.05, 0.8, 0.4, size=14, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

txt(s, "← 완성", 0.5, 2.8, 10.5, 0.4, size=14, color=GOLD)
txt(s, "목표 →", 10.6, 2.8, 2.2, 0.4, size=14, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)

# 로드맵
roadmap = [
    ("완료", "1~4단계", "기반 구축 완료",             ACCENT1),
    ("완료", "5단계",   "자율 시스템 핵심 요소 구현", GOLD),
    ("진행", "5단계 고도화", "자율성 강화·학습 루프 실전 적용", ACCENT2),
    ("목표", "6단계",   "Platform — 대시보드·다중 사용자·API", ACCENT3),
]
for i, (status, stage, desc, color) in enumerate(roadmap):
    y = 3.5 + i * 0.85
    badge(s, status, 0.5, y+0.05, color)
    txt(s, stage, 1.8, y, 2.5, 0.7, size=15, bold=True, color=color)
    txt(s, desc,  4.3, y, 8.5, 0.7, size=13, color=TEXT_LIGHT)


# ═══════════════════════════════════════════════
# 슬라이드 12 — 마무리
# ═══════════════════════════════════════════════
s = add_slide(); bg(s)

line_top = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.06))
line_top.fill.solid(); line_top.fill.fore_color.rgb = ACCENT2; line_top.line.fill.background()

txt(s, '"AI는 도구가 아니라 팀이다"',
    0.5, 0.5, 12.33, 1.0,
    size=32, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

summary = [
    ("1단계", "질문한다",          ACCENT1),
    ("2단계", "역할을 준다",        ACCENT1),
    ("3단계", "연결한다",           ACCENT1),
    ("4단계", "자동화한다",         ACCENT1),
    ("5단계", "스스로 판단한다  ←  지금 여기", GOLD),
    ("6단계", "생태계가 된다  ←  다음 목표",   ACCENT3),
]
for i, (stage, line, color) in enumerate(summary):
    y = 1.8 + i * 0.75
    txt(s, stage, 1.5, y, 1.8, 0.6, size=16, bold=True, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    txt(s, "  →  " + line, 3.4, y, 9.0, 0.6, size=18, bold=(i>=4), color=color)

line_bot = s.shapes.add_shape(1, Inches(1), Inches(6.8), Inches(11.33), Inches(0.03))
line_bot.fill.solid(); line_bot.fill.fore_color.rgb = ACCENT1; line_bot.line.fill.background()
txt(s, "Orchestration Kit v1.0  |  github.com/bernakilljos/orchestration",
    0.5, 6.9, 12.33, 0.4, size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════
import os
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "ai-system-stages-2026-04-16.pptx")
prs.save(out_path)
print(f"저장 완료: {out_path}")
print(f"슬라이드 수: {len(prs.slides)}장")
