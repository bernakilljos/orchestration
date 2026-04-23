"""
3종 PPT 생성 — 2026-04-23
  PPT-1: Claude Code 한글 설명 (마인드맵 스타일 포함)
  PPT-2: 사용자 작업 단계 (워크플로우)
  PPT-3: plugins 사용법

요구: pip install python-pptx
실행: python outputs/ppt/generate-3-decks.py
출력: outputs/ppt/{1-claude-code, 2-user-steps, 3-plugins}-2026-04-23.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from datetime import date

OUT_DIR = Path(__file__).parent
TODAY = date.today().isoformat()

# ── 컬러 팔레트 ──────────────────────────
APRICOT  = RGBColor(0xE8, 0x9B, 0x7A)
PEACH    = RGBColor(0xF4, 0xB8, 0x91)
BEIGE    = RGBColor(0xE8, 0xDC, 0xC4)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
CREAM    = RGBColor(0xFA, 0xF6, 0xF1)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x99, 0x99, 0x99)
INDIGO   = RGBColor(0x63, 0x66, 0xF1)
PINK     = RGBColor(0xEC, 0x48, 0x99)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
RED      = RGBColor(0xEF, 0x44, 0x44)

# 레이어 케이크 색
LAYER_COLORS = [
    RGBColor(0xF4, 0xD5, 0xC2),  # apricot
    RGBColor(0xD5, 0xC7, 0xE8),  # lilac
    RGBColor(0xC2, 0xD5, 0xE8),  # blue
    RGBColor(0xC7, 0xE8, 0xC2),  # green
    RGBColor(0xE8, 0xDC, 0xC2),  # yellow
    RGBColor(0xE8, 0xC2, 0xC2),  # coral
    RGBColor(0xE0, 0xE0, 0xE0),  # gray
]

FONT_KO = "맑은 고딕"
FONT_MONO = "Consolas"

# ── 헬퍼 ────────────────────────────────

def new_deck(width_in=13.333, height_in=7.5):
    """16:9 와이드 스크린"""
    p = Presentation()
    p.slide_width = Inches(width_in)
    p.slide_height = Inches(height_in)
    return p

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    return box

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KO):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_line(slide, x1, y1, x2, y2, color=GRAY, weight=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def cover_slide(prs, title, subtitle, footer):
    s = blank_slide(prs)
    # 배경
    bg = add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    # 중앙 살구색 원
    circle = add_rect(s, Inches(5.7), Inches(1.0), Inches(2.0), Inches(2.0),
                      APRICOT, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(5.7), Inches(1.0), Inches(2.0), Inches(2.0),
             "✦", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.5),
             title, size=44, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    # 부제목
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
             subtitle, size=20, color=APRICOT, align=PP_ALIGN.CENTER)
    # 푸터
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             footer, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    return s

def section_title_slide(prs, num, title, lead):
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    # 큰 번호
    add_text(s, Inches(0.5), Inches(2.5), Inches(3.0), Inches(2.5),
             num, size=200, bold=True, color=APRICOT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    add_text(s, Inches(4.0), Inches(2.8), Inches(8.5), Inches(1.0),
             title, size=40, bold=True, color=CHARCOAL)
    # 리드
    add_text(s, Inches(4.0), Inches(4.0), Inches(8.5), Inches(2.0),
             lead, size=18, color=CHARCOAL)
    return s

def content_slide(prs, title, body_lines, *, title_color=CHARCOAL):
    """일반 본문 — 제목 + bullet 라인들"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    # 제목
    add_text(s, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8),
             title, size=28, bold=True, color=title_color)
    # 구분선
    line = add_rect(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.04),
                    APRICOT, shape=MSO_SHAPE.RECTANGLE)
    # 본문
    if isinstance(body_lines, str):
        body_lines = body_lines.split("\n")
    add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5),
             "\n".join(body_lines), size=16, color=CHARCOAL)
    return s

def mindmap_slide(prs, title, center, categories):
    """방사형 마인드맵 — 중앙 + 카테고리 4~6 + 노드"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    # 제목
    add_text(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             title, size=24, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

    cx, cy = Inches(6.67), Inches(4.0)
    # 중앙 원
    cr = Inches(0.7)
    add_rect(s, cx - cr, cy - cr, cr*2, cr*2, APRICOT, shape=MSO_SHAPE.OVAL)
    add_text(s, cx - cr, cy - cr, cr*2, cr*2, center,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 4개 카테고리: 상하좌우 배치
    n = len(categories)
    if n == 4:
        positions = [
            (Inches(5.5), Inches(1.3)),   # top
            (Inches(0.5), Inches(3.7)),   # left
            (Inches(10.5), Inches(3.7)),  # right
            (Inches(5.5), Inches(6.0)),   # bottom
        ]
    else:
        # 6개 — 정육각형
        import math
        positions = []
        for i in range(n):
            angle = 2 * math.pi * i / n - math.pi/2
            px = cx + Emu(int(Inches(3.8).emu * math.cos(angle))) - Inches(1.2)
            py = cy + Emu(int(Inches(2.5).emu * math.sin(angle))) - Inches(0.5)
            positions.append((px, py))

    for (cat_name, nodes), (cx2, cy2) in zip(categories, positions):
        # 카테고리 박스 (검정)
        cw, ch = Inches(2.4), Inches(0.5)
        add_rect(s, cx2, cy2, cw, ch, CHARCOAL)
        add_text(s, cx2, cy2, cw, ch, cat_name,
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 연결선 (중심 → 카테고리)
        add_line(s, cx, cy, cx2 + cw/2, cy2 + ch/2, color=GRAY, weight=0.75)
        # 노드 (카테고리 아래)
        for i, node in enumerate(nodes[:3]):
            ny = cy2 + Inches(0.6) + Inches(0.55) * i
            add_rect(s, cx2, ny, cw, Inches(0.45), PEACH if i % 2 == 0 else BEIGE)
            add_text(s, cx2 + Inches(0.05), ny, cw - Inches(0.1), Inches(0.45),
                     node, size=10, color=CHARCOAL, anchor=MSO_ANCHOR.MIDDLE)
    return s

def layered_slide(prs, title, layers):
    """레이어 케이크 — 위→아래"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             title, size=24, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    n = len(layers)
    avail_h = Inches(6.2)
    layer_h = Emu(int(avail_h.emu / n))
    y = Inches(1.1)
    for i, (label, items) in enumerate(layers):
        color = LAYER_COLORS[i % len(LAYER_COLORS)]
        add_rect(s, Inches(1.5), y, Inches(10.5), layer_h, color, shape=MSO_SHAPE.RECTANGLE)
        # 좌측 아이콘 박스
        add_rect(s, Inches(0.5), y, Inches(0.9), layer_h,
                 RGBColor(int(color[0]*0.7), int(color[1]*0.7), int(color[2]*0.7)),
                 shape=MSO_SHAPE.RECTANGLE)
        # 라벨
        add_text(s, Inches(1.7), y, Inches(3.5), layer_h, label,
                 size=14, bold=True, color=CHARCOAL, anchor=MSO_ANCHOR.MIDDLE)
        # 항목
        items_text = "  •  ".join(items)
        add_text(s, Inches(5.3), y, Inches(6.5), layer_h,
                 items_text, size=11, color=CHARCOAL, anchor=MSO_ANCHOR.MIDDLE)
        y = y + layer_h
    return s

def steps_slide(prs, title, steps):
    """단계별 화살표 흐름"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             title, size=24, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    n = len(steps)
    if n <= 4:
        cols = n
        rows = 1
    else:
        cols = (n + 1) // 2
        rows = 2
    cell_w = Inches(11.5 / cols)
    cell_h = Inches(2.5)
    start_x = Inches(0.9)
    start_y = Inches(1.5)
    for i, (num, head, body) in enumerate(steps):
        r = i // cols
        c = i % cols
        x = start_x + cell_w * c
        y = start_y + cell_h * r + Inches(0.5) * r
        # 박스
        add_rect(s, x, y, cell_w - Inches(0.2), cell_h - Inches(0.2), WHITE,
                 line=APRICOT)
        # 번호 원
        nr = Inches(0.45)
        add_rect(s, x + Inches(0.15), y + Inches(0.15), nr, nr, APRICOT, shape=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.15), y + Inches(0.15), nr, nr, str(num),
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 헤드
        add_text(s, x + Inches(0.7), y + Inches(0.15), cell_w - Inches(0.9), Inches(0.5),
                 head, size=14, bold=True, color=CHARCOAL)
        # 본문
        add_text(s, x + Inches(0.2), y + Inches(0.7), cell_w - Inches(0.4), cell_h - Inches(0.9),
                 body, size=11, color=CHARCOAL)
    return s

def two_column_slide(prs, title, left_title, left_lines, right_title, right_lines):
    s = blank_slide(prs)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, CREAM, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
             title, size=26, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)
    # left
    add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.7), WHITE, line=APRICOT)
    add_text(s, Inches(0.7), Inches(1.45), Inches(5.7), Inches(0.5),
             left_title, size=18, bold=True, color=APRICOT)
    add_text(s, Inches(0.7), Inches(2.05), Inches(5.7), Inches(4.9),
             "\n".join(left_lines), size=13, color=CHARCOAL)
    # right
    add_rect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.7), WHITE, line=INDIGO)
    add_text(s, Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.5),
             right_title, size=18, bold=True, color=INDIGO)
    add_text(s, Inches(7.0), Inches(2.05), Inches(5.7), Inches(4.9),
             "\n".join(right_lines), size=13, color=CHARCOAL)
    return s

# ════════════════════════════════════════
#  PPT-1 : Claude Code 한글 설명
# ════════════════════════════════════════

def build_ppt1():
    prs = new_deck()
    cover_slide(prs,
        "Claude Code 완벽 가이드",
        "한 장의 그림으로 보는 Claude Code 의 모든 것",
        f"orchestration_v1 · {TODAY}")

    # 목차
    content_slide(prs, "목차",
        ["1.  Claude Code 란 무엇인가",
         "2.  6 레이어 아키텍처 — 케이크처럼 쌓인 시스템",
         "3.  4 가지 Work Modes — 사용자가 만나는 4 개 얼굴",
         "4.  Claude 모델 라인업 — Opus / Sonnet / Haiku",
         "5.  Context System — AI 의 기억 장치",
         "6.  Integrations & Connectors — 외부 세계와 연결",
         "7.  Create — 결과물을 만드는 3 가지 방식",
         "8.  For Teams — 팀 사용 옵션",
         "9.  마인드맵 — 한 장에 다 보기"])

    section_title_slide(prs, "01", "Claude Code 란?",
        "터미널에서 동작하는 AI 코딩 에이전트.\n"
        "단순 챗봇이 아니라 파일·MCP·하위 에이전트를 다루는 운영체제 레이어.")

    content_slide(prs, "1. Claude Code 의 핵심 정체성",
        ["•  Anthropic 의 공식 CLI — npm i -g @anthropic-ai/claude-code",
         "•  대화 ≠ 끝 — 파일 읽기·쓰기·편집·실행 모두 직접",
         "•  ~/.claude/, ./CLAUDE.md, .claude/ 세 영역으로 메모리 분리",
         "•  Skills, Hooks, MCP, Sub-agents 가 1급 시민 (first-class)",
         "•  가격: Pro $20 / Max5x $100 / Max20x $200 / API pay-per-use",
         "",
         "한 줄 요약: \"AI 동료가 당신 컴퓨터에 들어와 함께 일한다.\""])

    section_title_slide(prs, "02", "6 레이어 아키텍처",
        "위에서 아래로 — 사용자에 가까울수록 위쪽,\n인프라에 가까울수록 아래쪽.")

    layered_slide(prs, "2. Claude Code 의 7 층 케이크",
        [("TOP — Workflow Patterns",   ["Architect", "Creator", "PM"]),
         ("LAYER 5 — Agent Orchestration", ["Multi-Agent", "Sub-agents", "Hooks"]),
         ("LAYER 4 — Commands & Shortcuts", ["/help", "/compact", "@filename", "Esc·Esc"]),
         ("LAYER 3 — MCP Connections", ["USB-C for AI", "200+ tools", "GitHub·Slack·Notion"]),
         ("LAYER 2 — Skills Engine",   ["~/.claude/skills/", "auto-activate", "SKILL.md"]),
         ("LAYER 1 — Memory System",   ["CLAUDE.md", ".claudeignore", "Global·Project·Folder"]),
         ("FOUNDATION — The Runtime",  ["File access", "Hours-long tasks", "Full terminal"])])

    section_title_slide(prs, "03", "4 가지 Work Modes",
        "Claude 를 만나는 4 개의 입구 — 용도에 따라 다른 입구를 선택.")

    steps_slide(prs, "3. Work Modes — 4 개 얼굴",
        [(1, "Claude Chat",      "기본 대화창.\n시작하기 가장 쉬움."),
         (2, "Claude Cowork",    "내 PC 에 산다.\n파일 읽고 문서 만듦."),
         (3, "Claude Code",      "개발자의 무기.\n터미널에서 자율 실행."),
         (4, "Claude Projects",  "지속적 폴더.\n팀에 적합.")])

    section_title_slide(prs, "04", "모델 라인업",
        "한 종류 AI 가 아니라 4 종 — 작업 성격에 맞춰 사용.")

    steps_slide(prs, "4. Claude 모델 — 4 종",
        [(1, "Opus 4.7",          "최신·최강.\n생각이 필요한 모든 것에."),
         (2, "Sonnet 4.6",        "빠르고 날렵함.\n속도가 깊이보다 중요할 때."),
         (3, "Haiku 4.5",         "초고속.\n가볍고 저렴, 작은 작업에."),
         (4, "Extended Thinking", "답변 전 내부 사고.\n복잡 추론 시 켜기.")])

    section_title_slide(prs, "05", "Context System",
        "AI 의 \"기억\" 을 관리하는 5 가지 도구 —\n환각(hallucination) 을 막고 일관성 유지.")

    content_slide(prs, "5. Context System — 5 가지 도구",
        ["1.  CLAUDE.md         프로젝트의 헌법. AI 가 항상 읽음.",
         "                       Global / Project / Folder 3 단계 (Last wins)",
         "",
         "2.  Markdown Files    내 음성·규칙·예시를 .md 로.",
         "                       프롬프트 라이브러리 대체.",
         "",
         "3.  AskUserQuestion   AI 가 사용자에게 폼으로 묻기 (선택지·라벨).",
         "                       자유텍스트 답변보다 정확.",
         "",
         "4.  1M Token Window   책 한 권을 한 번에 읽음.",
         "                       대화 중 환각 격감.",
         "",
         "5.  Global Instructions  영구 규칙. 한 번 설정 → 영원히 적용."])

    section_title_slide(prs, "06", "Integrations & Connectors",
        "Claude 가 외부 세계와 만나는 4 가지 방식.")

    content_slide(prs, "6. Integrations — 외부 연결 4 종",
        ["•  Claude in Excel    셀 D14 가 실제 무엇인지 안다.",
         "                        스프레드시트 안의 AI.",
         "",
         "•  Claude in Chrome   웹 작업 자율 완료.",
         "                        브라우저 에이전트.",
         "",
         "•  Connectors         Slack·Drive·Notion·50+ 도구 1-click 연결.",
         "                        실제 앱을 읽음.",
         "",
         "•  Plugins            사전 제작 스킬팩 (Sales·Legal·Marketing·Data).",
         "                        한 번 클릭으로 설치."])

    section_title_slide(prs, "07", "Create",
        "결과물을 만드는 3 가지 방식 —\n단순 텍스트가 아닌 진짜 산출물.")

    steps_slide(prs, "7. Create — 결과물 3 종",
        [(1, "Artifacts",         "Claude 안의 인터랙티브 출력.\n계산기·차트·대시보드."),
         (2, "Skills",            "반복 가능한 워크플로우.\nAI 직원의 SOP."),
         (3, "Prompt Templates",  "팀 전체가 공유하는 시작 프롬프트.\n한 번 설정 → 영원히 사용.")])

    section_title_slide(prs, "08", "For Teams",
        "팀 사용 — 3 단계 옵션.")

    steps_slide(prs, "8. For Teams — 3 단계",
        [(1, "Team Plan",       "5~150 시트.\n공유 프로젝트, 데이터 학습 안함."),
         (2, "Enterprise",      "SSO·관리자·보안 컴플라이언스.\n법무팀 질문 시."),
         (3, "Shared Projects", "한 프로젝트, 전체 팀.\n같은 컨텍스트·같은 품질.")])

    section_title_slide(prs, "09", "마인드맵",
        "지금까지의 모든 것을 한 장으로.")

    mindmap_slide(prs, "9. Claude is eating up everything — 한 장 요약",
        center="Claude",
        categories=[
            ("Work Modes",     ["Chat", "Cowork", "Code"]),
            ("Models",         ["Opus 4.7", "Sonnet 4.6", "Haiku 4.5"]),
            ("Context",        ["CLAUDE.md", "1M Token", "AskUser"]),
            ("Integrations",   ["Excel", "Chrome", "Connectors"]),
            ("Create",         ["Artifacts", "Skills", "Prompts"]),
            ("For Teams",      ["Team Plan", "Enterprise", "Shared"]),
        ])

    # 마무리
    cover_slide(prs,
        "감사합니다",
        "다음: 사용자 작업 단계 / 플러그인 사용법",
        "참조: docs/screens/arch/  ·  guide.txt  ·  CLAUDE.md")

    out = OUT_DIR / f"1-claude-code-{TODAY}.pptx"
    prs.save(str(out))
    return out

# ════════════════════════════════════════
#  PPT-2 : 사용자 작업 단계
# ════════════════════════════════════════

def build_ppt2():
    prs = new_deck()
    cover_slide(prs,
        "사용자 작업 단계",
        "이 킷을 어떻게 쓰는가 — 7 단계 워크플로우",
        f"orchestration_v1 · {TODAY}")

    content_slide(prs, "목차",
        ["1.  큰 그림 — 어떤 일을 자동화하나",
         "2.  Step 1: 설치 (3 가지 방법)",
         "3.  Step 2: CLAUDE.md 확인·커스터마이즈",
         "4.  Step 3: Claude Code 실행 + 첫 작업",
         "5.  Step 4: 작업 의뢰 — 자연어로 말하면 자동 라우팅",
         "6.  Step 5: 워커 모니터링",
         "7.  Step 6: 결과물 확인",
         "8.  Step 7: git 커밋·푸시",
         "9.  트러블슈팅 + 자주 쓰는 슬래시 커맨드"])

    section_title_slide(prs, "00", "큰 그림",
        "사용자 = 팀장. AI 들 = 팀원.\n당신은 \"무엇을\" 만 말하고 \"어떻게\" 는 AI 가 처리.")

    two_column_slide(prs, "1. 큰 그림 — 누가 무엇을",
        "사용자가 하는 것",
        ["•  무엇을 만들지 결정",
         "•  방향 승인",
         "•  최종 검토",
         "•  결과 확인",
         "",
         "(설계·라우팅·실행은 AI 가)"],
        "킷이 자동으로 하는 것",
        ["•  Codex/Gemini/Claude 라우팅",
         "•  task-instruction.md 작성",
         "•  병렬 워커 실행",
         "•  핸드오프 (handoff-log)",
         "•  검증 → 채택",
         "•  outputs/ 저장"])

    section_title_slide(prs, "01", "설치",
        "3 가지 방법 — 본인 환경에 맞게 선택.")

    steps_slide(prs, "2. Step 1 — 설치",
        [(1, "GUI 위자드 (추천)",
            "OrchestrationKit-Setup.exe\n경로 선택 → 컴포넌트 → 설치"),
         (2, "git clone (CLI)",
            "git clone <repo>\nsetup\\setup.bat C:\\work\\myproject"),
         (3, "사일런트",
            "Setup.exe /VERYSILENT\n/DIR=\"C:\\work\\myproject\"\n자동화 스크립트용")])

    section_title_slide(prs, "02", "CLAUDE.md 확인",
        "프로젝트의 헌법. AI 가 항상 읽고 따름.\n자신의 규칙·금지 사항·경로를 명시.")

    content_slide(prs, "3. Step 2 — CLAUDE.md 커스터마이즈",
        ["기본 제공된 CLAUDE.md 가 이미 좋은 출발점.",
         "추가/수정할 만한 항목:",
         "",
         "  ·  Tech Stack       내 프로젝트의 언어·프레임워크",
         "  ·  Commands         build / test / deploy 명령",
         "  ·  Style            내 코딩 컨벤션",
         "  ·  금지 사항         optional chaining 금지 등",
         "  ·  Loading Order    스킬 로드 순서",
         "",
         "원칙: 500 줄 이하 · WHAT/WHY/HOW · 참조 중심 (중복 금지)",
         "",
         "관련: /init  /help"])

    section_title_slide(prs, "03", "Claude Code 실행",
        "터미널에서 'claude' 입력 → 첫 대화 시작.")

    content_slide(prs, "4. Step 3 — Claude Code 실행 + 첫 작업",
        ["[프로젝트 폴더에서]",
         "  $ claude",
         "",
         "[자동 발생하는 것]",
         "  ·  Session Start hook 실행 (.claude/hooks/hook-00-init)",
         "  ·  Orca Auto 활성화 (워커 spawn 준비)",
         "  ·  CLAUDE.md 로드 (프로젝트 컨텍스트)",
         "  ·  플러그인 18 개 자동 로드",
         "",
         "[첫 메시지 예시]",
         "  \"한 페이지 랜딩 만들어줘. Next.js 14 + Tailwind, 다크모드 지원\"",
         "",
         "→ Claude 가 task-instruction.md 작성 → Codex 워커 위임"])

    section_title_slide(prs, "04", "작업 의뢰",
        "자연어로 말하면 적절한 AI 에 자동 라우팅 —\n사용자는 \"누가 할까\" 신경 쓸 필요 없음.")

    two_column_slide(prs, "5. Step 4 — 자연어로 말하기",
        "Claude 가 직접",
        ["•  설계·아키텍처 결정",
         "•  PPT·디자인 (Gamma·Canva)",
         "•  500 줄 미만 코드 보완",
         "•  복잡 추론·트레이드오프",
         "",
         "예: \"이 구조 비교해줘\""],
        "Codex/Gemini 워커로 위임",
        ["•  500 줄 이상 신규 구현 → Codex (×4 병렬)",
         "•  코드 리뷰·검증 → Gemini Flash (저단가)",
         "•  대량 문서 요약 → Gemini",
         "",
         "예: \"전체 모듈 구현해줘\"",
         "→ task-instruction.md → codex-auto"])

    section_title_slide(prs, "05", "모니터링",
        "워커가 백그라운드에서 일하는 동안 진행 상황 확인.")

    steps_slide(prs, "6. Step 5 — 워커 모니터링",
        [(1, "/exec_status",
            "통합 대시보드.\n워커·큐·heartbeat 한 번에."),
         (2, "/check-agents",
            "워커 가용 여부.\n실행 중 작업 현황."),
         (3, "ls .claude/tasks/",
            "현재 큐 (대기·진행·완료).\ndone/ 폴더 = 완료."),
         (4, "ls .claude/state/",
            "heartbeat·quota·token-log.\n워커 살아있나?")])

    section_title_slide(prs, "06", "결과 확인",
        "워커 작업 결과는 자동으로 정해진 위치에 저장.")

    content_slide(prs, "7. Step 6 — 결과물 확인",
        ["[코드 결과]",
         "  src/ 또는 지정 폴더에 직접 생성",
         "",
         "[문서·디자인 결과]",
         "  outputs/ppt/<name>-<date>.pptx",
         "  outputs/pdf/<name>-<date>.pdf",
         "  outputs/arch/<pattern>-<slug>-<date>.png",
         "  outputs/artifacts/<type>-<slug>-<date>.html",
         "",
         "[태스크 로그]",
         "  .claude/tasks/done/task-*.md  (완료된 작업 기록)",
         "  .claude/state/token-log.jsonl  (토큰 소비 이력)",
         "",
         "[Gemini 검증 보고서]",
         "  docs/<date>/review-*.md"])

    section_title_slide(prs, "07", "커밋·푸시",
        "확인된 결과를 git 에 반영.")

    content_slide(prs, "8. Step 7 — git 커밋·푸시",
        ["[수동]",
         "  git add <files>",
         "  git commit -m \"feat: ...\"",
         "  git push",
         "",
         "[Claude 에게 맡기기 — 자연어]",
         "  \"방금 작업 커밋해줘\"  →  /commit  자동 호출",
         "  \"PR 만들어줘\"          →  /commit-push-pr  자동 호출",
         "",
         "[보안]",
         "  ·  .env / 시크릿 자동 제외 (gitignore)",
         "  ·  hook-08 가 위장 완료 차단 (빈 task done 이동 금지)"])

    section_title_slide(prs, "09", "트러블슈팅",
        "자주 발생하는 상황 + 자주 쓰는 커맨드.")

    two_column_slide(prs, "9. 자주 쓰는 / 트러블슈팅",
        "자주 쓰는 슬래시 커맨드",
        ["/help              플러그인 도움말",
         "/exec_status       통합 상태",
         "/check-agents      워커 가용",
         "/orcauto-start     자동 시작",
         "/orcauto-stop      자동 종료",
         "/loop-stop         루프 중단",
         "/godmode           최대 출력 모드",
         "/commit            커밋",
         "/arch-auto <주제>  다이어그램 자동",
         "/claude-status     Claude 가용성"],
        "트러블슈팅",
        ["·  워커 안 돔 → /check-agents 후",
         "                  /orcauto-start",
         "·  토큰 한도 → claude-quota-check.sh",
         "                자동 fallback 활성",
         "·  sync 안 됨 → bash sync-plugins.sh",
         "·  세션 끊김 → context-cache 자동 복구",
         "·  태스크 무한루프 → /loop-stop"])

    cover_slide(prs,
        "감사합니다",
        "다음: 플러그인 사용법",
        f"문의: github.com/bernakilljos/orchestration · {TODAY}")

    out = OUT_DIR / f"2-user-steps-{TODAY}.pptx"
    prs.save(str(out))
    return out

# ════════════════════════════════════════
#  PPT-3 : plugins 사용법
# ════════════════════════════════════════

def build_ppt3():
    prs = new_deck()
    cover_slide(prs,
        "Plugins 사용법",
        "21 개 플러그인 — 시나리오별·카테고리별 정리",
        f"orchestration_v1 · {TODAY}")

    content_slide(prs, "목차",
        ["1.  plugins/ 구조 — SoT 패턴",
         "2.  4 카테고리 — exec / mcp / design / review",
         "3.  exec_  계열 — 실행·라우팅",
         "4.  mcp_   계열 — 외부 도구 설치",
         "5.  design_ 계열 — 문서·시각 생성",
         "6.  시나리오 1: 큰 코드 구현",
         "7.  시나리오 2: 인터랙티브 결과물",
         "8.  시나리오 3: 발표·다이어그램",
         "9.  시나리오 4: 음성·자동화·로컬",
         "10. 새 플러그인 만들기 — _template/"])

    section_title_slide(prs, "01", "plugins/ 구조",
        "Single Source of Truth — 편집은 plugins/ 에서만.\n.claude/ 는 sync 결과물 (자동 생성).")

    content_slide(prs, "1. plugins/ 구조 — SoT 패턴",
        ["plugins/                     ←  편집은 여기만 (Source of Truth)",
         "  exec_orch/",
         "    plugin.json              ←  메타데이터·의존성",
         "    README.md                ←  사용법",
         "    SPEC.md                  ←  상세 스펙",
         "    commands/*.md            ←  슬래시 커맨드",
         "    skills/*.md              ←  자동 활성화 스킬",
         "    hooks/*.sh               ←  라이프사이클",
         "    scripts/                 ←  보조 스크립트",
         "    agents/*.md              ←  서브에이전트 정의",
         "",
         ".claude/                     ←  sync 결과물 (자동, 직접 편집 ❌)",
         "  commands/  skills/  hooks/  agents/",
         "",
         "동기화: bash .claude/scripts/sync-plugins.sh"])

    section_title_slide(prs, "02", "4 카테고리",
        "용도별로 4 그룹 — exec / mcp / design / review.")

    mindmap_slide(prs, "2. 4 카테고리 한 장 보기",
        center="Plugins",
        categories=[
            ("exec_",   ["orch", "claude (NEW)", "voice·learning"]),
            ("mcp_",    ["dev·collab", "data·docs", "media·web"]),
            ("design_", ["ppt·pdf", "excel·word", "web·video"]),
            ("review_", ["qa", "(미래)", "(미래)"]),
        ])

    section_title_slide(prs, "03", "exec_  계열",
        "실행·라우팅·세션 관리.")

    content_slide(prs, "3. exec_ 계열 — 실행 코어",
        ["•  exec_orch          멀티AI 라우팅 엔진 (Claude+Codex+Gemini)",
         "                        /exec_status  /check-agents  /orcauto-start",
         "",
         "•  exec_claude  ★NEW  Claude 깊이 활용 (이번에 추가)",
         "                        /claude-status  /claude-ask  /claude-artifact",
         "                        /claude-connectors  /claude-thinking",
         "",
         "•  exec_learning      세션 학습·실패 패턴 축적",
         "                        /learn  /recall  /summarize",
         "",
         "•  exec_session_guard 토큰 소진 대비 자동 스냅샷",
         "                        /guard-save  /token-stats",
         "",
         "•  exec_voice         음성 STT·TTS·회의록",
         "                        /transcribe  /speak  /meeting"])

    section_title_slide(prs, "04", "mcp_  계열",
        "MCP 서버 일괄 설치 허브 — 그룹별로 한 번에.")

    content_slide(prs, "4. mcp_ 계열 — 외부 도구 설치 허브",
        ["•  mcp_dev          GitHub·GitLab·Docker·K8s·AWS·Firebase·Vercel",
         "•  mcp_collab       Slack·Notion·Jira·Trello·Gmail·Google Calendar",
         "•  mcp_data         MySQL·PostgreSQL·MongoDB·BigQuery·Sheets·Airtable",
         "•  mcp_docs         PDF·DOCX·OCR (Tesseract)",
         "•  mcp_media        Whisper(STT)·TTS·FFmpeg",
         "•  mcp_web          Playwright·Puppeteer·Selenium·Apify",
         "•  mcp_queue        Kafka·RabbitMQ·Redis·SQS  (스펙)",
         "•  mcp_social       YouTube Data API v3       (스펙)",
         "",
         "사용: /mcp_dev-install   →  설치",
         "       /mcp_dev-status   →  상태 확인",
         "       /plug_all          →  7 그룹 일괄 설치"])

    section_title_slide(prs, "05", "design_ 계열",
        "문서·시각 결과물 자동 생성.")

    content_slide(prs, "5. design_ 계열 — 결과물 생성",
        ["•  design_ppt    PPT (Canva·Figma·Gamma 통합)",
         "                  /make-ppt  /design_ppt  /ai-system-stages",
         "                  /arch-mindmap  /arch-auto  ★ NEW (다이어그램)",
         "",
         "•  design_pdf    PDF 양식·서명·암호화·변환",
         "                  /pdf-fill  /pdf-generate  /pdf-secure  /pdf-sign",
         "                  /arch-layered  /arch-cheatsheet  ★ NEW",
         "",
         "•  design_excel  Excel + 차트 + Sheets",
         "                  /excel-make  /excel-status",
         "",
         "•  design_word   Word + python-docx + Mermaid + PDF",
         "                  /word-make  /word-status",
         "",
         "•  design_web    랜딩·블로그·포트폴리오·SEO 메타"])

    section_title_slide(prs, "06", "시나리오 1",
        "큰 코드 구현 — 500 줄 이상 신규 모듈.")

    steps_slide(prs, "6. 시나리오 1 — 큰 코드 구현 워크플로우",
        [(1, "사용자",
            "\"전체 결제 모듈 만들어줘.\nNext.js 14, Stripe, Webhook\""),
         (2, "Claude (Opus)",
            "task-instruction.md 작성.\n아키텍처 결정·파일 분할."),
         (3, "Codex × 4",
            "병렬 구현 (codex-auto).\n각 워커가 모듈별 1 차 코드."),
         (4, "Claude (Sonnet)",
            "결과 보완·통합.\n에러 처리·테스트."),
         (5, "Gemini Flash",
            "검증·리뷰·보안 점검.\n저단가 마지막 게이트."),
         (6, "Claude → 사용자",
            "최종 보고 + git push 제안.")])

    section_title_slide(prs, "07", "시나리오 2",
        "인터랙티브 결과물 — 단순 PDF 가 아닌 클릭 가능한 HTML.")

    content_slide(prs, "7. 시나리오 2 — 인터랙티브 결과물",
        ["[사용자]   \"분기별 매출 대시보드 만들어줘. 인터랙티브로\"",
         "",
         "[Claude]   skill-claude-artifact 활성화",
         "             →  타입 추론: dashboard",
         "             →  Chart.js 단일 HTML 생성",
         "             →  outputs/artifacts/dashboard-매출-2026-04-23.html",
         "",
         "[결과]     단일 HTML 파일 (외부 의존: CDN 만)",
         "             ·  더블클릭 → 브라우저에서 즉시 동작",
         "             ·  모바일 반응형",
         "             ·  다크/라이트 자동 감지",
         "             ·  100KB 이하",
         "",
         "관련 커맨드: /claude-artifact dashboard <topic>"])

    section_title_slide(prs, "08", "시나리오 3",
        "발표·다이어그램 — 한 장에 핵심을 담는 그림.")

    content_slide(prs, "8. 시나리오 3 — 발표·다이어그램",
        ["[사용자]   \"우리 킷 전체 구조 그림으로 만들어줘\"",
         "",
         "[Claude]   /arch-auto orchestration_v1 자동 실행",
         "             →  주제 분석: \"여러 영역 전체뷰\" → mindmap 선택",
         "             →  Mermaid mindmap 생성 → PNG 렌더",
         "             →  outputs/arch/mindmap-orchestration_v1-2026-04-23.png",
         "",
         "[패턴 선택지]",
         "  /arch-mindmap     방사형 (전체 영역)        — 발표·SNS",
         "  /arch-layered     레이어 케이크 (계층)      — 기술 문서",
         "  /arch-cheatsheet  3 컬럼 (빠른 참조)        — 온보딩",
         "  /arch-auto        주제 보고 자동 판단",
         "",
         "[발표 자료 풀세트]",
         "  /ai-system-stages   AI 시스템 6 단계 PPT 자동 생성"])

    section_title_slide(prs, "09", "시나리오 4",
        "음성 인터페이스 + 로컬 LLM 으로 비용 0 운영.")

    content_slide(prs, "9. 시나리오 4 — 음성·로컬·자동화",
        ["[음성으로 작업 의뢰]",
         "  /voice-task          말한 내용 → task-instruction.md 자동 작성",
         "  /transcribe          음성 파일 → 텍스트 (Whisper)",
         "  /speak               결과 읽어주기 (edge-tts)",
         "",
         "[로컬 LLM 으로 비용 0]",
         "  /exec_offline-setup    Ollama + ChromaDB + Phoenix 일괄 설치",
         "  /exec_offline-route    API vs 로컬 라우팅 (비용 임계 기준)",
         "  /exec_offline-model    Llama·Gemma·Mistral 다운로드",
         "",
         "[스케줄·자동화]",
         "  /exec_scheduler-cron    YAML 선언형 크론 잡",
         "  /exec_scheduler-workflow DAG 워크플로우 정의",
         "  /vibe-loop              멀티에이전트 자동 루프"])

    section_title_slide(prs, "10", "새 플러그인 만들기",
        "_template/ 복사 → 편집 → sync.")

    content_slide(prs, "10. 새 플러그인 만들기 — 5 단계",
        ["1.  cp -r plugins/_template plugins/<my_plugin>",
         "",
         "2.  plugin.json 수정",
         "    {\"name\":\"my_plugin\", \"display\":\"...\", \"prefix\":\"my_\", ...}",
         "",
         "3.  commands/ skills/ 채우기",
         "    각 .md 에 frontmatter (name·description) + 본문",
         "",
         "4.  bash .claude/scripts/sync-plugins.sh --dry  (미리보기)",
         "    bash .claude/scripts/sync-plugins.sh        (실제)",
         "",
         "5.  검증",
         "    python .claude/scripts/validate-plugin-schema.py",
         "    /help my_plugin  → 정상 등록 확인",
         "",
         "참조: .claude/rules/plugin-structure.md"])

    cover_slide(prs,
        "감사합니다",
        "전체 슬래시 커맨드: /help · 플러그인별: /help <name>",
        f"orchestration_v1 · {TODAY}")

    out = OUT_DIR / f"3-plugins-{TODAY}.pptx"
    prs.save(str(out))
    return out

# ════════════════════════════════════════
#  실행
# ════════════════════════════════════════

if __name__ == "__main__":
    import sys
    sys.stdout.reconfigure(encoding="utf-8")
    print("PPT 3-deck generation start...")
    f1 = build_ppt1()
    print(f"  [OK] PPT-1: {f1}")
    f2 = build_ppt2()
    print(f"  [OK] PPT-2: {f2}")
    f3 = build_ppt3()
    print(f"  [OK] PPT-3: {f3}")
    print("Done.")
